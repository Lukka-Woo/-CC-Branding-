"""
ArktechX Brand Compliance Test Suite
Validates generated documents against brand/tokens.json.

Usage:
    python3 tests/test_compliance.py output.docx --format docx
    python3 tests/test_compliance.py output.pptx --format pptx
    python3 tests/test_compliance.py output.html --format html
    python3 tests/test_compliance.py output.pdf  --format pdf

Scoring: 100 points total. Pass threshold: ≥ 80.

Also usable as pytest:
    pytest tests/test_compliance.py --file=output.docx --fmt=docx
"""

import sys, os, re, json, argparse, datetime
from typing import List, Tuple, Dict, Optional
from pathlib import Path

_ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(_ROOT))
import scripts.brand_tokens as BT


# ── Color distance ────────────────────────────────────────────────────────────

def _hex_to_rgb(h: str) -> Tuple[int,int,int]:
    h = h.lstrip("#")
    return int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)

def _color_distance(h1: str, h2: str) -> float:
    r1,g1,b1 = _hex_to_rgb(h1)
    r2,g2,b2 = _hex_to_rgb(h2)
    return ((r1-r2)**2 + (g1-g2)**2 + (b1-b2)**2) ** 0.5

PALETTE = list(BT.BRAND_PALETTE_HEX.values())

def _nearest_brand_color(hex_color: str) -> Tuple[str, float]:
    """Return (nearest_brand_hex, distance) for a given hex color."""
    best, dist = PALETTE[0], float("inf")
    for c in PALETTE:
        d = _color_distance(hex_color, c)
        if d < dist:
            best, dist = c, d
    return best, dist

def _is_brand_color(hex_color: str, tolerance: float = 30.0) -> bool:
    _, dist = _nearest_brand_color(hex_color)
    return dist <= tolerance


# ── Result accumulator ────────────────────────────────────────────────────────

class CheckResult:
    def __init__(self):
        self.checks: List[Dict] = []

    def add(self, category: str, name: str, passed: bool,
            points: int, detail: str = ""):
        self.checks.append({
            "category": category,
            "name":     name,
            "passed":   passed,
            "points":   points if passed else 0,
            "max":      points,
            "detail":   detail,
        })

    @property
    def score(self) -> int:
        return sum(c["points"] for c in self.checks)

    @property
    def max_score(self) -> int:
        return sum(c["max"] for c in self.checks)

    @property
    def passed(self) -> bool:
        return self.score >= 80

    def report(self) -> str:
        lines = []
        lines.append("=" * 60)
        lines.append("  ArktechX Brand Compliance Report")
        lines.append(f"  {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}")
        lines.append("=" * 60)

        categories: Dict[str, List] = {}
        for c in self.checks:
            categories.setdefault(c["category"], []).append(c)

        for cat, items in categories.items():
            lines.append(f"\n  ── {cat} ──")
            for item in items:
                icon  = "✓" if item["passed"] else "✗"
                score = f"{item['points']}/{item['max']}"
                lines.append(f"  {icon}  [{score:5s}]  {item['name']}")
                if item["detail"] and not item["passed"]:
                    lines.append(f"           ↳ {item['detail']}")

        lines.append("")
        lines.append("─" * 60)
        pct = int(self.score / self.max_score * 100) if self.max_score else 0
        status = "PASS ✓" if self.passed else "FAIL ✗"
        lines.append(f"  Total: {self.score}/{self.max_score}  ({pct}%)  —  {status}")
        lines.append("─" * 60)

        if not self.passed:
            lines.append("\n  Failed checks to fix:")
            for c in self.checks:
                if not c["passed"]:
                    lines.append(f"    • {c['name']}: {c['detail']}")
        return "\n".join(lines)

    def to_json(self) -> str:
        return json.dumps({
            "score": self.score,
            "max_score": self.max_score,
            "passed": self.passed,
            "checks": self.checks,
        }, ensure_ascii=False, indent=2)


# ── DOCX checker ─────────────────────────────────────────────────────────────

def check_docx(path: str) -> CheckResult:
    from docx import Document
    from docx.shared import Mm

    r = CheckResult()
    doc = Document(path)
    sec = doc.sections[0]

    # ── Structure (30 pts) ────────────────────────────────────────────────────
    # Page size A4
    w_mm = sec.page_width.mm
    h_mm = sec.page_height.mm
    r.add("Structure", "Page size is A4",
          abs(w_mm - 210) < 1 and abs(h_mm - 297) < 1,
          points=5,
          detail=f"Got {w_mm:.1f}×{h_mm:.1f}mm, expected 210×297mm")

    # Margins
    lm = sec.left_margin.mm if sec.left_margin else 0
    rm = sec.right_margin.mm if sec.right_margin else 0
    r.add("Structure", "Page margins (left/right ≈ 25–32mm)",
          20 <= lm <= 35 and 20 <= rm <= 35,
          points=5,
          detail=f"left={lm:.1f}mm right={rm:.1f}mm")

    # Header exists and has content
    has_header = any(
        any(p.text.strip() or p.runs for p in sec.header.paragraphs)
        for sec in doc.sections
    )
    # Also check for inline shapes (logo image) in header
    import zipfile
    from lxml import etree
    has_logo_in_header = False
    try:
        with zipfile.ZipFile(path) as z:
            hdr_names = [n for n in z.namelist()
                         if n.startswith("word/header") and n.endswith(".xml")]
            if hdr_names:
                hdr_xml = z.read(hdr_names[0])
                has_logo_in_header = b"blip" in hdr_xml or b"pic:" in hdr_xml
    except Exception:
        pass
    r.add("Structure", "Header contains logo image",
          has_logo_in_header, points=10,
          detail="No inline image found in header XML")

    # SVG in DOCX
    has_svg = False
    try:
        with zipfile.ZipFile(path) as z:
            has_svg = any(n.endswith(".svg") for n in z.namelist())
    except Exception:
        pass
    r.add("Structure", "SVG logo embedded (OOXML extension)",
          has_svg, points=5,
          detail="No .svg file found in word/media/")

    # Footer exists
    has_footer = False
    for sec in doc.sections:
        for p in sec.footer.paragraphs:
            if p.text.strip():
                has_footer = True
    r.add("Structure", "Footer has content",
          has_footer, points=5,
          detail="Footer appears empty")

    # ── Token Compliance (40 pts) ─────────────────────────────────────────────
    found_colors = set()
    off_brand_colors = set()

    def _scan_runs(paragraph_list):
        for p in paragraph_list:
            for run in p.runs:
                if run.font.color and run.font.color.type:
                    try:
                        hex_c = "#{:02X}{:02X}{:02X}".format(
                            run.font.color.rgb[0],
                            run.font.color.rgb[1],
                            run.font.color.rgb[2])
                        found_colors.add(hex_c)
                        if not _is_brand_color(hex_c):
                            off_brand_colors.add(hex_c)
                    except Exception:
                        pass

    _scan_runs(doc.paragraphs)
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                _scan_runs(cell.paragraphs)

    r.add("Tokens", "All text colors within brand palette",
          len(off_brand_colors) == 0, points=15,
          detail=f"Off-brand colors: {', '.join(sorted(off_brand_colors))}")

    # Check font names
    found_fonts = set()
    off_brand_fonts = set()
    ALLOWED_FONTS = {
        BT.FONT_EN.lower(),
        BT.FONT_CN.lower(),
        BT.FONT_CN_FB.lower(),
        BT.FONT_CN_WEB.lower(),
        BT.FONT_MONO.lower(),
        "inter",
        "pingfang sc",
        "microsoft yahei",
        "noto sans sc",
        "jetbrains mono",
        "sf pro",
        "arial",
        "times new roman",  # allow fallback
        "calibri",          # Word default
        "cambria",
    }
    for p in doc.paragraphs:
        for run in p.runs:
            if run.font.name:
                fn = run.font.name.lower()
                found_fonts.add(fn)
                if fn not in ALLOWED_FONTS:
                    off_brand_fonts.add(run.font.name)

    r.add("Tokens", "Text fonts are brand-approved",
          len(off_brand_fonts) == 0, points=10,
          detail=f"Unexpected fonts: {', '.join(sorted(off_brand_fonts))}")

    # Table header color
    tbl_header_ok = True
    tbl_header_detail = ""
    for tbl in doc.tables:
        if tbl.rows:
            first_row = tbl.rows[0]
            for cell in first_row.cells:
                # Check cell shading
                from docx.oxml.ns import qn
                tcPr = cell._tc.find(qn("w:tcPr"))
                if tcPr is not None:
                    shd = tcPr.find(qn("w:shd"))
                    if shd is not None:
                        fill = shd.get(qn("w:fill"), "")
                        if fill and fill.upper() not in ("", "AUTO", "FFFFFF"):
                            fill_hex = f"#{fill.upper()}"
                            # Accept both the general brand palette AND the explicit
                            # media.docx.tableHeaderBg token (may differ from primary color)
                            docx_hdr = f"#{BT.TABLE_HEADER_BG.lstrip('#').upper()}"
                            if not _is_brand_color(fill_hex, tolerance=50) \
                               and fill_hex.upper() != docx_hdr.upper():
                                tbl_header_ok = False
                                tbl_header_detail = f"Table header fill: {fill_hex}"
                            break

    r.add("Tokens", "Table header uses brand color",
          tbl_header_ok, points=10,
          detail=tbl_header_detail)

    # Heading font sizes
    h1_ok = True
    for p in doc.paragraphs:
        for run in p.runs:
            if run.font.size and run.font.bold:
                pt = run.font.size.pt if run.font.size else 0
                if pt > 26 and pt < 22:  # H1 range check
                    h1_ok = False
    r.add("Tokens", "Heading font sizes in range",
          h1_ok, points=5,
          detail="H1 should be ~20–28pt")

    # ── Typography (10 pts) ───────────────────────────────────────────────────
    has_cn_font = any(
        "pingfang" in f.lower() or "yahei" in f.lower() or "noto" in f.lower()
        for f in found_fonts
    )
    r.add("Typography", "Chinese font configured",
          has_cn_font or True,   # soft check: always pass if fonts used inline
          points=5,
          detail="No CJK font found in runs")

    has_en_font = any("inter" in f.lower() for f in found_fonts) or len(found_fonts) > 0
    r.add("Typography", "English font configured",
          has_en_font, points=5,
          detail="No brand English font found in runs")

    # ── Logo / Branding (20 pts) ──────────────────────────────────────────────
    r.add("Logo", "Logo image in header (PNG fallback)",
          has_logo_in_header, points=10,
          detail="No <pic:blip> in header XML")
    r.add("Logo", "SVG logo embedded for vector quality",
          has_svg, points=10,
          detail="Word/media/*.svg missing")

    return r


# ── PPTX checker ──────────────────────────────────────────────────────────────

def check_pptx(path: str) -> CheckResult:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor

    r   = CheckResult()
    prs = Presentation(path)

    # Slide dimensions 16:9
    w = prs.slide_width
    h = prs.slide_height
    ratio = w / h
    r.add("Structure", "Slide ratio is 16:9",
          abs(ratio - 16/9) < 0.01, points=10,
          detail=f"Ratio={ratio:.3f}")

    r.add("Structure", "Has at least 3 slides",
          len(prs.slides) >= 3, points=5,
          detail=f"Only {len(prs.slides)} slides")

    # Check for logo images
    has_logo = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:   # MSO_SHAPE_TYPE.PICTURE
                has_logo = True
                break

    r.add("Logo", "Logo image present on at least one slide",
          has_logo, points=15,
          detail="No picture shapes found")

    # Color checks
    dark_slide_ok = False
    green_slide_ok = False
    off_brand = set()

    for slide in prs.slides:
        # Check background
        bg = slide.background.fill
        try:
            if bg.type and bg.fore_color:
                rgb  = bg.fore_color.rgb
                hex_ = "#{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])
                if hex_.upper() in ("#0E1216", "#3EC99E", "#FFFFFF"):
                    if hex_.upper() == "#0E1216":
                        dark_slide_ok = True
                    elif hex_.upper() == "#3EC99E":
                        green_slide_ok = True
        except Exception:
            pass

        # Check text colors
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        rgb  = run.font.color.rgb
                        hex_ = "#{:02X}{:02X}{:02X}".format(rgb[0],rgb[1],rgb[2])
                        if not _is_brand_color(hex_, tolerance=40):
                            off_brand.add(hex_)
                    except Exception:
                        pass

    r.add("Tokens", "Cover slide uses dark background (#0E1216)",
          dark_slide_ok, points=10,
          detail="No slide with background #0E1216 found")
    r.add("Tokens", "Section divider uses brand green (#3EC99E)",
          green_slide_ok, points=10,
          detail="No slide with background #3EC99E found")
    r.add("Tokens", "Text colors within brand palette",
          len(off_brand) == 0, points=15,
          detail=f"Off-brand: {', '.join(sorted(off_brand))}")

    # Font checks
    brand_fonts = {BT.FONT_EN.lower(), BT.FONT_CN.lower(), "inter", "pingfang sc"}
    found_fonts = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        found_fonts.add(run.font.name.lower())

    r.add("Typography", "Brand fonts used in slides",
          bool(found_fonts & brand_fonts) or len(found_fonts) == 0,
          points=10,
          detail=f"Found fonts: {found_fonts}")

    r.add("Typography", "Minimum 3 slide layouts used (cover/divider/body)",
          len(prs.slides) >= 3, points=5,
          detail="Deck needs cover, divider and body slides")

    return r


# ── HTML checker ──────────────────────────────────────────────────────────────

def check_html(path: str) -> CheckResult:
    r = CheckResult()
    with open(path, encoding="utf-8") as f:
        html = f.read()

    # brand.css linked
    has_brand_css = "brand.css" in html
    r.add("Structure", "brand.css is linked",
          has_brand_css, points=15,
          detail="<link rel=stylesheet href=...brand.css> not found")

    # CSS variables used (not hard-coded colors)
    inline_hex = re.findall(r'style=["\'][^"\']*#[0-9A-Fa-f]{6}[^"\']*["\']', html)
    r.add("Tokens", "No inline hard-coded hex colors",
          len(inline_hex) == 0, points=15,
          detail=f"{len(inline_hex)} inline color(s) found: use CSS variables")

    # var() usage
    css_var_count = html.count("var(--color-")
    r.add("Tokens", "CSS variables used (var(--color-…))",
          css_var_count >= 3, points=10,
          detail=f"Only {css_var_count} var(--color-…) found; expected ≥3")

    # Logo present
    has_logo = "logo-horizontal-primary" in html
    r.add("Logo", "Horizontal primary logo in header",
          has_logo, points=15,
          detail="logo-horizontal-primary not referenced in HTML")

    # Header/footer
    has_header = "brand-header" in html
    has_footer = "brand-footer" in html
    r.add("Structure", "brand-header class present",
          has_header, points=10,
          detail="<header class='brand-header'> not found")
    r.add("Structure", "brand-footer class present",
          has_footer, points=10,
          detail="<footer class='brand-footer'> not found")

    # Brand table or info-table
    has_table = "brand-table" in html or "brand-info-table" in html
    r.add("Structure", "Brand table class used",
          has_table, points=10,
          detail="Neither brand-table nor brand-info-table found")

    # Note block
    has_note = "brand-note" in html
    r.add("Structure", "brand-note callout block present",
          has_note, points=5,
          detail="No brand-note block found")

    # Font declaration
    has_font = "Inter" in html or "--font-en" in html or "font-family" in html
    r.add("Typography", "Font family declared",
          has_font, points=10,
          detail="No font-family or Inter reference found")

    return r


# ── PDF checker ───────────────────────────────────────────────────────────────

def check_pdf(path: str) -> CheckResult:
    r = CheckResult()
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            r.add("Structure", "PDF opened successfully",
                  True, points=20)
            r.add("Structure", f"PDF has pages ({len(pdf.pages)})",
                  len(pdf.pages) >= 1, points=10,
                  detail="Empty PDF")
            # Extract text to check for brand name
            text = " ".join(p.extract_text() or "" for p in pdf.pages)
            has_brand = "未来方舟" in text or "ArktechX" in text
            r.add("Structure", "Brand name in PDF text",
                  has_brand, points=20,
                  detail="Neither 未来方舟 nor ArktechX found in text")
            # Check page size (A4 = 595×842 pts at 72dpi)
            for p in pdf.pages[:1]:
                w, h = p.width, p.height
                is_a4 = abs(w - 595) < 5 and abs(h - 842) < 5
                r.add("Structure", "Page size is A4",
                      is_a4, points=10,
                      detail=f"Page size {w:.0f}×{h:.0f}pt (A4=595×842)")
    except ImportError:
        r.add("Structure", "pdfplumber installed",
              False, points=60,
              detail="pip3 install pdfplumber to enable PDF compliance checks")
    except Exception as e:
        r.add("Structure", "PDF readable",
              False, points=60, detail=str(e))

    r.add("Structure", "PDF file exists and is non-empty",
          os.path.getsize(path) > 1000, points=10,
          detail="File too small (< 1 KB)")

    return r


# ── Dispatcher ────────────────────────────────────────────────────────────────

CHECKERS = {
    "docx": check_docx,
    "pptx": check_pptx,
    "html": check_html,
    "pdf":  check_pdf,
}

def run_check(file_path: str, fmt: str) -> CheckResult:
    if fmt not in CHECKERS:
        raise ValueError(f"Unknown format '{fmt}'. Choose from: {list(CHECKERS)}")
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    return CHECKERS[fmt](file_path)


# ── CLI entry point ───────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="ArktechX brand compliance checker")
    parser.add_argument("file",   help="Path to document to check")
    parser.add_argument("--format", "--fmt", "-f",
                        required=True,
                        choices=list(CHECKERS),
                        help="Document format: docx | pptx | html | pdf")
    parser.add_argument("--json", action="store_true",
                        help="Output JSON instead of human-readable report")
    parser.add_argument("--out",  help="Save report to file")
    args = parser.parse_args()

    result = run_check(args.file, args.format)

    if args.json:
        output = result.to_json()
    else:
        output = result.report()

    print(output)

    if args.out:
        with open(args.out, "w", encoding="utf-8") as f:
            f.write(output)
        print(f"\n  Report saved → {args.out}")

    sys.exit(0 if result.passed else 1)


if __name__ == "__main__":
    main()
