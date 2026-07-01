"""
BrandPptx — Brand-consistent PPTX builder for 未来方舟 / ArktechX  v2.0

Usage:
    from scripts.pptx_builder import BrandPptx

    prs = BrandPptx()
    prs.add_cover("AI 智能化解决方案", "ArktechX 2025")
    prs.add_cover_light("轻色封面版本", "Sub-Title")
    prs.add_divider("第一章：产品概述")
    prs.add_body_slide("核心功能", ["要点一", "要点二"])
    prs.add_big_stats("核心指标", [("98%", "准确率", "自动核算"), ...])
    prs.add_three_cards("三大优势", [{"title": "...", "body": "..."}, ...])
    prs.add_timeline("路线图", [("2024 Q1", "里程碑A", "说明"), ...])
    # add_timeline is adaptive: ≤5 → horizontal, 6+ → vertical two-column
    prs.add_quote("引用内容", "作者姓名", "职位")
    prs.add_table_slide("数据对比", ["列1","列2"], [["A","B"], ...])
    prs.add_closing("谢谢", "contact@futureark.ai")
    prs.save("output.pptx")

Gradient title note:
    add_cover_light / add_cover / add_closing / add_divider all use the
    brand title gradient: primary-500 → primary-600 → secondary-500.
    This mirrors the VS template's gradFill text approach.
"""

import os
import re
import math
import copy
from typing import List, Optional, Tuple, Dict, Any
import scripts.brand_tokens as BT

# Brand root directory (two levels up from this file: scripts/ → brand 3/)
# Used by the icon resolver to locate assets/icons/ without requiring callers
# to pass absolute paths for brand-wide shared icons.
_BRAND_ROOT: str = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
_ICONS_DIR:  str = os.path.join(_BRAND_ROOT, "assets", "icons")

# ── Decorative element paths — A/B/C 三系列 ──────────────────────────────────
# 规则：每页最多使用一个装饰，避免页面过于花哨。
#   A 系列（Primary #3EC99E）— 白色/浅色背景页面（最常用）
#   B 系列（Success  #5CC13C）— 灰色背景 / 整页图片页面（较少用）
#   C 系列（Secondary #C8E13C）— 深色背景页面（如 add_cover / add_closing）
# 同一系列的四种装饰：underline / circle / arrow1 / arrow2
_DECO_DIR: str = os.path.join(_BRAND_ROOT, "assets", "装饰性元素")
_DECO_SERIES: dict = {
    "a": {
        "underline": (os.path.join(_DECO_DIR, "A下划_Main Dec:重点:@4x.png"), 9.87),
        "circle":    (os.path.join(_DECO_DIR, "A画圈_答案@4x.png"),           3.56),
        "arrow1":    (os.path.join(_DECO_DIR, "A箭头1@4x.png"),               3.15),
        "arrow2":    (os.path.join(_DECO_DIR, "A箭头2@4x.png"),               3.63),
    },
    "b": {
        "underline": (os.path.join(_DECO_DIR, "B下划@4x.png"),  9.87),
        "circle":    (os.path.join(_DECO_DIR, "B画圈@4x.png"),  3.56),
        "arrow1":    (os.path.join(_DECO_DIR, "B箭头1@4x.png"), 3.14),
        "arrow2":    (os.path.join(_DECO_DIR, "B箭头2@4x.png"), 3.63),
    },
    "c": {
        "underline": (os.path.join(_DECO_DIR, "C下划@4x.png"),  10.32),
        "circle":    (os.path.join(_DECO_DIR, "C画圈@4x.png"),  3.61),
        "arrow1":    (os.path.join(_DECO_DIR, "C箭头1@4x.png"), 3.10),
        "arrow2":    (os.path.join(_DECO_DIR, "C箭头2@4x.png"), 3.63),
    },
}
_DECO_A: dict = _DECO_SERIES["a"]   # backward-compat alias
_DECO_CHAR_W_MM: float = 8.1        # 26pt bold CJK 估算字宽 (mm/字)

# ── Brand inline arrow images (bullet arrow replacement) ──────────────────────
# Original PNG natural direction: upper-right (~NE, 60° clockwise from east).
# CW +60° → east (→ replacement);  CCW −30° → north (↑ replacement).
_BULLET_ARROW_PATHS: dict = {
    "primary": os.path.join(_DECO_DIR, "primary arrow.png"),   # lime-green gradient
    "white":   os.path.join(_DECO_DIR, "white arrow.png"),     # white, for dark backgrounds
}
_BULLET_ARROW_ROT: dict = {
    "right": 60.0,   # CW 60° → pointing east
    "up":   -30.0,   # CCW 30° → pointing north
}

from pptx import Presentation
from pptx.util import Inches, Pt, Mm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _est_text_h_mm(text: str, sz_pt: float, w_mm: float, ls_pt: float = None) -> float:
    """Estimate rendered text block height in mm for CJK-dominant mixed text.

    Used for content-adaptive card sizing: compute required textbox height before
    drawing so the card background can be sized to match.

    sz_pt:  font size in points
    w_mm:   textbox inner width in mm (after horizontal padding)
    ls_pt:  line spacing in points (defaults to 1.4× font size)
    Returns height in mm including a 2mm buffer for PPTX internal margin.
    """
    if not text or not text.strip():
        return 0.0
    ls = ls_pt if ls_pt else sz_pt * 1.4
    # CJK full-width char ≈ sz × 0.353 mm; use 0.85 factor (conservative estimate)
    char_w = sz_pt * 0.353 * 0.85
    chars_per_line = max(1.0, w_mm / char_w)
    total_lines = 0
    for para in text.split('\n'):
        stripped = para.strip()
        total_lines += max(1, math.ceil(len(stripped) / chars_per_line)) if stripped else 1
    return total_lines * ls * 0.353 + 2.0  # +2 mm for PPTX default top+bottom inset


# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

ML = Mm(21)              # left / right margin
MR = Mm(21)
CW = SLIDE_W - ML - MR  # usable content width

HEADER_H  = Mm(36)
FOOTER_H  = Mm(13)
CONTENT_Y = HEADER_H
CONTENT_H = SLIDE_H - CONTENT_Y - FOOTER_H

C2_GAP = Mm(6)
C2_W   = (CW - C2_GAP) // 2
C3_GAP = Mm(5)
C3_W   = (CW - 2 * C3_GAP) // 3

# ── Grid layout engine ────────────────────────────────────────────────────────
# All add_slide() calls use these two gap constants; per-method local ROW_GAPs
# still exist in legacy add_* methods and are untouched.
ROW_GAP   = Mm(4)   # vertical gap between rows in a grid preset
COL_GAP   = Mm(6)   # horizontal gap between columns in a grid preset
INNER_PAD = Mm(5)   # uniform inner padding inside every container slot


def _grid_cols(ax, ay, aw, ah, weights):
    """Split (ax, ay, aw, ah) into vertical columns by weight ratios.
    Last column absorbs any rounding remainder so total width = aw exactly.
    Returns list of (l, t, w, h) per column, left-to-right."""
    n        = len(weights)
    usable_w = aw - COL_GAP * (n - 1)
    total    = sum(weights)
    slots, cur_x = [], ax
    for i, wt in enumerate(weights):
        col_w = int(usable_w * wt / total) if i < n - 1 else ax + aw - cur_x
        slots.append((cur_x, ay, col_w, ah))
        cur_x += col_w + COL_GAP
    return slots


def _grid_rows(ax, ay, aw, ah, weights):
    """Split (ax, ay, aw, ah) into horizontal rows by weight ratios.
    Last row absorbs any rounding remainder so total height = ah exactly.
    Returns list of (l, t, w, h) per row, top-to-bottom."""
    n        = len(weights)
    usable_h = ah - ROW_GAP * (n - 1)
    total    = sum(weights)
    slots, cur_y = [], ay
    for i, wt in enumerate(weights):
        row_h = int(usable_h * wt / total) if i < n - 1 else ay + ah - cur_y
        slots.append((ax, cur_y, aw, row_h))
        cur_y += row_h + ROW_GAP
    return slots


# ── Preset helpers for compound layouts ──────────────────────────────────────
def _preset_top_2col(ax, ay, aw, ah):
    rows = _grid_rows(ax, ay, aw, ah, [1, 2])
    return [rows[0]] + _grid_cols(ax, rows[1][1], aw, rows[1][3], [1, 1])


def _preset_2col_bot(ax, ay, aw, ah):
    rows = _grid_rows(ax, ay, aw, ah, [2, 1])
    return _grid_cols(ax, rows[0][1], aw, rows[0][3], [1, 1]) + [rows[1]]


def _preset_top_3col(ax, ay, aw, ah):
    rows = _grid_rows(ax, ay, aw, ah, [1, 2])
    return [rows[0]] + _grid_cols(ax, rows[1][1], aw, rows[1][3], [1, 1, 1])


def _preset_left_2row(ax, ay, aw, ah):
    col_w   = (aw - COL_GAP) // 2
    right_l = ax + col_w + COL_GAP
    right_w = ax + aw - right_l            # absorb rounding
    return _grid_rows(ax, ay, col_w, ah, [1, 1]) + [(right_l, ay, right_w, ah)]


def _preset_2row_right(ax, ay, aw, ah):
    col_w   = (aw - COL_GAP) // 2
    right_l = ax + col_w + COL_GAP
    right_w = ax + aw - right_l
    return [(ax, ay, col_w, ah)] + _grid_rows(right_l, ay, right_w, ah, [1, 1])


# ── Preset registry ───────────────────────────────────────────────────────────
#
#  Each callable signature: (ax, ay, aw, ah) → list[(l, t, w, h)]
#  ax/ay = top-left origin of the content area
#  aw/ah = available width / height (already excludes header, footer, callout)
#  Slots are returned in reading order (left→right, top→bottom).
#
GRID_PRESETS = {
    # Single area
    "full":        lambda ax, ay, aw, ah: [(ax, ay, aw, ah)],

    # Column splits  ─ all same height
    "2col":        lambda ax, ay, aw, ah: _grid_cols(ax, ay, aw, ah, [1, 1]),
    "3col":        lambda ax, ay, aw, ah: _grid_cols(ax, ay, aw, ah, [1, 1, 1]),

    # Row splits  ─ all same width
    "2row":        lambda ax, ay, aw, ah: _grid_rows(ax, ay, aw, ah, [1, 1]),
    "3row":        lambda ax, ay, aw, ah: _grid_rows(ax, ay, aw, ah, [1, 1, 1]),

    # Full-width top row  +  column bottom rows
    # slot 0 = top full-width (1/3 h),  slots 1-2 = bottom 2-col (2/3 h)
    "top_2col":    _preset_top_2col,
    # slots 0-1 = top 2-col (2/3 h),    slot 2  = bottom full-width (1/3 h)
    "2col_bot":    _preset_2col_bot,
    # slot 0 = top full-width (1/3 h),  slots 1-3 = bottom 3-col (2/3 h)
    "top_3col":    _preset_top_3col,

    # One column subdivided into rows, other column full height
    # slots 0-1 = left col (2 equal rows),  slot 2  = right col (full height)
    "left_2row":   _preset_left_2row,
    # slot 0  = left col (full height),     slots 1-2 = right col (2 equal rows)
    "2row_right":  _preset_2row_right,
}

# Expected slot count per preset — used by add_slide() to validate input.
GRID_SLOT_COUNT = {
    "full":        1,
    "2col":        2,
    "3col":        3,
    "2row":        2,
    "3row":        3,
    "top_2col":    3,
    "2col_bot":    3,
    "top_3col":    4,
    "left_2row":   3,
    "2row_right":  3,
}

# Logo (horizontal) dimensions for footer
_LOGO_H = Mm(7)
_LOGO_W = int(_LOGO_H * BT.LOGO_HORIZONTAL_ASPECT)

# Brand gradient stops (matches VS template palette, mapped to our tokens)
TITLE_GRADIENT = [
    (0,   BT.PRIMARY_500_HEX),    # #3EC99E
    (48,  BT.SUCCESS_HEX),        # #5CC13C  (semantic/success)
    (100, BT.SECONDARY_500_HEX),  # #C8E13C
]


# ── Low-level primitives ──────────────────────────────────────────────────────

def _apply_radius(shape, radius_mm, w, h):
    """Set rounded-rect adj on an existing shape from an absolute mm radius.
    radius_mm=BT.RADIUS_PILL_MM renders as a full circle/capsule."""
    geom = shape._element.find(qn("p:spPr")).find(qn("a:prstGeom"))
    if geom is None:
        return
    avLst = geom.find(qn("a:avLst"))
    if avLst is None:
        avLst = etree.SubElement(geom, qn("a:avLst"))
    gd = avLst.find(qn("a:gd"))
    if gd is None:
        gd = etree.SubElement(avLst, qn("a:gd"))
    gd.set("name", "adj")
    adj_val = min(50000, int(radius_mm * 36000 / min(w, h) * 100000))
    gd.set("fmla", f"val {adj_val}")


def _apply_shape_gradient(shape, stops, angle_deg=135):
    """Replace a shape's fill with a linear gradient.
    stops: [(pct_0_to_100, '#hexcolor'), ...]  angle_deg: clockwise from left."""
    spPr = shape._element.find(qn("p:spPr"))
    for fill_tag in [qn("a:solidFill"), qn("a:gradFill"), qn("a:pattFill"),
                     qn("a:noFill"),    qn("a:blipFill")]:
        el = spPr.find(fill_tag)
        if el is not None:
            insert_idx = list(spPr).index(el)
            spPr.remove(el)
            break
    else:
        insert_idx = len(list(spPr))

    gf = etree.Element(qn("a:gradFill"))
    gf.set("rotWithShape", "1")
    gsLst = etree.SubElement(gf, qn("a:gsLst"))
    for pct, hex_color in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(int(pct * 1000)))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", hex_color.lstrip("#"))
    lin = etree.SubElement(gf, qn("a:lin"))
    lin.set("ang", str(int(angle_deg * 60000)))
    lin.set("scaled", "0")
    spPr.insert(insert_idx, gf)


def _rect(slide, l, t, w, h, fill=None, line=None, lw_pt=0.75, radius_mm=0, fill_alpha=100):
    """Plain rectangle. radius_mm > 0 → rounded corners. fill_alpha 0–100 → semi-transparent fill."""
    s = slide.shapes.add_shape(5 if radius_mm > 0 else 1, int(l), int(t), int(w), int(h))
    if radius_mm > 0:
        _apply_radius(s, radius_mm, int(w), int(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
        if fill_alpha < 100:
            spPr = s._element.find(qn("p:spPr"))
            sf   = spPr.find(qn("a:solidFill"))
            if sf is not None:
                sc = sf.find(qn("a:srgbClr"))
                if sc is not None:
                    al = etree.SubElement(sc, qn("a:alpha"))
                    al.set("val", str(int(fill_alpha * 1000)))
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = _rgb(line)
        s.line.width = Pt(lw_pt)
    else:
        s.line.fill.background()
    return s


def _rect_gradient_h(slide, l, t, w, h, stops):
    """Rectangle with horizontal linear gradient fill. stops: [(pct_0-100, '#hex'), ...]"""
    _ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    def _Q(tag): return f"{{{_ns}}}{tag}"

    s = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    s.line.fill.background()

    spPr = s._element.find(qn("p:spPr"))

    _fill_tags = {_Q(t) for t in
                  ["solidFill", "gradFill", "pattFill", "noFill", "blipFill", "grpFill"]}
    for _child in list(spPr):
        if _child.tag in _fill_tags:
            spPr.remove(_child)

    gf = etree.Element(_Q("gradFill"))
    gf.set("rotWithShape", "1")
    gsLst = etree.SubElement(gf, _Q("gsLst"))
    for pct, hex_c in stops:
        gs = etree.SubElement(gsLst, _Q("gs"))
        gs.set("pos", str(int(pct * 1000)))
        srgb = etree.SubElement(gs, _Q("srgbClr"))
        srgb.set("val", hex_c.lstrip("#"))
    lin = etree.SubElement(gf, _Q("lin"))
    lin.set("ang", "0")    # 0 = left → right
    lin.set("scaled", "0")

    prstGeom = spPr.find(_Q("prstGeom"))
    if prstGeom is not None:
        spPr.insert(list(spPr).index(prstGeom) + 1, gf)
    else:
        spPr.append(gf)

    return s


def _set_run_fonts(run, en_font=BT.FONT_EN, cn_font=BT.FONT_CN):
    """Set latin + east-asian font on a run."""
    run.font.name = en_font
    rPr = run._r.get_or_add_rPr()
    for tag, face in [("a:latin", en_font), ("a:ea", cn_font)]:
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", face)


def _apply_gradient_to_run(run, stops=None, angle_deg=0):
    """
    Replace the run's color with a gradient fill (DrawingML gradFill).
    stops: list of (pct_0_to_100, '#hexcolor')
    angle_deg: 0 = left→right, 90 = top→bottom
    """
    if stops is None:
        stops = TITLE_GRADIENT
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    Q = lambda tag: f"{{{ns}}}{tag}"

    rPr = run._r.get_or_add_rPr()

    # Remove existing fill elements
    for tag in ["solidFill", "gradFill", "pattFill", "noFill", "blipFill"]:
        el = rPr.find(Q(tag))
        if el is not None:
            rPr.remove(el)

    gf = etree.Element(Q("gradFill"))
    gf.set("rotWithShape", "1")

    gsLst = etree.SubElement(gf, Q("gsLst"))
    for pct, hex_color in stops:
        gs = etree.SubElement(gsLst, Q("gs"))
        gs.set("pos", str(int(pct * 1000)))
        srgb = etree.SubElement(gs, Q("srgbClr"))
        srgb.set("val", hex_color.lstrip("#"))
        alpha_el = etree.SubElement(srgb, Q("alpha"))
        alpha_el.set("val", "100000")

    lin = etree.SubElement(gf, Q("lin"))
    lin.set("ang", str(int(angle_deg * 60000)))

    # Insert before any font spec elements so rendering is correct
    rPr.insert(0, gf)


def _txb(slide, text, l, t, w, h,
         sz=16, bold=False, color=None,
         align=PP_ALIGN.LEFT, wrap=True,
         ls_pt=None, valign=None, alpha=None,
         cn_font=BT.FONT_CN, en_font=BT.FONT_EN):
    """Add a plain text box. color=None keeps default (white).
    alpha: opacity 0–100 (e.g. alpha=8 → 8% opaque / 92% transparent).
    """
    if color is None:
        color = BT.NEUTRAL_700_HEX
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    if valign is not None:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    if ls_pt:
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        etree.SubElement(lnSpc, qn("a:spcPts"),
                         attrib={"val": str(int(ls_pt * 100))})
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    if alpha is not None:
        # Inject <a:alpha val="N"/> into the srgbClr element (N = alpha% × 1000)
        _rPr = run._r.get_or_add_rPr()
        _sf  = _rPr.find(qn("a:solidFill"))
        if _sf is not None:
            _sc = _sf.find(qn("a:srgbClr"))
            if _sc is not None:
                etree.SubElement(_sc, qn("a:alpha")).set("val", str(int(alpha * 1000)))
    _set_run_fonts(run, en_font=en_font, cn_font=cn_font)
    return tb


def _clone_slide_from_master(dest_prs, master_path, slide_idx):
    """Clone master PPTX slide[slide_idx] into dest_prs, remapping image relationships."""
    _IMG_RT = ('http://schemas.openxmlformats.org/officeDocument/2006/'
               'relationships/image')
    src_prs   = Presentation(master_path)
    src_slide = src_prs.slides[slide_idx]
    dest_slide = dest_prs.slides.add_slide(dest_prs.slide_layouts[6])

    # Build src→dest rId map for images
    rId_map = {}
    for src_rId, rel in src_slide.part.rels.items():
        if rel.reltype == _IMG_RT and not rel.is_external:
            new_rId = dest_slide.part.relate_to(rel._target, rel.reltype)
            rId_map[src_rId] = new_rId

    # Deep-copy spTree XML; remap rIds in one string pass
    src_xml = etree.tostring(src_slide.shapes._spTree, encoding='unicode')
    for old_id, new_id in rId_map.items():
        src_xml = src_xml.replace(f'r:embed="{old_id}"', f'r:embed="{new_id}"')
        src_xml = src_xml.replace(f'r:id="{old_id}"',    f'r:id="{new_id}"')

    new_tree = etree.fromstring(src_xml)
    dst_tree = dest_slide.shapes._spTree
    dst_tree.getparent().replace(dst_tree, new_tree)
    return dest_slide


def _replace_shape_text(shape, parts):
    """Replace a shape's text with multi-color parts: [(text, hex_color), ...].
    '\\n' within a segment creates a new paragraph.
    Preserves the original font name and size from the first run.
    """
    tf    = shape.text_frame
    txBody = tf._txBody
    _ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    _q  = lambda t: f'{{{_ns}}}{t}'

    # Save rPr template (font name + size) from first run; strip fill
    rPr_tpl = None
    for p in txBody.findall(_q('p')):
        for r in p.findall(_q('r')):
            rPr = r.find(_q('rPr'))
            if rPr is not None:
                rPr_tpl = copy.deepcopy(rPr)
                for ft in [_q('solidFill'), _q('gradFill'), _q('pattFill'), _q('noFill')]:
                    fe = rPr_tpl.find(ft)
                    if fe is not None:
                        rPr_tpl.remove(fe)
            break
        if rPr_tpl is not None:
            break

    # Save first paragraph's pPr (alignment, spacing)
    first_pPr = None
    fps = txBody.findall(_q('p'))
    if fps:
        pPr = fps[0].find(_q('pPr'))
        if pPr is not None:
            first_pPr = copy.deepcopy(pPr)

    # Remove all existing paragraphs
    for p in txBody.findall(_q('p')):
        txBody.remove(p)

    # Rebuild with new parts
    cur_p       = None
    first_done  = False
    for text, color in parts:
        for seg_i, seg in enumerate(text.split('\n')):
            if cur_p is None or seg_i > 0:
                cur_p = etree.SubElement(txBody, _q('p'))
                if not first_done and first_pPr is not None:
                    cur_p.insert(0, copy.deepcopy(first_pPr))
                first_done = True
            if not seg:
                continue
            r_el  = etree.SubElement(cur_p, _q('r'))
            rPr_n = copy.deepcopy(rPr_tpl) if rPr_tpl is not None else etree.Element(_q('rPr'))
            rPr_n.set('lang', 'zh-CN')
            rPr_n.set('dirty', '0')
            sf   = etree.SubElement(rPr_n, _q('solidFill'))
            srgb = etree.SubElement(sf,    _q('srgbClr'))
            srgb.set('val', color.lstrip('#'))
            r_el.insert(0, rPr_n)
            t_el = etree.SubElement(r_el, _q('t'))
            t_el.text = seg

    if not txBody.findall(_q('p')):          # safety: at least one paragraph
        etree.SubElement(txBody, _q('p'))


def _txb_runs(slide, parts, l, t, w, h, sz=42, bold=True, ls_pt=None):
    """Multi-color text box. parts: [(text, hex_color), ...].
    '\\n' in any text segment creates a new paragraph."""
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    cur_para = tf.paragraphs[0]
    cur_para.alignment = PP_ALIGN.LEFT

    def _apply_ls(para):
        if ls_pt:
            pPr = para._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            etree.SubElement(lnSpc, qn("a:spcPts"),
                             attrib={"val": str(int(ls_pt * 100))})

    _apply_ls(cur_para)
    for text, color in parts:
        segments = text.split("\n")
        for i, seg in enumerate(segments):
            if i > 0:
                cur_para = tf.add_paragraph()
                cur_para.alignment = PP_ALIGN.LEFT
                _apply_ls(cur_para)
            if seg:
                run = cur_para.add_run()
                run.text = seg
                run.font.size = Pt(sz)
                run.font.bold = bold
                run.font.color.rgb = _rgb(color)
                _set_run_fonts(run)
    return tb


def _txb_gradient(slide, text, l, t, w, h,
                  sz=40, bold=True,
                  stops=None, angle_deg=0,
                  align=PP_ALIGN.LEFT, wrap=True,
                  ls_pt=None,
                  cn_font=BT.FONT_CN, en_font=BT.FONT_EN):
    """Add a text box with gradient-filled text (like VS template title)."""
    if stops is None:
        stops = TITLE_GRADIENT
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if ls_pt:
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        etree.SubElement(lnSpc, qn("a:spcPts"),
                         attrib={"val": str(int(ls_pt * 100))})
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    # Set fonts first, then replace color with gradient
    _set_run_fonts(run, en_font=en_font, cn_font=cn_font)
    _apply_gradient_to_run(run, stops=stops, angle_deg=angle_deg)
    return tb


def _set_slide_bg(slide, hex_color: str):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(hex_color)


def _add_logo_h(slide, right_mm=9, bottom_mm=3, h_mm=7, reverse=False):
    """Add horizontal logo to bottom-right."""
    if reverse:
        path = BT.LOGO_HORIZONTAL_REVERSE_PNG or BT.LOGO_HORIZONTAL_PRIMARY_PNG
    else:
        path = BT.LOGO_HORIZONTAL_PRIMARY_PNG
    if not (path and os.path.exists(path)):
        return None
    lh = Mm(h_mm)
    lw = int(lh * BT.LOGO_HORIZONTAL_ASPECT)
    return slide.shapes.add_picture(
        path,
        int(SLIDE_W - Mm(right_mm) - lw),
        int(SLIDE_H - Mm(bottom_mm) - lh),
        width=lw, height=lh)


def _add_logo_stacked(slide, reverse=True, l_mm=12, t_mm=10, h_mm=26):
    """Add stacked logo (cover / section pages)."""
    path = BT.LOGO_STACKED_REVERSE_PNG if reverse else BT.LOGO_STACKED_PRIMARY_PNG
    if not (path and os.path.exists(path)):
        return None
    lh = Mm(h_mm)
    from PIL import Image as _PILImage
    img = _PILImage.open(path)
    lw = int(lh * img.width / img.height)
    return slide.shapes.add_picture(path, Mm(l_mm), Mm(t_mm), width=lw, height=lh)


# ── Decorative element helpers ───────────────────────────────────────────────

def _send_to_back(slide, shape):
    """Move shape to bottom of z-order so it renders behind all existing shapes."""
    sp     = shape._element
    spTree = slide.shapes._spTree
    spTree.remove(sp)
    # index 0 = cNvGrpSpPr, index 1 = grpSpPr → shapes start at index 2
    spTree.insert(2, sp)


def _deco_img(slide, key: str, l_mm: float, t_mm: float, w_mm: float,
              send_back: bool = False, series: str = "a"):
    """Place a brand decorative PNG at the given position.
    key:       'underline' | 'circle' | 'arrow1' | 'arrow2'
    series:    'a' (white bg, primary)  | 'b' (gray/image bg, success)
               'c' (dark bg, secondary)
    send_back: move behind all existing shapes (use for circle to show text through it).
    """
    entry = _DECO_SERIES.get(series, _DECO_A).get(key)
    if not entry:
        return None
    path, ratio = entry
    if not os.path.exists(path):
        return None
    pic = slide.shapes.add_picture(path, Mm(l_mm), Mm(t_mm),
                                   width=Mm(w_mm), height=Mm(w_mm / ratio))
    if send_back:
        _send_to_back(slide, pic)
    return pic


_DECO_MIN_CHARS: int = 4   # 装饰笔画最小有效宽度（字符数）
# 规则：装饰的视觉宽度 ≥ 4 字，避免太细显得单薄。
# char_count 仍表示目标词的实际字数（用于确定锚点中心），
# 渲染宽度取 max(char_count, _DECO_MIN_CHARS)，并以目标词中心为基准居中。
# 示例：char_start=0, char_count=2（"八大"）→ 渲染宽=4字，居中后从 "八"左边 1字处开始，
#        但左边界不超过 ML；char_start=7, char_count=4（"三大困局"）→ 正好4字，保持原位。


def _apply_title_deco(slide, deco: dict, y_title_mm: float, title_sz_pt: float = 26.0,
                      anchor_top_mm: float = None, anchor_bot_mm: float = None):
    """Apply one decorative element to the header title area.

    deco keys:
      type       : "underline" | "circle"          (required)
      char_start : 0-based index of first target char (default 0)
      char_count : number of chars to cover          (default 3)
      char_w_mm  : override per-char width in mm     (default _DECO_CHAR_W_MM ≈ 8.1mm)
      series     : "a" | "b" | "c"                  (default "a")

    Width rule (enforced automatically):
      Rendered width = max(char_count, _DECO_MIN_CHARS) chars, centered on the target word.
      This keeps the decorative stroke visually consistent (~4-char width) regardless
      of whether the highlighted word is 2 chars or 6.

    Placement:
      underline → 1mm below title baseline, centered on the target chars
      circle    → vertically centered between anchor_top_mm (label bottom) and
                  anchor_bot_mm (subtitle top), so the gap above equals the gap below.
                  Rendered behind the title textbox so text reads through it.
    """
    deco_type  = deco.get("type", "underline")
    char_start = deco.get("char_start", 0)
    char_count = deco.get("char_count", 3)
    char_w     = deco.get("char_w_mm", _DECO_CHAR_W_MM)
    series     = deco.get("series", "a")
    ml_mm      = 21.0                            # ML = Mm(21) → 21mm
    line_h_mm  = title_sz_pt * 0.353 * 1.1      # ≈10.1mm for 26pt
    title_bot  = y_title_mm + line_h_mm

    # Enforce minimum visual width, centered on the target word
    eff_chars  = max(char_count, _DECO_MIN_CHARS)
    word_cx    = char_start * char_w + char_count * char_w / 2   # center of target word (from ML)
    raw_l      = ml_mm + word_cx - eff_chars * char_w / 2        # centered placement
    l_mm       = max(ml_mm, raw_l)                               # never bleed past left margin
    w_mm       = eff_chars * char_w

    if deco_type == "underline":
        _deco_img(slide, "underline",
                  l_mm = l_mm,
                  t_mm = title_bot + 1.0,
                  w_mm = w_mm, series = series)

    elif deco_type == "circle":
        ratio = _DECO_SERIES.get(series, _DECO_A).get("circle", (None, 3.56))[1]
        cw    = w_mm + 8.0                              # extra left/right padding for circle
        cl    = max(ml_mm, l_mm - 4.0)
        ch    = cw / ratio                              # circle height

        if anchor_top_mm is not None and anchor_bot_mm is not None:
            # Vertically center circle between label-bottom and subtitle-top,
            # so gap_above == gap_below.
            ct = (anchor_top_mm + anchor_bot_mm - ch) / 2
        else:
            ct = title_bot - ch                         # legacy: bottom-aligned

        _deco_img(slide, "circle",
                  l_mm = cl,
                  t_mm = ct,
                  w_mm = cw, series = series,
                  send_back = True)                      # render below title text


# ── Composite helpers: header / footer / card ─────────────────────────────────

def _header(slide, title, subtitle="", label="", title_deco=None):
    """Standard page header: label → title → subtitle (or accent line when no subtitle).

    title_deco: None | dict — add ONE decorative element below/around the title.
      {"type": "underline"|"circle", "char_start": int, "char_count": int}
      Maximum one decoration per slide. circle renders behind the title text.

    i18n: title accepts str | {"cn": ..., "en": ...}.
      lang="en"        → display EN string, wider box (CW*0.92), no word-wrap
      lang="bilingual" → display CN string; EN string fills subtitle slot
      lang="cn"        → display as-is (str or dict["cn"])
    """
    import scripts.i18n as _i18n
    from scripts.i18n import get_cn, get_en, T as _T

    is_en_only   = _i18n.LANG == "en"
    is_bilingual = _i18n.LANG == "bilingual"

    # Resolve what goes into the title text box
    if is_en_only:
        display_title = get_en(title) or _T(title)
    elif is_bilingual:
        display_title = get_cn(title) or _T(title)
    else:
        display_title = _T(title)

    # EN-only: wider box so longer English phrases fit on one line
    title_w = CW * 0.92 if is_en_only else CW * 0.82

    y_title_mm = 5.5 if not label else 11.0
    y_title    = Mm(y_title_mm)
    if label:
        _txb(slide, label, l=ML, t=Mm(4.5), w=Mm(80), h=Mm(6),
             sz=9, color=BT.PRIMARY_500_HEX)
    _txb(slide, display_title, l=ML, t=y_title, w=title_w, h=Mm(18),
         sz=26, bold=True, color=BT.NEUTRAL_900_HEX,
         wrap=not is_en_only)   # EN: single-line, no wrap

    # Bilingual: EN from title dict takes the subtitle slot (overrides passed subtitle)
    if is_bilingual:
        effective_subtitle = get_en(title) or subtitle
    else:
        effective_subtitle = subtitle

    if title_deco:
        # Compute anchor bounds so circle can be vertically centered between
        # label-bottom and subtitle-top (equal gap above and below).
        _anchor_top = (4.5 + 6.0) if label else y_title_mm
        _anchor_bot = (27.0 if label else 24.0) if effective_subtitle else (30.0 if label else 25.0)
        _apply_title_deco(slide, title_deco, y_title_mm=y_title_mm,
                          anchor_top_mm=_anchor_top, anchor_bot_mm=_anchor_bot)
    if effective_subtitle:
        _txb(slide, effective_subtitle, l=ML, t=Mm(24 if not label else 27),
             w=CW * 0.80, h=Mm(9), sz=12, color=BT.NEUTRAL_400_HEX)
    else:
        # No subtitle: thin gray divider below title
        _t = Mm(30) if label else Mm(25)
        _rect(slide, l=ML, t=_t, w=CW, h=Mm(0.3),
              fill=BT.NEUTRAL_200_HEX)


def _footer(slide, layout_label=""):
    """Standard page footer: logo + optional layout label."""
    _add_logo_h(slide, right_mm=9, bottom_mm=3, h_mm=7)
    if layout_label:
        _txb(slide, layout_label,
             l=SLIDE_W - Mm(72), t=SLIDE_H - FOOTER_H + Mm(1),
             w=Mm(60), h=Mm(10),
             sz=8, color=BT.NEUTRAL_200_HEX, align=PP_ALIGN.RIGHT)


def _card(slide, l, t, w, h, bg=None, border=None, radius_mm=None):
    """Card-shaped rounded rectangle. radius_mm=None → BT.RADIUS_SM_MM (4mm)."""
    if bg is None:
        bg = BT.NEUTRAL_100_HEX
    if radius_mm is None:
        radius_mm = BT.RADIUS_SM_MM
    use_round = radius_mm > 0
    s = slide.shapes.add_shape(5 if use_round else 1, int(l), int(t), int(w), int(h))
    if use_round:
        _apply_radius(s, radius_mm, int(w), int(h))
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(bg)
    if border:
        s.line.color.rgb = _rgb(border)
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s


def _is_numeric_val(s: str) -> bool:
    """True if value text is already a standalone number (e.g. '01', '8').
    When True, the sequence-number badge is suppressed to avoid duplication."""
    return bool(re.fullmatch(r'\d+', s.strip()))


def _seq_badge(slide, x, y, seq: int, color: str, size_mm: float = 9):
    """Sequence number badge — decorative colored number, no circle background.
    Uses Alibaba PuHuiTi 2.0 (CN font) for clean sans-serif digits.
    """
    s = Mm(size_mm)
    _txb(slide, f"{seq:02d}",
         l=x, t=y, w=s, h=s,
         sz=16, bold=True, color=color,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=False,
         en_font=BT.FONT_CN, cn_font=BT.FONT_CN)


def _brand_arrow(slide, x, y, size_mm=5.5, direction="right", style="primary"):
    """Insert a brand arrow PNG icon at position (x, y), bounding square = size_mm×size_mm.
    direction: "right" (CW 60°) | "up" (CCW 30°) | float (custom CW degrees).
    style: "primary" (lime-green gradient) | "white" (for dark backgrounds).
    """
    img_path = _BULLET_ARROW_PATHS.get(style, _BULLET_ARROW_PATHS["primary"])
    if not os.path.exists(img_path):
        return None
    s = Mm(size_mm)
    pic = slide.shapes.add_picture(img_path, int(x), int(y), int(s), int(s))
    rot = _BULLET_ARROW_ROT.get(direction, direction) if isinstance(direction, str) else float(direction)
    pic.rotation = rot
    return pic


def _render_content_with_arrows(slide, content, x, y, w, h,
                                  sz=14, color=None, ls_pt=22,
                                  arrow_style="primary", arrow_size_mm=5.5):
    """
    Render free-form content string, replacing standalone '→ text' lines with
    brand arrow PNG + bold text rows. Other lines rendered as normal text blocks.

    A 'standalone arrow line' is any line whose first non-space character is '→'.
    Mid-sentence '→' (e.g. 'A → B') are NOT affected.

    arrow_style: "primary" | "white"  (matches slide background tone)
    arrow_size_mm: arrow icon square size in mm (default 5.5, matches 14pt line)
    """
    color = color or BT.NEUTRAL_700_HEX
    if not content:
        return

    lines = content.split('\n')
    ARROW_SZ  = Mm(arrow_size_mm)
    ARROW_GAP = Mm(2.5)
    LINE_H_MM = ls_pt * 0.353          # approx mm per text line
    ROW_H     = Mm(max(arrow_size_mm + 1.5, LINE_H_MM))
    TEXT_W_MM = w / 36000

    # Group lines into sequential runs: ("text", joined_str) or ("arrow", text_after_arrow)
    runs = []
    text_buf = []
    for line in lines:
        if line.strip().startswith('→'):
            if text_buf:
                runs.append(("text", '\n'.join(text_buf)))
                text_buf = []
            arrow_text = line.strip()[1:].lstrip()
            runs.append(("arrow", arrow_text))
        else:
            text_buf.append(line)
    if text_buf:
        runs.append(("text", '\n'.join(text_buf)))

    cur_y  = y
    max_y  = y + h

    for kind, seg in runs:
        if cur_y + Mm(3) > max_y:
            break

        if kind == "text":
            block_h_mm = _est_text_h_mm(seg, sz, TEXT_W_MM, ls_pt=ls_pt)
            block_h = min(Mm(block_h_mm), max_y - cur_y)
            if block_h > Mm(1.5):
                _txb(slide, seg, l=x, t=cur_y, w=w, h=block_h,
                     sz=sz, color=color, ls_pt=ls_pt)
            cur_y += block_h

        elif kind == "arrow":
            arrow_top = cur_y + max(0, (ROW_H - ARROW_SZ) / 2)
            _brand_arrow(slide, x, arrow_top,
                         size_mm=arrow_size_mm, direction="right", style=arrow_style)
            text_x = x + ARROW_SZ + ARROW_GAP
            text_w = w - ARROW_SZ - ARROW_GAP
            avail  = min(ROW_H, max_y - cur_y)
            if avail > Mm(2):
                _txb(slide, seg, l=text_x, t=cur_y, w=text_w, h=avail,
                     sz=sz, bold=True, color=color, ls_pt=ls_pt)
            cur_y += ROW_H + Mm(1)


def _render_card_inner(slide, layout, x, y, w, h, data, acc_color,
                       txt_color=None, body_color=None, pad_s=None, pad_t=None):
    """
    Render card interior content for non-vertical-stack layouts.
    Caller must draw the card background rect first.

    layout values:
      "horizontal_split" — large metric RIGHT | tag+title+body LEFT
      "icon_left"        — circular badge top-left | title beside | body full-width below
      "quote_card"       — decorative quote mark + quote text + bottom attribution

    x,y,w,h : card bounding box in EMU
    data    : dict; keys vary by layout:
      horizontal_split: "metric", "tag", "title", "body",
                        "metric_color" (default acc_color),
                        "metric_sz"   (default 28, use 44 for big_stats),
                        "title_sz"    (default 12),
                        "body_sz"     (default 11)
      icon_left       : "icon" (1-2 chars for badge), "tag", "title", "body"
      quote_card      : "quote" (or "body"), "attribution"
    acc_color  : accent hex — tag text, icon badge bg, quote mark color
    txt_color  : title color (default NEUTRAL_900)
    body_color : body text color (default NEUTRAL_700)
    pad_s      : side padding in EMU (default Mm(5))
    pad_t      : top  padding in EMU (default Mm(6))
    """
    txt_color  = txt_color  or BT.NEUTRAL_900_HEX
    body_color = body_color or BT.NEUTRAL_700_HEX
    pad_s = pad_s if pad_s is not None else Mm(5)
    pad_t = pad_t if pad_t is not None else Mm(6)
    pad_b = Mm(7)
    inner_l = x + pad_s
    inner_w = w - 2 * pad_s

    # ── horizontal_split ──────────────────────────────────────────────────────
    if layout == "horizontal_split":
        SPLIT_GAP  = Mm(4)
        RIGHT_FRAC = 0.38
        right_w    = int(inner_w * RIGHT_FRAC)
        left_w     = inner_w - right_w - SPLIT_GAP
        right_x    = inner_l + left_w + SPLIT_GAP
        mtr_sz     = data.get("metric_sz",  28)
        title_sz   = data.get("title_sz",   12)
        body_sz    = data.get("body_sz",    11)
        mtr_color  = data.get("metric_color", acc_color)
        metric     = data.get("metric", data.get("val", ""))

        # RIGHT: large metric, right-aligned, anchored near top — no wrap
        mtr_h = Mm(int(mtr_sz * 0.353) + 6)
        if metric:
            _txb(slide, metric,
                 l=right_x, t=y + pad_t, w=right_w, h=mtr_h,
                 sz=mtr_sz, bold=True, align=PP_ALIGN.RIGHT,
                 color=mtr_color, wrap=False)

        # LEFT: tag → title → body (fills remaining height)
        y_off = y + pad_t
        tag = data.get("tag", "")
        if tag:
            _txb(slide, tag, l=inner_l, t=y_off, w=left_w, h=Mm(5),
                 sz=7, bold=True, color=acc_color)
            y_off += Mm(6)
        title = data.get("title", "")
        t_h = Mm(10)   # 1-line headroom for 11-14pt titles
        if title:
            _txb(slide, title, l=inner_l, t=y_off, w=left_w, h=t_h,
                 sz=title_sz, bold=True, color=txt_color)
            y_off += t_h + Mm(2)
        body = data.get("body", "")
        body_h = y + h - y_off - pad_b
        if body and body_h > Mm(4):
            tb = _txb(slide, body, l=inner_l, t=y_off, w=left_w, h=body_h,
                      sz=body_sz, color=body_color, ls_pt=round(body_sz * 1.4))
            tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # ── icon_left ─────────────────────────────────────────────────────────────
    elif layout == "icon_left":
        ICON_D = Mm(20)

        # Draw badge circle (top-left of inner zone)
        _rect(slide, l=inner_l, t=y + pad_t, w=ICON_D, h=ICON_D,
              fill=acc_color, radius_mm=BT.RADIUS_PILL_MM)
        icon_ch = data.get("icon", "")
        if icon_ch:
            # Full badge box + MIDDLE anchor = true vertical center
            _txb(slide, icon_ch[:2].upper(),
                 l=inner_l, t=y + pad_t, w=ICON_D, h=ICON_D,
                 sz=10, bold=True, color=BT.WHITE_HEX,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=False)

        # Content right of icon (title / tag — narrow zone)
        cnt_x = inner_l + ICON_D + Mm(4)
        cnt_w = x + w - cnt_x - pad_s
        y_off = y + pad_t
        tag = data.get("tag", "")
        if tag and cnt_w > Mm(8):
            _txb(slide, tag, l=cnt_x, t=y_off, w=cnt_w, h=Mm(5),
                 sz=7, bold=True, color=acc_color)
            y_off += Mm(6)
        title = data.get("title", "")
        if title and cnt_w > Mm(8):
            _txb(slide, title, l=cnt_x, t=y_off, w=cnt_w, h=Mm(14),
                 sz=13, bold=True, color=txt_color)
            y_off += Mm(15)

        # Body starts below icon zone, spans full inner width
        body_y = max(y_off, y + pad_t + ICON_D + Mm(4))
        body_h = y + h - body_y - pad_b
        body = data.get("body", "")
        if body and body_h > Mm(4):
            tb = _txb(slide, body, l=inner_l, t=body_y, w=inner_w, h=body_h,
                      sz=12, color=body_color, ls_pt=17)
            tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # ── quote_card ────────────────────────────────────────────────────────────
    elif layout == "quote_card":
        # Decorative opening quote mark (decorative, not content)
        _txb(slide, "“",
             l=inner_l, t=y + Mm(4), w=Mm(14), h=Mm(11),
             sz=28, bold=True, color=acc_color)
        quote  = data.get("quote", data.get("body", ""))
        attrib = data.get("attribution", "")
        attrib_h = Mm(10) if attrib else Mm(0)
        q_top    = y + Mm(13)
        q_h      = y + h - q_top - attrib_h - pad_b
        if quote and q_h > Mm(4):
            tb = _txb(slide, quote,
                      l=inner_l + Mm(4), t=q_top,
                      w=inner_w - Mm(4), h=q_h,
                      sz=13, color=body_color, ls_pt=19)
            tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if attrib:
            _txb(slide, f"— {attrib}",
                 l=inner_l, t=y + h - attrib_h - Mm(4), w=inner_w, h=attrib_h,
                 sz=10, color=BT.NEUTRAL_400_HEX, align=PP_ALIGN.RIGHT)


# ── Pill design system ────────────────────────────────────────────────────────

# Named style presets: (bg_hex, fg_hex)
# Use as: bg, fg = _PILL_STYLES["primary"]
_PILL_STYLES: dict = {
    # Filled / bold (bg = brand color, text = white or dark)
    "primary":        (BT.PRIMARY_500_HEX,    BT.WHITE_HEX),
    "secondary":      (BT.SECONDARY_500_HEX,  BT.NEUTRAL_900_HEX),
    "success":        (BT.SUCCESS_HEX,         BT.WHITE_HEX),
    "warning":        (BT.WARNING_HEX,         BT.NEUTRAL_900_HEX),
    "teal":           (BT.TEAL_HEX,            BT.WHITE_HEX),
    "purple":         (BT.PURPLE_HEX,          BT.WHITE_HEX),
    "danger":         (BT.DANGER_HEX,          BT.WHITE_HEX),
    "dark":           (BT.NEUTRAL_900_HEX,     BT.WHITE_HEX),
    # Soft / light (bg = *_100 tint, text = accent)
    "primary-soft":   (BT.PRIMARY_100_HEX,     BT.PRIMARY_500_HEX),
    "secondary-soft": (BT.SECONDARY_100_HEX,   BT.SECONDARY_500_HEX),
    "neutral":        (BT.NEUTRAL_100_HEX,     BT.NEUTRAL_700_HEX),
    "neutral-dark":   (BT.NEUTRAL_200_HEX,     BT.NEUTRAL_700_HEX),
}

# Colors that need dark text for readability (light/warm backgrounds)
_LIGHT_PILL_BGS = frozenset({
    BT.SECONDARY_500_HEX,  # #C8E13C yellow-green
    BT.WARNING_HEX,        # #FFB928 orange-yellow
})


def _pill_fg(bg: str) -> str:
    """Return white or dark foreground text based on bg luminance."""
    return BT.NEUTRAL_900_HEX if bg in _LIGHT_PILL_BGS else BT.WHITE_HEX


def _pill(slide, text: str, l: int, t: int,
          bg: str = None, fg: str = None,
          font_sz: int = 8, bold: bool = True,
          h: int = None, max_w: int = None) -> int:
    """
    Render a single labeled pill at (l, t).
    Returns pill width in EMU.  Use for layout chaining:
        next_x = l + _pill(slide, text, l, t, ...) + Mm(gap)

    bg / fg: hex strings; omit to get PRIMARY filled style.
    h:       pill height in EMU; default Mm(7).
    max_w:   cap pill width (prevents overflow in constrained slots).
    """
    if bg is None:
        bg = BT.PRIMARY_500_HEX
    if fg is None:
        fg = _pill_fg(bg)
    if h is None:
        h = Mm(7)
    PAD_X = Mm(3.5)
    CJK_W = Mm(font_sz * 0.44)     # 8pt ≈ 3.5mm/CJK char; 11pt ≈ 4.8mm
    pw    = int(len(text) * CJK_W + 2 * PAD_X)
    if max_w and pw > max_w:
        pw = max_w
    _rect(slide, l=l, t=t, w=pw, h=h,
          fill=bg, radius_mm=BT.RADIUS_PILL_MM)
    # Full-height box + MIDDLE anchor → text sits in exact vertical center of pill
    _txb(slide, text, l=l + PAD_X, t=t,
         w=pw - 2 * PAD_X, h=h,
         sz=font_sz, bold=bold, color=fg,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=False)
    return pw


# ── Callout / Note blocks ─────────────────────────────────────────────────────

# Presentational gradient stop colors for callout backgrounds.
# Computed as 12% tint of brand color on white — deliberately NOT exported as
# brand tokens (too contextual, too light for general use).
#   Primary  #3EC99E × 12% + white × 88% → #E8F9F3
#   Secondary #C8E13C × 12% + white × 88% → #F8FBE8
_CGRAD_L = "#E8F9F3"
_CGRAD_R = "#F8FBE8"

_CALLOUT_STYLES: dict = {
    # (bg_hex_or_None, pill_style, default_label, grad_stops_or_None)
    # All styles use a horizontal gradient for visual consistency.
    # note/info/tip: green-tinted gradient (primary palette)
    # warning: warm orange → lime-100 gradient — retains semantic colour while matching rhythm
    # danger:  solid red bg (no gradient — severity intentionally flat/stark)
    "note":    (None,              "primary", "注",   [(0, _CGRAD_L), (100, _CGRAD_R)]),
    "info":    (None,              "primary", "说明", [(0, _CGRAD_L), (100, _CGRAD_R)]),
    "tip":     (None,              "primary", "提示", [(0, _CGRAD_R), (100, _CGRAD_L)]),
    "warning": (None,              "warning", "注意", [(0, BT.CARD_ORANGE_BG), (100, BT.SECONDARY_100_HEX)]),
    "danger":  (BT.CARD_DANGER_BG, "danger",  "警告", None),
}

_CALLOUT_LABELS_EN = {
    "note":    "Note",
    "info":    "Info",
    "tip":     "Tip",
    "warning": "Warning",
    "danger":  "Danger",
}

CALLOUT_H = Mm(12)   # standard callout block height


def _callout(slide, text: str, l: int, t: int, w: int,
             label: str = None, style: str = "note",
             font_sz: int = 10) -> int:
    """
    Render a callout / annotation block at (l, t, w).
    Returns height in EMU for layout chaining.

    label:  pill text — None uses the style default ("注"/"说明" etc.); "" suppresses pill
    style:  "note" | "info" | "tip" | "warning" | "danger"
    """
    bg_hex, pill_style, default_label, grad_stops = _CALLOUT_STYLES.get(
        style, _CALLOUT_STYLES["note"])
    # Use English labels when lang="en" and no explicit label is given
    import scripts.i18n as _i18n
    if _i18n.LANG == "en" and label is None:
        default_label = _CALLOUT_LABELS_EN.get(style, "Note")
    PAD   = Mm(5)

    shape = _rect(slide, l=l, t=t, w=w, h=CALLOUT_H,
                  fill=bg_hex or _CGRAD_L, radius_mm=BT.RADIUS_SM_MM)
    if grad_stops:
        _apply_shape_gradient(shape, grad_stops, angle_deg=0)

    pill_label = default_label if label is None else label
    text_x     = l + PAD
    if pill_label:
        _ps    = _PILL_STYLES.get(pill_style, _PILL_STYLES["neutral"])
        pw     = _pill(slide, pill_label, l=l + PAD, t=t + Mm(2.5),
                       bg=_ps[0], fg=_ps[1], font_sz=8)
        text_x = l + PAD + pw + Mm(3)

    _txb(slide, text,
         l=text_x, t=t + Mm(2),
         w=w - (text_x - l) - PAD, h=CALLOUT_H - Mm(4),
         sz=font_sz, color=BT.NEUTRAL_700_HEX, ls_pt=16)

    return CALLOUT_H


def _flow_pills(slide, items: list, x0: int, y0: int, max_w: int,
                pill_h: int = None, font_sz: int = 9,
                pill_bg: str = None, pill_fg: str = None,
                arr_c: str = None) -> int:
    """
    Render items as pill→arrow→pill inline flow.
    Auto-wraps to a new line when items exceed max_w.
    Returns the bottom y-coordinate of the last rendered row.
    """
    if not items:
        return y0
    if pill_h is None:
        pill_h = Mm(8)
    if pill_bg is None:
        pill_bg = BT.PRIMARY_100_HEX
    if pill_fg is None:
        pill_fg = BT.PRIMARY_500_HEX
    if arr_c is None:
        arr_c = BT.NEUTRAL_400_HEX

    ARR_W   = Mm(7)
    ROW_GAP = Mm(3)

    def _pw(text: str) -> int:
        return int(len(text) * Mm(font_sz * 0.44) + 2 * Mm(3.5))

    # build wrapped rows
    rows: list = []
    cur_row: list = []
    cur_w = 0
    for item in items:
        pw    = _pw(item)
        arrow = ARR_W if cur_row else 0
        if cur_row and cur_w + arrow + pw > max_w:
            rows.append(cur_row)
            cur_row = [item]
            cur_w   = pw
        else:
            cur_row.append(item)
            cur_w += arrow + pw
    if cur_row:
        rows.append(cur_row)

    cur_y = y0
    for row in rows:
        cur_x = x0
        for j, item in enumerate(row):
            if j > 0:
                _txb(slide, "→", l=cur_x, t=cur_y, w=ARR_W, h=pill_h,
                     sz=font_sz - 1, color=arr_c, align=PP_ALIGN.CENTER,
                     bold=False, wrap=False)
                cur_x += ARR_W
            cur_x += _pill(slide, item, l=cur_x, t=cur_y,
                           bg=pill_bg, fg=pill_fg, font_sz=font_sz, h=pill_h)
        cur_y += pill_h + ROW_GAP

    return cur_y - ROW_GAP


# ── Container renderers ───────────────────────────────────────────────────────
#
#  Every renderer has the signature:
#      _render_ct_*(slide, slot: dict, l, t, w, h) -> None
#
#  l/t/w/h are the pre-computed pixel coordinates for this slot.
#  The renderer owns everything inside that rectangle.
#  INNER_PAD is applied internally — callers never add extra padding.
#

def _render_ct_placeholder(slide, slot, l, t, w, h):
    """Dashed placeholder box — used during drafting or when a type is TBD."""
    label = slot.get("label", "Placeholder")
    _card(slide, l=l, t=t, w=w, h=h,
          bg=BT.NEUTRAL_100_HEX, border=BT.NEUTRAL_200_HEX)
    _txb(slide, f"[ {label} ]",
         l=l, t=t + (h - Mm(8)) // 2, w=w, h=Mm(8),
         sz=11, color=BT.NEUTRAL_400_HEX, align=PP_ALIGN.CENTER)


def _render_ct_blank(slide, slot, l, t, w, h):
    """Empty spacer — no visual output."""
    pass


def _render_ct_text(slide, slot, l, t, w, h):
    """
    Text / bullet container.

    slot keys:
        label   : str               — optional bold section sub-header at top
        content : str               — free paragraph text
        bullets : list              — str items, or dicts with "pill"/"section"/"text"
                                      (same format as add_body_slide bullets)
        dark    : bool              — black bg + white text (NEUTRAL_900 card, accent label)
    content and bullets are mutually exclusive; if both are given, bullets wins.
    """
    from scripts.i18n import T

    is_dark = slot.get("dark", False)
    if is_dark:
        _card(slide, l=l, t=t, w=w, h=h, bg=BT.NEUTRAL_900_HEX)
        label_color   = BT.SECONDARY_500_HEX
        content_color = BT.NEUTRAL_200_HEX
    else:
        label_color   = BT.NEUTRAL_900_HEX
        content_color = BT.NEUTRAL_700_HEX

    cur_y   = t + INNER_PAD
    avail_w = w - 2 * INNER_PAD
    avail_h = h - 2 * INNER_PAD

    # Optional container sub-header
    raw_label = slot.get("label", "")
    if raw_label:
        _txb(slide, T(raw_label),
             l=l + INNER_PAD, t=cur_y, w=avail_w, h=Mm(8),
             sz=12, bold=True, color=label_color, wrap=False)
        cur_y   += Mm(8) + Mm(4)
        avail_h -= Mm(8) + Mm(4)

    content = slot.get("content", "")
    bullets = slot.get("bullets") or []

    if bullets:
        _has_pills = any(
            isinstance(b, dict) and (b.get("pill") or b.get("section"))
            for b in bullets
        )
        if not _has_pills:
            # Fast path: single textbox for a uniform bullet list
            tb = slide.shapes.add_textbox(
                int(l + INNER_PAD), int(cur_y), int(avail_w), int(avail_h))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                text = T(b) if isinstance(b, str) else T(b.get("text", ""))
                p    = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                pPr  = p._p.get_or_add_pPr()
                lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
                etree.SubElement(lnSpc, qn("a:spcPts"), attrib={"val": "2100"})
                spc_b = etree.SubElement(pPr, qn("a:spcBef"))
                etree.SubElement(spc_b, qn("a:spcPts"), attrib={"val": "700"})
                run  = p.add_run()
                run.text = f"• {text}"
                run.font.size = Pt(15)
                run.font.color.rgb = _rgb(content_color)
                _set_run_fonts(run)
        else:
            # Per-element path: mixed pill / section-header / plain bullets
            BULLET_H = Mm(9)
            B_GAP    = Mm(2.5)
            bx       = l + INNER_PAD
            for b in bullets:
                if isinstance(b, dict) and b.get("section"):
                    cur_y += Mm(3)
                    _rect(slide, l=bx, t=cur_y + Mm(2),
                          w=Mm(3), h=Mm(5),
                          fill=BT.PRIMARY_500_HEX, radius_mm=1.0)
                    _txb(slide, T(b["section"]),
                         l=bx + Mm(5), t=cur_y, w=avail_w - Mm(5), h=BULLET_H,
                         sz=12, bold=True, color=BT.PRIMARY_500_HEX)
                    cur_y += BULLET_H + Mm(3)
                    continue
                if isinstance(b, str):
                    _txb(slide, f"• {T(b)}",
                         l=bx, t=cur_y, w=avail_w, h=BULLET_H,
                         sz=15, color=content_color)
                else:
                    pill_text = b.get("pill", "")
                    body_t    = T(b.get("text", ""))
                    bstyle    = b.get("style", "primary-soft")
                    ps        = _PILL_STYLES.get(bstyle, _PILL_STYLES["primary-soft"])
                    tx        = bx
                    if pill_text:
                        pw = _pill(slide, T(pill_text), l=bx, t=cur_y + Mm(1),
                                   bg=ps[0], fg=ps[1], font_sz=9)
                        tx = bx + pw + Mm(3)
                    if body_t:
                        _txb(slide, body_t, l=tx, t=cur_y,
                             w=avail_w - (tx - bx), h=BULLET_H,
                             sz=14, color=content_color, ls_pt=20)
                cur_y += BULLET_H + B_GAP

    elif content:
        _txb(slide, T(content),
             l=l + INNER_PAD, t=cur_y, w=avail_w, h=avail_h,
             sz=15, color=content_color, ls_pt=22, wrap=True)


# ── Cards container ───────────────────────────────────────────────────────────

# 6-slot default palette — cycles when a card has no explicit "color" field
_SLOT_BG  = [BT.PRIMARY_100_HEX, BT.NEUTRAL_100_HEX, BT.SECONDARY_100_HEX,
             BT.CARD_ORANGE_BG,  BT.CARD_TEAL_BG,    BT.CARD_PURPLE_BG]
_SLOT_ACC = [BT.PRIMARY_500_HEX, BT.SUCCESS_HEX,     BT.SECONDARY_500_HEX,
             BT.WARNING_HEX,     BT.TEAL_HEX,         BT.PURPLE_HEX]


def _card_slot_color(c: dict, slot_idx: int):
    """Return (bg, acc, txt, body_txt) for one card dict at a given slot index."""
    from scripts.components.base import resolve_card_color
    if c.get("dark"):
        return (BT.NEUTRAL_900_HEX, BT.SECONDARY_500_HEX,
                BT.WHITE_HEX,       BT.NEUTRAL_400_HEX)
    if c.get("danger"):
        return (BT.CARD_DANGER_BG, BT.DANGER_HEX,
                BT.NEUTRAL_900_HEX, BT.NEUTRAL_700_HEX)
    if c.get("color"):
        bg, acc, txt = resolve_card_color(c["color"])
        return (bg, acc, txt, BT.NEUTRAL_700_HEX)
    bg  = _SLOT_BG [slot_idx % 6]
    acc = _SLOT_ACC[slot_idx % 6]
    return (bg, acc, BT.NEUTRAL_900_HEX, BT.NEUTRAL_700_HEX)


def _draw_card_vs(slide, c: dict, x, y, cw, ch, acc, txt_color, body_color):
    """Render one vertical_stack card with typography scaled to card height."""
    from scripts.i18n import T
    compact  = ch < Mm(70)
    pad_s    = Mm(4)  if compact else Mm(5)
    pad_t    = Mm(5)  if compact else Mm(6)
    title_sz = 13     if compact else 16
    title_h  = Mm(16) if compact else Mm(24)
    title_adv= Mm(17) if compact else Mm(25)
    bar_w    = Mm(20) if compact else Mm(28)
    bar_h_v  = Mm(0.8)if compact else Mm(1)
    bar_adv  = Mm(4)  if compact else Mm(5)
    body_sz  = 11     if compact else 13
    tag_sz   = 8      if compact else 9

    inner_l = x + pad_s
    inner_w = cw - 2 * pad_s
    y_off   = y + pad_t

    tag = T(c.get("tag", ""))
    if tag:
        _txb(slide, tag, l=inner_l, t=y_off, w=inner_w, h=Mm(6),
             sz=tag_sz, bold=True, color=acc)
        y_off += Mm(7)

    _txb(slide, T(c.get("title", "")),
         l=inner_l, t=y_off, w=inner_w, h=title_h,
         sz=title_sz, bold=True, color=txt_color)
    y_off += title_adv

    _rect(slide, l=inner_l, t=y_off, w=bar_w, h=bar_h_v, fill=acc)
    y_off += bar_h_v + bar_adv

    body   = T(c.get("body", ""))
    body_h = ch - (y_off - y) - Mm(8)
    if body and body_h > Mm(6):
        tb = _txb(slide, body, l=inner_l, t=y_off, w=inner_w, h=body_h,
                  sz=body_sz, color=body_color, ls_pt=16)
        tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _render_ct_cards(slide, slot: dict, l, t, w, h):
    """
    Card grid container.  Auto-selects column/row layout by item count.

    slot keys:
        items       : list of card dicts (title, body, tag, color, dark, layout)
        intro_text  : str  — 5-card mode: replaces top-left slot with summary text
        intro_label : str  — small green label above intro_text
        intro_flow  : list — pill→arrow→pill chain in the intro slot

    Auto layout rules (per CLAUDE.md):
        1-3 items  →  1 row × n columns   (three-card style, tall cards)
        4 items    →  1 row × 4 columns   (four-card style, tighter gap)
        5 items    →  intro slot + 2 row × 3 col  (requires intro_text or intro_flow)
        6 items    →  2 rows × 3 columns  (six-card style, seq badges)
    """
    items       = slot.get("items") or []
    n           = len(items)
    intro_text  = slot.get("intro_text",  "")
    intro_label = slot.get("intro_label", "")
    intro_flow  = slot.get("intro_flow")  or []
    if n == 0:
        return

    # ── Dark card budget (CLAUDE.md: 0 if <5 cards; max 1 if ≥5) ────────────
    dark_budget = (BT.MAX_DARK_PER_SLIDE if n >= BT.MIN_CARDS_FOR_DARK else 0
                   ) if BT.DARK_ACCENT_CARDS_ENABLED else 0
    dark_used, adj = 0, []
    for c in items:
        if c.get("dark") and dark_used < dark_budget:
            adj.append(c); dark_used += 1
        elif c.get("dark"):
            adj.append({k: v for k, v in c.items() if k != "dark"})
        else:
            adj.append(c)

    # ── Grid dimensions ───────────────────────────────────────────────────────
    use_intro   = (bool(intro_text) or bool(intro_flow)) and n == 5
    slot_offset = 1 if use_intro else 0

    if n <= 3:
        n_cols, n_rows, col_gap = n, 1, COL_GAP
    elif n == 4:
        n_cols, n_rows, col_gap = 4, 1, Mm(4)
    else:                         # 5 or 6
        n_cols, n_rows, col_gap = 3, 2, C3_GAP

    PAD_T    = Mm(4)
    cards_t  = t + PAD_T
    cards_h  = h - PAD_T
    card_w   = (w - (n_cols - 1) * col_gap) // n_cols
    card_h   = ((cards_h - ROW_GAP) // 2) if n_rows == 2 else cards_h
    show_seq = (n_rows == 2)   # seq badges only on 6-card grid

    # ── Intro slot (5-card balanced layout) ───────────────────────────────────
    if use_intro:
        _pad  = Mm(5)
        _iw   = card_w - 2 * _pad
        cur_y = cards_t + _pad
        if intro_label:
            _txb(slide, intro_label, l=l + _pad, t=cur_y, w=_iw, h=Mm(6),
                 sz=8, bold=True, color=BT.PRIMARY_500_HEX)
            cur_y += Mm(8)
        if intro_text:
            _txb(slide, intro_text, l=l + _pad, t=cur_y, w=_iw, h=Mm(20),
                 sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=17)
            cur_y += Mm(22)
        if intro_flow:
            _flow_pills(slide, intro_flow,
                        x0=l + _pad, y0=cur_y, max_w=_iw, font_sz=9)

    # ── Render cards ──────────────────────────────────────────────────────────
    for i, c in enumerate(adj):
        si   = i + slot_offset
        col  = si % n_cols
        row  = si // n_cols
        cx   = l + col * (card_w + col_gap)
        cy   = cards_t + row * (card_h + ROW_GAP)
        # Last cell absorbs rounding so cards reach exactly the slot edge
        cw_i = (l + w - cx)       if col == n_cols - 1 else card_w
        ch_i = (cards_t + cards_h - cy) if row == n_rows - 1 else card_h

        bg, acc, txt_color, body_color = _card_slot_color(c, si)
        _card(slide, l=cx, t=cy, w=cw_i, h=ch_i, bg=bg)

        if show_seq:
            _seq_badge(slide, x=cx + cw_i - Mm(13), y=cy + Mm(4),
                       seq=i + 1, color=acc)

        inner_layout = c.get("layout", "vertical_stack")
        if inner_layout != "vertical_stack":
            _render_card_inner(slide, inner_layout, cx, cy, cw_i, ch_i,
                               c, acc, txt_color=txt_color,
                               body_color=body_color,
                               pad_s=Mm(4), pad_t=Mm(5))
        else:
            _draw_card_vs(slide, c, cx, cy, cw_i, ch_i, acc, txt_color, body_color)


# ── Strips container ──────────────────────────────────────────────────────────

def _draw_strip_item(slide, item: dict, style: str,
                     strip_l, cur_y, strip_w, strip_h, idx: int):
    """
    Render one strip card at the given coordinates.

    Extracted so that add_strip_list, add_two_panel_strips, and _render_ct_strips
    all call the same implementation.

    style: "accent" | "numbered" | "tag_left"
    """
    from scripts.components.base import resolve_card_color
    from scripts.i18n import T

    item_title = T(item.get("title", ""))
    item_body  = T(item.get("body",  ""))
    item_tag   = T(item.get("tag",   ""))
    color_key  = item.get("color", "primary")
    num_label  = item.get("num", f"{idx + 1:02d}")
    bg, acc, txt = resolve_card_color(color_key)

    is_dark_strip = item.get("dark", False)
    if is_dark_strip:
        strip_bg  = BT.NEUTRAL_900_HEX
        strip_txt = BT.WHITE_HEX
        strip_acc = BT.SECONDARY_500_HEX   # lime accent for dark strip
        strip_body_color = BT.NEUTRAL_200_HEX
    else:
        strip_bg       = bg
        strip_txt      = txt
        strip_acc      = acc
        strip_body_color = BT.NEUTRAL_700_HEX

    if style == "accent":
        _card(slide, l=strip_l, t=cur_y, w=strip_w, h=strip_h, bg=strip_bg)

        # Left icon/number indicator — decorative bold number/icon in accent color
        # Alibaba PuHuiTi 2.0 for clean sans-serif digits; vertically centered
        ICON_W   = Mm(16)
        icon_str = T(item.get("icon", "")) or num_label
        _txb(slide, icon_str,
             l=strip_l + Mm(2),
             t=cur_y,
             w=ICON_W, h=strip_h,
             sz=20, bold=True, color=strip_acc,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=False,
             en_font=BT.FONT_CN, cn_font=BT.FONT_CN)

        PAD_L     = Mm(2) + ICON_W + Mm(4)   # shift content right of icon
        PAD_T     = Mm(4)
        PAD_R     = Mm(6)
        content_w = strip_w - PAD_L - PAD_R

        tag_reserve = 0
        if item_tag:
            est_tag_w = int(len(item_tag) * Mm(8 * 0.44) + 2 * Mm(3.5))
            tag_t     = cur_y + (strip_h - Mm(8)) // 2
            tag_l     = strip_l + strip_w - PAD_R - est_tag_w
            _pill(slide, item_tag, l=tag_l, t=tag_t,
                  bg=strip_acc, fg=_pill_fg(strip_acc), font_sz=8, h=Mm(8))
            tag_reserve = est_tag_w + Mm(4)

        title_h = Mm(8)
        body_h  = strip_h - PAD_T - title_h - Mm(2) - Mm(3)

        if item_title:
            _txb(slide, item_title,
                 l=strip_l + PAD_L, t=cur_y + PAD_T,
                 w=content_w - tag_reserve, h=title_h,
                 sz=13, bold=True, color=strip_txt)
        if item_body and body_h > Mm(4):
            tb = _txb(slide, item_body,
                      l=strip_l + PAD_L, t=cur_y + PAD_T + title_h + Mm(2),
                      w=content_w, h=body_h,
                      sz=11, color=strip_body_color, ls_pt=16)
            try:
                tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass

    elif style == "numbered":
        BADGE_D  = Mm(10) if strip_h < Mm(20) else Mm(12)
        badge_sz = max(7, int(BADGE_D / Mm(1)) - 2)
        _card(slide, l=strip_l, t=cur_y, w=strip_w, h=strip_h,
              bg=BT.WHITE_HEX, border=BT.NEUTRAL_200_HEX)
        badge_l  = strip_l + Mm(5)
        badge_t  = cur_y + (strip_h - BADGE_D) // 2
        _rect(slide, l=badge_l, t=badge_t, w=BADGE_D, h=BADGE_D,
              fill=acc, radius_mm=BT.RADIUS_PILL_MM)
        # Full badge box + MIDDLE anchor = true vertical center
        _txb(slide, num_label,
             l=badge_l, t=badge_t, w=BADGE_D, h=BADGE_D,
             sz=badge_sz, bold=True, color=BT.WHITE_HEX,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=False,
             en_font=BT.FONT_CN, cn_font=BT.FONT_CN)

        CONTENT_L = badge_l + BADGE_D + Mm(5)
        content_w = strip_l + strip_w - CONTENT_L - Mm(5)
        title_h   = Mm(7)
        body_h    = strip_h - Mm(4) - title_h - Mm(2) - Mm(3)
        y_title   = (cur_y + int((strip_h - title_h) // 2)
                     if not item_body or body_h < Mm(4)
                     else cur_y + Mm(3))

        if item_title:
            _txb(slide, item_title,
                 l=CONTENT_L, t=y_title, w=content_w, h=title_h,
                 sz=13, bold=True, color=BT.NEUTRAL_900_HEX)
        if item_body and body_h > Mm(4):
            tb = _txb(slide, item_body,
                      l=CONTENT_L, t=y_title + title_h + Mm(2),
                      w=content_w, h=body_h,
                      sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=16)
            try:
                tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass

    elif style == "tag_left":
        _card(slide, l=strip_l, t=cur_y, w=strip_w, h=strip_h, bg=bg)
        TAG_PAD   = Mm(5)
        tag_label = item_tag or (item_title[:10] if item_title else "")
        est_tag_w = int(len(tag_label) * Mm(9 * 0.44) + 2 * Mm(4))
        tag_t     = cur_y + (strip_h - Mm(8)) // 2
        _pill(slide, tag_label, l=strip_l + TAG_PAD, t=tag_t,
              bg=acc, fg=_pill_fg(acc), font_sz=9, h=Mm(8))

        CONTENT_L = strip_l + TAG_PAD + est_tag_w + Mm(5)
        content_w = strip_l + strip_w - CONTENT_L - Mm(5)
        title_h   = Mm(7)
        body_h    = strip_h - Mm(4) - title_h - Mm(2) - Mm(3)
        y_title   = (cur_y + int((strip_h - title_h) // 2)
                     if not item_body or body_h < Mm(4)
                     else cur_y + Mm(3))

        if item_title:
            _txb(slide, item_title,
                 l=CONTENT_L, t=y_title, w=content_w, h=title_h,
                 sz=13, bold=True, color=txt)
        if item_body and body_h > Mm(4):
            tb = _txb(slide, item_body,
                      l=CONTENT_L, t=y_title + title_h + Mm(2),
                      w=content_w, h=body_h,
                      sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=16)
            try:
                tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass


def _render_ct_strips(slide, slot: dict, l, t, w, h):
    """
    Strip list container — stacked horizontal strip cards filling the slot.

    slot keys:
        style     : "accent" | "numbered" | "tag_left"   (default "accent")
        items     : list of strip dicts (title, body, tag, color, num)
        label     : str — optional bold section sub-header rendered at top
        panel_bg  : False (default) — no bg
                    True            — NEUTRAL_100 tray background
                    "bordered"      — WHITE card + NEUTRAL_200 border
    """
    from scripts.i18n import T

    style     = slot.get("style",    "accent")
    items     = slot.get("items")    or []
    raw_label = slot.get("label",    "")
    panel_bg  = slot.get("panel_bg", False)
    n         = len(items)
    if n == 0:
        return

    # ── Panel background (optional visual grouping) ───────────────────────────
    if panel_bg == "bordered":
        _card(slide, l=l, t=t, w=w, h=h,
              bg=BT.WHITE_HEX, border=BT.NEUTRAL_200_HEX)
    elif panel_bg:
        _card(slide, l=l, t=t, w=w, h=h, bg=BT.NEUTRAL_100_HEX)

    # Inner padding: generous when panel_bg is present, minimal otherwise
    PAD = INNER_PAD if panel_bg else Mm(2)

    cur_y   = t + PAD
    avail_w = w - 2 * PAD
    avail_h = h - 2 * PAD

    # ── Optional section sub-header ───────────────────────────────────────────
    if raw_label:
        _txb(slide, T(raw_label),
             l=l + PAD, t=cur_y, w=avail_w, h=Mm(8),
             sz=12, bold=True, color=BT.NEUTRAL_900_HEX, wrap=False)
        cur_y   += Mm(8) + Mm(4)
        avail_h -= Mm(8) + Mm(4)

    # ── Strip height: fill available height evenly ────────────────────────────
    STRIP_GAP = Mm(3)
    strip_h   = int((avail_h - (n - 1) * STRIP_GAP) // n)

    for idx, item in enumerate(items):
        _draw_strip_item(slide, item, style, l + PAD, cur_y, avail_w, strip_h, idx)
        cur_y += strip_h + STRIP_GAP


# ── Stats container ───────────────────────────────────────────────────────────

def _draw_stat_card(slide, item: dict, cx, cy, cw, ch,
                    acc, txt_color, body_color, n_cols: int,
                    card_layout: str):
    """
    Render one stat card with bilingual-aware layout.

    horizontal_split (default):
      Right zone  — large metric number (Alibaba PuHuiTi, accent colour)
      Left zone   — [cn_label + en_label (bilingual)] bottom-aligned with metric
                    [cn_desc  + en_desc  (bilingual)] below, full inner width

    vertical_stack:
      metric  → cn_label [+ en_label] → separator bar → cn_desc [+ en_desc]
    """
    from scripts.i18n import T, get_cn, get_en, is_bilingual

    bilingual = is_bilingual()
    PAD_S = Mm(5)
    PAD_T = Mm(6)
    PAD_B = Mm(5)

    raw_label = item.get("label", "")
    raw_desc  = item.get("desc",  "")
    val       = str(item.get("val", ""))

    cn_label = get_cn(raw_label) if bilingual and isinstance(raw_label, dict) else T(raw_label)
    en_label = get_en(raw_label) if bilingual and isinstance(raw_label, dict) else ""
    cn_desc  = get_cn(raw_desc)  if bilingual and isinstance(raw_desc, dict)  else T(raw_desc)
    en_desc  = get_en(raw_desc)  if bilingual and isinstance(raw_desc, dict)  else ""

    inner_l = cx + PAD_S
    inner_w = cw - 2 * PAD_S

    # ── horizontal_split ──────────────────────────────────────────────────────
    if card_layout == "horizontal_split":
        # Metric font size scales with number of columns
        METRIC_SZ   = {1: 52, 2: 44, 3: 36}.get(n_cols, 40)
        # Tight box ≈ exact cap-height; MIDDLE anchor centers text precisely
        METRIC_H    = Mm(int(METRIC_SZ * 0.353) + 1)

        # Title group heights
        LABEL_SZ    = 14
        LABEL_H     = Mm(7)
        EN_LABEL_SZ = 10
        EN_LABEL_H  = Mm(5)  if en_label else Mm(0)
        LBL_GAP     = Mm(2)  if en_label else Mm(0)
        TITLE_GH    = LABEL_H + LBL_GAP + EN_LABEL_H

        # Header row: taller of metric or title group
        HEADER_H    = max(METRIC_H, TITLE_GH)

        # Right zone width: fraction of inner width
        RIGHT_FRAC  = {1: 0.22, 2: 0.32, 3: 0.38}.get(n_cols, 0.32)
        right_w     = int(inner_w * RIGHT_FRAC)
        left_w      = inner_w - right_w - Mm(4)
        right_x     = inner_l + left_w + Mm(4)

        header_bot  = cy + PAD_T + HEADER_H

        # ── Metric (right zone, bottom-aligned with header) ────────────────
        metric_t = int(header_bot - METRIC_H)
        if val:
            _txb(slide, val,
                 l=right_x, t=metric_t, w=right_w, h=METRIC_H,
                 sz=METRIC_SZ, bold=True, color=acc,
                 align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, wrap=False,
                 en_font=BT.FONT_CN, cn_font=BT.FONT_CN)

        # ── Title group (left zone, bottom-aligned with metric) ────────────
        lbl_top = int(header_bot - TITLE_GH)
        if cn_label:
            _txb(slide, cn_label,
                 l=inner_l, t=lbl_top, w=left_w, h=LABEL_H,
                 sz=LABEL_SZ, bold=True, color=txt_color)
        if en_label:
            _txb(slide, en_label,
                 l=inner_l, t=int(lbl_top + LABEL_H + LBL_GAP), w=left_w, h=EN_LABEL_H,
                 sz=EN_LABEL_SZ, color=BT.NEUTRAL_400_HEX)

        # ── Description block (below header, full inner width) ─────────────
        desc_t    = int(header_bot + Mm(4))
        total_desc_h = cy + ch - PAD_B - desc_t
        if total_desc_h < Mm(4):
            return

        if cn_desc and en_desc:
            cn_desc_h = int(total_desc_h * 0.56)
            en_desc_h = total_desc_h - cn_desc_h - Mm(2)
        else:
            cn_desc_h = total_desc_h
            en_desc_h = Mm(0)

        if cn_desc and cn_desc_h > Mm(3):
            tb = _txb(slide, cn_desc,
                      l=inner_l, t=desc_t, w=inner_w, h=cn_desc_h,
                      sz=11, color=body_color, ls_pt=16)
            try: tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception: pass

        if en_desc and en_desc_h > Mm(3):
            tb = _txb(slide, en_desc,
                      l=inner_l, t=int(desc_t + cn_desc_h + Mm(2)),
                      w=inner_w, h=en_desc_h,
                      sz=9, color=BT.NEUTRAL_400_HEX, ls_pt=13)
            try: tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception: pass

    # ── vertical_stack ────────────────────────────────────────────────────────
    else:
        METRIC_SZ = {1: 52, 2: 44, 3: 36}.get(n_cols, 40)
        METRIC_H  = Mm(int(METRIC_SZ * 0.353) + 1)
        LABEL_SZ  = 13

        y_off = cy + PAD_T
        if val:
            _txb(slide, val,
                 l=inner_l, t=y_off, w=inner_w, h=METRIC_H,
                 sz=METRIC_SZ, bold=True, color=acc,
                 valign=MSO_ANCHOR.MIDDLE, wrap=False,
                 en_font=BT.FONT_CN, cn_font=BT.FONT_CN)
            y_off += METRIC_H + Mm(3)

        if cn_label:
            _txb(slide, cn_label,
                 l=inner_l, t=y_off, w=inner_w, h=Mm(7),
                 sz=LABEL_SZ, bold=True, color=txt_color)
            y_off += Mm(7)
        if en_label:
            _txb(slide, en_label,
                 l=inner_l, t=y_off, w=inner_w, h=Mm(5),
                 sz=10, color=BT.NEUTRAL_400_HEX)
            y_off += Mm(5)

        # Separator bar
        y_off += Mm(3)
        _rect(slide, l=inner_l, t=y_off, w=Mm(20), h=Mm(1), fill=acc)
        y_off += Mm(4)

        desc_h = cy + ch - PAD_B - y_off
        if cn_desc and desc_h > Mm(4):
            if en_desc:
                cn_h = int(desc_h * 0.56)
                en_h = desc_h - cn_h - Mm(2)
            else:
                cn_h, en_h = desc_h, Mm(0)
            tb = _txb(slide, cn_desc,
                      l=inner_l, t=y_off, w=inner_w, h=cn_h,
                      sz=10, color=body_color, ls_pt=15)
            try: tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception: pass
            if en_desc and en_h > Mm(3):
                tb = _txb(slide, en_desc,
                          l=inner_l, t=int(y_off + cn_h + Mm(2)), w=inner_w, h=en_h,
                          sz=9, color=BT.NEUTRAL_400_HEX, ls_pt=13)
                try: tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception: pass


def _render_ct_stats(slide, slot: dict, l, t, w, h):
    """
    KPI metric cards filling the slot.

    slot keys:
        items       : list of stat dicts
                        val   : str   — the big number/symbol
                        label : str | {"cn":..,"en":..} — metric label
                        desc  : str | {"cn":..,"en":..} — explanation text
                        color : str   — palette key
        card_layout : "horizontal_split" (default) | "vertical_stack"

    Auto grid:
        1 item  → full width, single card
        2 items → 1 row × 2 cols
        3 items → 1 row × 3 cols
        4 items → 2 rows × 2 cols

    Bilingual mode (configure(lang="bilingual")):
        horizontal_split — cn_label + en_label form a title group bottom-aligned
                           with the metric; cn_desc + en_desc fill below.
        vertical_stack   — stacked: metric → cn_label → en_label → bar → cn_desc → en_desc
    """
    items       = slot.get("items") or []
    card_layout = slot.get("card_layout", "horizontal_split")
    n           = len(items)
    if n == 0:
        return

    if n <= 3:
        n_cols, n_rows = n, 1
    else:
        n_cols, n_rows = 2, 2

    cw  = (w - (n_cols - 1) * COL_GAP) // n_cols
    ch  = (h - (n_rows - 1) * ROW_GAP) // n_rows if n_rows > 1 else h

    for i, item in enumerate(items):
        col  = i % n_cols
        row  = i // n_cols
        cx   = l + col * (cw + COL_GAP)
        cy   = t + row * (ch + ROW_GAP)
        cw_i = (l + w - cx) if col == n_cols - 1 else cw
        ch_i = (t + h - cy) if row == n_rows - 1 else ch

        bg, acc, txt, body_color = _card_slot_color(item, i)
        _card(slide, l=cx, t=cy, w=cw_i, h=ch_i, bg=bg)
        _draw_stat_card(slide, item, cx, cy, cw_i, ch_i,
                        acc, txt, body_color, n_cols, card_layout)


# ── Flow container ────────────────────────────────────────────────────────────

def _render_ct_flow(slide, slot: dict, l, t, w, h):
    """
    Process / workflow steps.

    slot keys:
        items     : list of str or {"label": ..., "desc": ...}
        style     : "pills"    — pill→arrow→pill chain  (default)
                    "numbered" — numbered circles with label+desc
        direction : "horizontal" (default) | "vertical"
    """
    from scripts.i18n import T

    items     = slot.get("items") or []
    style     = slot.get("style",     "pills")
    direction = slot.get("direction", "horizontal")
    n         = len(items)
    if n == 0:
        return

    def _item_label(item):
        return T(item) if isinstance(item, (str, dict)) and not isinstance(item, dict) \
               else T(item.get("label", ""))

    def _item_desc(item):
        return "" if isinstance(item, str) else T(item.get("desc", ""))

    # ── pills ─────────────────────────────────────────────────────────────────
    if style == "pills":
        labels = [_item_label(it) for it in items]

        if direction == "horizontal":
            # Use _flow_pills — auto-wrap if needed
            _flow_pills(slide, labels,
                        x0=l + INNER_PAD, y0=t + INNER_PAD,
                        max_w=w - 2 * INNER_PAD,
                        pill_h=Mm(9), font_sz=10,
                        pill_bg=BT.PRIMARY_100_HEX,
                        pill_fg=BT.PRIMARY_500_HEX,
                        arr_c=BT.NEUTRAL_400_HEX)
        else:
            # Vertical: stacked pills + down arrows
            PILL_H    = Mm(9)
            ARR_H     = Mm(7)
            ITEM_H    = PILL_H + ARR_H
            pill_w    = w - 2 * INNER_PAD
            cur_y     = t + INNER_PAD
            for i, lbl in enumerate(labels):
                est_pw = min(pill_w, int(len(lbl) * Mm(10 * 0.44) + 2 * Mm(5)))
                _pill(slide, lbl,
                      l=l + INNER_PAD, t=cur_y,
                      bg=BT.PRIMARY_100_HEX, fg=BT.PRIMARY_500_HEX,
                      font_sz=10, h=PILL_H, w=est_pw)
                cur_y += PILL_H
                if i < n - 1:
                    _txb(slide, "↓",
                         l=l + INNER_PAD, t=cur_y,
                         w=Mm(10), h=ARR_H,
                         sz=11, color=BT.NEUTRAL_400_HEX,
                         align=PP_ALIGN.CENTER)
                    cur_y += ARR_H

    # ── numbered ──────────────────────────────────────────────────────────────
    elif style == "numbered":
        BADGE_D = Mm(12)

        if direction == "horizontal":
            # Horizontal: evenly spaced columns, badge on top + title + desc below
            item_w  = (w - (n - 1) * COL_GAP) // n
            item_h  = h
            ACC_COLORS = [_SLOT_ACC[i % 6] for i in range(n)]

            for i, item in enumerate(items):
                ix = l + i * (item_w + COL_GAP)
                # Badge circle
                badge_l = ix + (item_w - BADGE_D) // 2
                _rect(slide, l=badge_l, t=t, w=BADGE_D, h=BADGE_D,
                      fill=ACC_COLORS[i], radius_mm=BT.RADIUS_PILL_MM)
                # Full badge box + MIDDLE anchor = true vertical center
                _txb(slide, f"{i+1:02d}",
                     l=badge_l, t=t, w=BADGE_D, h=BADGE_D,
                     sz=10, bold=True, color=BT.WHITE_HEX,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=False,
                     en_font=BT.FONT_CN, cn_font=BT.FONT_CN)

                # Arrow connector between badges (except last)
                if i < n - 1:
                    ax = ix + item_w
                    _txb(slide, "→",
                         l=ax, t=t, w=COL_GAP, h=BADGE_D,
                         sz=9, color=BT.NEUTRAL_400_HEX,
                         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=False)

                label_y  = t + BADGE_D + Mm(3)
                label_h  = Mm(8)
                desc_y   = label_y + label_h + Mm(1)
                desc_h   = item_h - BADGE_D - Mm(3) - label_h - Mm(1)

                if _item_label(item):
                    _txb(slide, _item_label(item),
                         l=ix, t=label_y, w=item_w, h=label_h,
                         sz=12, bold=True, color=BT.NEUTRAL_900_HEX,
                         align=PP_ALIGN.CENTER)
                if _item_desc(item) and desc_h > Mm(4):
                    tb = _txb(slide, _item_desc(item),
                              l=ix, t=desc_y, w=item_w, h=desc_h,
                              sz=10, color=BT.NEUTRAL_700_HEX,
                              align=PP_ALIGN.CENTER, ls_pt=15)
                    try:
                        tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    except Exception:
                        pass

        else:
            # Vertical: badge + label on each row, connector line between badges
            ITEM_GAP = Mm(4)
            item_h   = (h - (n - 1) * ITEM_GAP) // n
            ACC_COLORS = [_SLOT_ACC[i % 6] for i in range(n)]

            for i, item in enumerate(items):
                iy = t + i * (item_h + ITEM_GAP)
                cx = l + BADGE_D + Mm(5)
                cw = w - BADGE_D - Mm(5)
                label_h    = Mm(7)
                desc_avail = item_h - BADGE_D - Mm(2)
                has_desc   = bool(_item_desc(item)) and desc_avail > Mm(4)

                # Always top-align: badge and label both start at item top
                badge_t = iy
                lbl_y   = iy + int((BADGE_D - label_h) // 2)   # vertically center label within badge height
                if has_desc:
                    desc_y = iy + BADGE_D + Mm(2)
                    desc_h = iy + item_h - desc_y
                else:
                    desc_y = desc_h = 0

                _rect(slide, l=l, t=badge_t, w=BADGE_D, h=BADGE_D,
                      fill=ACC_COLORS[i], radius_mm=BT.RADIUS_PILL_MM)
                _txb(slide, f"{i+1:02d}",
                     l=l, t=badge_t, w=BADGE_D, h=BADGE_D,
                     sz=10, bold=True, color=BT.WHITE_HEX,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=False,
                     en_font=BT.FONT_CN, cn_font=BT.FONT_CN)

                # Connector: badge bottom → next badge top (formula valid for both cases)
                if i < n - 1:
                    line_t = badge_t + BADGE_D
                    line_h = item_h + ITEM_GAP - BADGE_D
                    _rect(slide,
                          l=l + (BADGE_D - Mm(1)) // 2,
                          t=line_t, w=Mm(1), h=line_h,
                          fill=BT.NEUTRAL_200_HEX)

                if _item_label(item):
                    _txb(slide, _item_label(item),
                         l=cx, t=lbl_y, w=cw, h=label_h,
                         sz=13, bold=True, color=BT.NEUTRAL_900_HEX)
                if has_desc:
                    tb = _txb(slide, _item_desc(item),
                              l=cx, t=desc_y, w=cw, h=desc_h,
                              sz=10, color=BT.NEUTRAL_700_HEX, ls_pt=15)
                    try:
                        tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    except Exception:
                        pass


# ── Image container ───────────────────────────────────────────────────────────

def _render_ct_image(slide, slot: dict, l, t, w, h):
    """
    Image display container.

    slot keys:
        path    : str  — file path (absolute, or relative to CWD)
        fit     : "contain" (default) | "cover"
        caption : str  — optional text shown below the image
        dark    : bool — dark caption text background
    """
    from scripts.i18n import T
    import os

    path    = slot.get("path", "")
    fit     = slot.get("fit",  "contain")
    caption = T(slot.get("caption", ""))
    is_dark = slot.get("dark", False)

    CAP_H = Mm(8) if caption else Mm(0)
    img_h = h - CAP_H
    img_w = w

    if not path or not os.path.exists(path):
        _render_ct_placeholder(slide,
                               {"label": slot.get("label", path or "Image")},
                               l, t, w, h)
        return

    try:
        if fit == "contain":
            # Add with only width first to get auto height
            pic = slide.shapes.add_picture(path,
                                           int(l), int(t),
                                           width=int(img_w))
            if pic.height > img_h:
                # Too tall — redo with height constraint
                sp = pic._element
                sp.getparent().remove(sp)
                pic = slide.shapes.add_picture(path,
                                               int(l), int(t),
                                               height=int(img_h))
            # Center horizontally and vertically in the slot
            pic.left = int(l + (img_w - pic.width)  // 2)
            pic.top  = int(t + (img_h - pic.height) // 2)

        else:  # cover — fill slot, crop via clipping in position
            pic = slide.shapes.add_picture(path,
                                           int(l), int(t),
                                           width=int(img_w),
                                           height=int(img_h))
    except Exception:
        _render_ct_placeholder(slide,
                               {"label": f"⚠ {os.path.basename(path)}"},
                               l, t, w, h)
        return

    if caption:
        cap_t = t + h - CAP_H
        if is_dark:
            _card(slide, l=l, t=cap_t, w=w, h=CAP_H, bg=BT.NEUTRAL_900_HEX)
            cap_color = BT.WHITE_HEX
        else:
            cap_color = BT.NEUTRAL_400_HEX
        _txb(slide, caption,
             l=l + INNER_PAD, t=cap_t + Mm(1),
             w=w - 2 * INNER_PAD, h=CAP_H - Mm(2),
             sz=9, color=cap_color, wrap=False)


# ── Tags container ────────────────────────────────────────────────────────────

def _render_ct_tags(slide, slot: dict, l, t, w, h):
    """
    Horizontally wrapping pill tag cloud.

    slot keys:
        items  : list of str or {"text": ..., "style": "primary"|"success"|...}
        label  : str — optional bold section header above tags
        style  : global default pill style (overridden per-item)
    """
    from scripts.i18n import T

    items        = slot.get("items") or []
    raw_label    = slot.get("label", "")
    default_style = slot.get("style", "primary-soft")
    if not items:
        return

    cur_y   = t + INNER_PAD
    avail_w = w - 2 * INNER_PAD

    if raw_label:
        _txb(slide, T(raw_label),
             l=l + INNER_PAD, t=cur_y, w=avail_w, h=Mm(7),
             sz=11, bold=True, color=BT.NEUTRAL_900_HEX, wrap=False)
        cur_y += Mm(7) + Mm(4)

    PILL_H = Mm(9)
    H_GAP  = Mm(3)
    V_GAP  = Mm(3)
    px     = l + INNER_PAD

    for item in items:
        if isinstance(item, str):
            text  = T(item)
            pstyle = default_style
        else:
            text  = T(item.get("text", ""))
            pstyle = item.get("style", default_style)

        ps = _PILL_STYLES.get(pstyle, _PILL_STYLES.get("primary-soft",
             (BT.PRIMARY_100_HEX, BT.PRIMARY_500_HEX)))
        est_pw = int(len(text) * Mm(10 * 0.44) + 2 * Mm(4.5))

        # Wrap to next row if needed
        if px + est_pw > l + w - INNER_PAD and px > l + INNER_PAD:
            px     = l + INNER_PAD
            cur_y += PILL_H + V_GAP

        if cur_y + PILL_H > t + h:
            break

        pw = _pill(slide, text, l=px, t=cur_y,
                   bg=ps[0], fg=ps[1], font_sz=10, h=PILL_H)
        px += pw + H_GAP


def _render_ct_table(slide, slot: dict, l, t, w, h):
    """
    Compact data table container.

    slot keys:
        headers : list[str]          — column headers (optional)
        rows    : list[list[str]]    — data rows
        label   : str                — small section label above table
        col_widths : list[float]     — relative column widths, e.g. [2,1,1]
    """
    from scripts.i18n import T

    headers    = slot.get("headers") or []
    rows       = slot.get("rows") or []
    raw_label  = slot.get("label", "")
    col_ratios = slot.get("col_widths") or []

    cur_y   = t + INNER_PAD
    avail_w = w - 2 * INNER_PAD

    if raw_label:
        _txb(slide, T(raw_label),
             l=l + INNER_PAD, t=cur_y, w=avail_w, h=Mm(6),
             sz=9, bold=True, color=BT.PRIMARY_500_HEX, wrap=False)
        cur_y += Mm(6) + Mm(2)

    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    if n_cols == 0:
        return

    # Compute column widths
    if col_ratios and len(col_ratios) >= n_cols:
        total = sum(col_ratios[:n_cols])
        col_ws = [int(avail_w * r / total) for r in col_ratios[:n_cols]]
    else:
        col_ws = [avail_w // n_cols] * n_cols
    # Correct rounding
    col_ws[-1] = avail_w - sum(col_ws[:-1])

    ROW_H  = Mm(8)
    PAD_X  = Mm(2)

    # Header row
    if headers:
        cx = l + INNER_PAD
        for ci, hdr in enumerate(headers):
            _rect(slide, l=cx, t=cur_y, w=col_ws[ci], h=ROW_H,
                  fill=BT.PRIMARY_500_HEX)
            _txb(slide, T(str(hdr)),
                 l=cx + PAD_X, t=cur_y, w=col_ws[ci] - PAD_X, h=ROW_H,
                 sz=9, bold=True, color=BT.WHITE_HEX,
                 valign=MSO_ANCHOR.MIDDLE, wrap=False)
            cx += col_ws[ci]
        cur_y += ROW_H

    # Data rows
    for ri, row in enumerate(rows):
        if cur_y + ROW_H > t + h - INNER_PAD:
            break
        row_bg = BT.NEUTRAL_100_HEX if ri % 2 == 0 else BT.WHITE_HEX
        cx = l + INNER_PAD
        for ci in range(n_cols):
            _rect(slide, l=cx, t=cur_y, w=col_ws[ci], h=ROW_H,
                  fill=row_bg, line=BT.NEUTRAL_200_HEX, lw_pt=0.5)
            cell = str(row[ci]) if ci < len(row) else ""
            _txb(slide, T(cell),
                 l=cx + PAD_X, t=cur_y, w=col_ws[ci] - PAD_X, h=ROW_H,
                 sz=9, color=BT.NEUTRAL_700_HEX,
                 valign=MSO_ANCHOR.MIDDLE, wrap=False)
            cx += col_ws[ci]
        cur_y += ROW_H


def _render_ct_quote(slide, slot: dict, l, t, w, h):
    """
    Quote / testimonial block.

    slot keys:
        quote       : str  — quote body text
        attribution : str  — source / name line
        dark        : bool — dark card (default False → primary-100 bg)
    """
    from scripts.i18n import T

    quote       = T(slot.get("quote", ""))
    attribution = T(slot.get("attribution", ""))
    is_dark     = slot.get("dark", False)

    if is_dark:
        _card(slide, l=l, t=t, w=w, h=h, bg=BT.NEUTRAL_900_HEX)
        text_color = BT.NEUTRAL_200_HEX
        mark_color = BT.PRIMARY_500_HEX
        attr_color = BT.PRIMARY_500_HEX
    else:
        _card(slide, l=l, t=t, w=w, h=h, bg=BT.PRIMARY_100_HEX)
        text_color = BT.NEUTRAL_900_HEX
        mark_color = BT.PRIMARY_500_HEX
        attr_color = BT.PRIMARY_500_HEX

    PAD = INNER_PAD + Mm(1)

    # Decorative opening quotation mark
    _txb(slide, "“",
         l=l + PAD, t=t + Mm(1), w=Mm(14), h=Mm(14),
         sz=40, bold=True, color=mark_color, wrap=False,
         en_font=BT.FONT_CN, cn_font=BT.FONT_CN)

    # Quote text
    attr_h = Mm(9) if attribution else 0
    qt     = t + Mm(10)
    qh     = h - Mm(10) - attr_h - Mm(2)
    tb = _txb(slide, quote,
              l=l + PAD, t=qt, w=w - 2 * PAD, h=qh,
              sz=13, color=text_color, ls_pt=19)
    try:
        tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    # Attribution line
    if attribution:
        _txb(slide, f"— {attribution}",
             l=l + PAD, t=t + h - attr_h - Mm(1), w=w - 2 * PAD, h=attr_h,
             sz=10, bold=True, color=attr_color, wrap=False)


def _render_ct_checklist(slide, slot: dict, l, t, w, h):
    """
    Checklist with status icons.

    slot keys:
        items  : list of str  or  {"text": ..., "status": "done"|"none"|"pending"}
        label  : str  — optional section header
    """
    from scripts.i18n import T

    items     = slot.get("items") or []
    raw_label = slot.get("label", "")

    cur_y   = t + INNER_PAD
    avail_w = w - 2 * INNER_PAD

    if raw_label:
        _txb(slide, T(raw_label),
             l=l + INNER_PAD, t=cur_y, w=avail_w, h=Mm(6),
             sz=9, bold=True, color=BT.PRIMARY_500_HEX, wrap=False)
        cur_y += Mm(6) + Mm(3)

    ROW_H   = Mm(8)
    DOT_D   = Mm(5)
    GAP     = Mm(3)
    V_GAP   = Mm(2)

    STATUS = {
        "done":    (BT.SUCCESS_HEX,        "✓"),
        "none":    (BT.DANGER_HEX,         "✗"),
        "pending": (BT.NEUTRAL_400_HEX,    "○"),
    }

    for item in items:
        if cur_y + ROW_H > t + h - INNER_PAD:
            break

        if isinstance(item, str):
            text, status = T(item), "done"
        else:
            text   = T(item.get("text", ""))
            status = item.get("status", "done")

        dot_color, icon_ch = STATUS.get(status, STATUS["done"])
        dot_t = cur_y + (ROW_H - DOT_D) // 2

        # Status dot
        _rect(slide, l=l + INNER_PAD, t=dot_t, w=DOT_D, h=DOT_D,
              fill=dot_color, radius_mm=BT.RADIUS_PILL_MM)
        _txb(slide, icon_ch,
             l=l + INNER_PAD, t=dot_t, w=DOT_D, h=DOT_D,
             sz=7, bold=True, color=BT.WHITE_HEX,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, wrap=False)

        # Item text
        _txb(slide, text,
             l=l + INNER_PAD + DOT_D + GAP, t=cur_y,
             w=avail_w - DOT_D - GAP, h=ROW_H,
             sz=11, color=BT.NEUTRAL_700_HEX,
             valign=MSO_ANCHOR.MIDDLE)

        cur_y += ROW_H + V_GAP


def _render_ct_progress(slide, slot: dict, l, t, w, h):
    """
    Horizontal progress bars.

    slot keys:
        items  : list of {"label": str, "value": 0-100, "color": hex (optional)}
                 or str (label only, value=0)
        label  : str  — optional section header
    """
    from scripts.i18n import T

    items     = slot.get("items") or []
    raw_label = slot.get("label", "")

    cur_y   = t + INNER_PAD
    avail_w = w - 2 * INNER_PAD

    if raw_label:
        _txb(slide, T(raw_label),
             l=l + INNER_PAD, t=cur_y, w=avail_w, h=Mm(6),
             sz=9, bold=True, color=BT.PRIMARY_500_HEX, wrap=False)
        cur_y += Mm(6) + Mm(3)

    n       = len(items)
    if n == 0:
        return

    avail_h  = t + h - INNER_PAD - cur_y
    item_h   = avail_h // n

    LABEL_W  = int(avail_w * 0.36)
    PCT_W    = Mm(10)
    BAR_W    = avail_w - LABEL_W - PCT_W - Mm(3)
    BAR_H    = Mm(4)
    LBL_H    = Mm(6)

    ACC = [_SLOT_ACC[i % 6] for i in range(n)]

    for i, item in enumerate(items):
        if isinstance(item, str):
            lbl, value, color = T(item), 0, ACC[i]
        else:
            lbl   = T(item.get("label", ""))
            value = max(0, min(100, int(item.get("value", 0))))
            color = item.get("color") or ACC[i]

        mid_y   = cur_y + (item_h - LBL_H) // 2
        bar_t   = cur_y + (item_h - BAR_H) // 2
        bar_l   = l + INNER_PAD + LABEL_W + Mm(2)

        # Label
        _txb(slide, lbl,
             l=l + INNER_PAD, t=mid_y, w=LABEL_W, h=LBL_H,
             sz=11, color=BT.NEUTRAL_700_HEX,
             valign=MSO_ANCHOR.MIDDLE, wrap=False)

        # Track (background)
        _rect(slide, l=bar_l, t=bar_t, w=BAR_W, h=BAR_H,
              fill=BT.NEUTRAL_200_HEX, radius_mm=2)

        # Fill
        fill_w = max(int(BAR_H * 0.8), int(BAR_W * value / 100))
        _rect(slide, l=bar_l, t=bar_t, w=fill_w, h=BAR_H,
              fill=color, radius_mm=2)

        # Percentage label
        _txb(slide, f"{value}%",
             l=bar_l + BAR_W + Mm(2), t=mid_y, w=PCT_W, h=LBL_H,
             sz=10, bold=True, color=BT.NEUTRAL_900_HEX,
             align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, wrap=False)

        cur_y += item_h


# ── Container dispatcher ─────────────────────────────────────────────────────

_CONTAINER_RENDERERS: dict = {
    "text":        _render_ct_text,
    "cards":       _render_ct_cards,
    "strips":      _render_ct_strips,
    "stats":       _render_ct_stats,
    "flow":        _render_ct_flow,
    "image":       _render_ct_image,
    "tags":        _render_ct_tags,
    "table":       _render_ct_table,
    "quote":       _render_ct_quote,
    "checklist":   _render_ct_checklist,
    "progress":    _render_ct_progress,
    "placeholder": _render_ct_placeholder,
    "blank":       _render_ct_blank,
}


def _render_container(slide, slot: dict, l: int, t: int, w: int, h: int):
    """Route a slot dict to the correct container renderer."""
    ctype    = slot.get("type", "placeholder")
    renderer = _CONTAINER_RENDERERS.get(ctype)
    if renderer is None:
        raise ValueError(
            f"Unknown container type '{ctype}'. "
            f"Available: {sorted(_CONTAINER_RENDERERS)}."
        )
    renderer(slide, slot, int(l), int(t), int(w), int(h))


# ── BrandPptx ─────────────────────────────────────────────────────────────────

class BrandPptx:
    """
    Brand-consistent PPTX builder v2.0.
    Automatically uses templates/ppt/brand_master.pptx as the base template
    if it exists, so that theme colors are properly embedded.
    """

    _TEMPLATE_PATH = os.path.abspath(
        os.path.join(os.path.dirname(__file__), "..", "templates", "ppt", "brand_master.pptx")
    )

    def __init__(self, use_template: bool = True):
        if use_template and os.path.exists(self._TEMPLATE_PATH):
            self._prs = Presentation(self._TEMPLATE_PATH)
            self._clear_slides()
        else:
            self._prs = Presentation()
        self._prs.slide_width  = SLIDE_W
        self._prs.slide_height = SLIDE_H
        self._blank = self._prs.slide_layouts[6]  # blank layout
        # Theme mode: read once at construction from brand_config.json via BT.
        # "light" → dividers use light gradient; "dark" → solid green bg.
        self._theme_mode: str = BT.THEME_MODE

    def _clear_slides(self):
        """Remove all content slides, keeping master and layouts."""
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        sldIdLst = self._prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(f"{{{r_ns}}}id")
            if rId:
                try:
                    self._prs.part.drop_rel(rId)
                except Exception:
                    pass
            sldIdLst.remove(sldId)

    def _new_slide(self):
        return self._prs.slides.add_slide(self._blank)

    # ── Cover Slide — Dark ────────────────────────────────────────────────────

    def add_cover(self, title: str, subtitle: str = "", tagline: str = ""):
        """
        Dark cover (#0E1216 bg) with gradient title text.
        Stacked reverse logo top-left, green accent bar left, bottom bar.
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.NEUTRAL_900_HEX)

        # Logo — aligned to left margin (ML)
        _add_logo_stacked(slide, reverse=True, l_mm=21, t_mm=11, h_mm=26)

        # Gradient title
        _txb_gradient(slide, title,
                      l=ML, t=Mm(65), w=Mm(205), h=Mm(38),
                      sz=42, bold=True, align=PP_ALIGN.LEFT,
                      stops=TITLE_GRADIENT)

        # Subtitle
        if subtitle:
            _txb(slide, subtitle,
                 l=ML, t=Mm(107), w=Mm(205), h=Mm(16),
                 sz=16, color=BT.NEUTRAL_400_HEX)

        # Tagline
        if tagline:
            _txb(slide, tagline,
                 l=ML, t=Mm(126), w=Mm(205), h=Mm(10),
                 sz=11, color=BT.PRIMARY_100_HEX)

        # Bottom bar
        _rect(slide, l=0, t=SLIDE_H - Mm(3.5), w=SLIDE_W, h=Mm(3.5),
              fill=BT.PRIMARY_500_HEX)
        return slide

    # ── Cover Slide — Light (VS Template Style) ───────────────────────────────

    def add_cover_light(self, title: str, subtitle: str = "",
                        date_or_meta: str = "", tagline: str = ""):
        """
        Light cover (#FAFAFA bg) with large gradient title — mirrors the VS template.
        Primary logo top-right, thin separator line, gradient 96pt title.
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.BG_PAGE_HEX)   # #F8FAFC — page-level bg

        # Thin left accent bar
        _rect(slide, l=0, t=0, w=Mm(3), h=SLIDE_H, fill=BT.PRIMARY_500_HEX)

        # Horizontal separator line
        _rect(slide, l=Mm(16), t=Mm(90), w=SLIDE_W - Mm(20), h=Mm(0.5),
              fill=BT.NEUTRAL_200_HEX)

        # Gradient title (large, 96pt equivalent ≈ 44pt at 1333px canvas)
        _txb_gradient(slide, title,
                      l=Mm(22), t=Mm(94), w=Mm(200), h=Mm(42),
                      sz=44, bold=True, align=PP_ALIGN.LEFT,
                      stops=TITLE_GRADIENT, ls_pt=52)

        # Subtitle
        if subtitle:
            _txb(slide, subtitle,
                 l=Mm(22), t=Mm(139), w=Mm(200), h=Mm(14),
                 sz=15, bold=False, color=BT.NEUTRAL_900_HEX)

        # Meta / date
        if date_or_meta:
            _txb(slide, date_or_meta,
                 l=Mm(22), t=Mm(155), w=Mm(200), h=Mm(10),
                 sz=12, color=BT.NEUTRAL_400_HEX)

        # Tagline
        if tagline:
            _txb(slide, tagline,
                 l=Mm(22), t=Mm(168), w=Mm(200), h=Mm(10),
                 sz=11, color=BT.NEUTRAL_400_HEX)

        # Primary stacked logo — top left
        _add_logo_stacked(slide, reverse=False, l_mm=13, t_mm=14, h_mm=24)

        return slide

    # ── Section Divider ───────────────────────────────────────────────────────

    def add_divider(self, chapter_title: str, chapter_num: str = ""):
        """
        Theme-aware section divider.
        • light mode (self._theme_mode == "light") → add_divider_light: horizontal gradient bg,
          PRIMARY_500 text, stacked primary logo.
        • dark mode  (self._theme_mode == "dark")  → _add_divider_dark: solid PRIMARY_500 bg,
          white text, reversed logo.
        Gen scripts always call add_divider(); brand_config.json controls which renders.
        """
        if self._theme_mode == "light":
            return self.add_divider_light(chapter_title, chapter_num)
        return self._add_divider_dark(chapter_title, chapter_num)

    def _add_divider_dark(self, chapter_title: str, chapter_num: str = ""):
        """Dark variant: solid PRIMARY_500 green background, white text, reversed logo."""
        slide = self._new_slide()
        _set_slide_bg(slide, BT.PRIMARY_500_HEX)

        _rect(slide, l=0, t=0, w=SLIDE_W, h=SLIDE_H, fill=BT.PRIMARY_500_HEX)

        if chapter_num:
            _txb(slide, chapter_num,
                 l=ML, t=Mm(60), w=SLIDE_W - 2*ML, h=Mm(14),
                 sz=13, color=BT.PRIMARY_100_HEX,
                 align=PP_ALIGN.CENTER)

        _txb(slide, chapter_title,
             l=ML, t=Mm(78), w=SLIDE_W - 2*ML, h=Mm(44),
             sz=38, bold=True, color=BT.WHITE_HEX, align=PP_ALIGN.CENTER,
             ls_pt=46)

        _add_logo_stacked(slide, reverse=True,
                          l_mm=int((SLIDE_W / Mm(1) - 30)),
                          t_mm=int((SLIDE_H / Mm(1) - 28)),
                          h_mm=18)
        return slide

    def add_divider_light(self, chapter_title: str, chapter_num: str = ""):
        """
        Light section divider: horizontal gradient bg (PRIMARY_100 → SUCCESS_100 → SECONDARY_100),
        PRIMARY_600 dark-green text, stacked primary color logo bottom-right.
        """
        slide = self._new_slide()
        # Slide background set to PRIMARY_500 for compliance check; gradient shape renders on top.
        _set_slide_bg(slide, BT.PRIMARY_500_HEX)

        # Gradient background shape spanning full slide
        grad = _rect_gradient_h(slide, 0, 0, SLIDE_W, SLIDE_H, [
            (0,   BT.PRIMARY_100_HEX),    # #EAFAF5 — light mint
            (50,  "#F4FCF0"),             # success-100 equiv — light grass-green
            (100, BT.SECONDARY_100_HEX),  # #F8FBE7 — light lime
        ])
        _send_to_back(slide, grad)

        if chapter_num:
            _txb(slide, chapter_num,
                 l=ML, t=Mm(60), w=SLIDE_W - 2 * ML, h=Mm(14),
                 sz=13, color=BT.PRIMARY_500_HEX,
                 align=PP_ALIGN.CENTER)

        _txb(slide, chapter_title,
             l=ML, t=Mm(78), w=SLIDE_W - 2 * ML, h=Mm(44),
             sz=38, bold=True, color=BT.PRIMARY_500_HEX, align=PP_ALIGN.CENTER,
             ls_pt=46)

        # Stacked primary logo (colour, not reversed) — bottom-right
        _add_logo_stacked(slide, reverse=False,
                          l_mm=int((SLIDE_W / Mm(1) - 30)),
                          t_mm=int((SLIDE_H / Mm(1) - 28)),
                          h_mm=18)
        return slide

    # ── Standard Body Slide ───────────────────────────────────────────────────

    def add_body_slide(self,
                       title: str,
                       bullets: Optional[List] = None,
                       body_text: str = "",
                       subtitle: str = "",
                       label: str = "",
                       slide_label: str = "",
                       title_deco=None,
                       note: str = "",
                       note_style: str = "note"):
        """Standard white body slide with header, bullets or free text.

        note / note_style: optional callout block at the bottom of the slide.
        note_style: "note" | "info" | "tip" | "warning" | "danger"
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide, layout_label=slide_label)

        _note_reserve = CALLOUT_H + Mm(6) if note else 0
        content_top   = CONTENT_Y + Mm(4)
        content_h     = CONTENT_H - Mm(4) - _note_reserve

        if bullets:
            _has_pills = any(
                isinstance(b, dict) and (b.get("pill") or b.get("section"))
                for b in bullets
            )
            if not _has_pills:
                # Fast path: single textbox for uniform bullet list
                tb = slide.shapes.add_textbox(ML, content_top, CW, content_h)
                tf = tb.text_frame
                tf.word_wrap = True
                for i, bullet in enumerate(bullets):
                    text = bullet if isinstance(bullet, str) else bullet.get("text", "")
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    pPr = p._p.get_or_add_pPr()
                    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
                    etree.SubElement(lnSpc, qn("a:spcPts"), attrib={"val": "2100"})
                    spc_bef = etree.SubElement(pPr, qn("a:spcBef"))
                    etree.SubElement(spc_bef, qn("a:spcPts"), attrib={"val": "800"})
                    run = p.add_run()
                    run.text = f"• {text}"
                    run.font.size = Pt(17)
                    run.font.color.rgb = _rgb(BT.NEUTRAL_700_HEX)
                    _set_run_fonts(run)
            else:
                # Per-element rendering
                # Supported dict formats:
                #   {"section": "Label"}               ← content sub-header (NOT a pill)
                #   {"pill": "Tag", "text": "..."}     ← pill prefix + body text
                #   {"pill": "Tag", "text": "", "style": "primary-soft"}  ← pill only
                BULLET_H = Mm(9)
                ROW_GAP  = Mm(2.5)
                cur_y    = content_top
                for b in bullets:
                    # ── section sub-header ────────────────────────────────────
                    if isinstance(b, dict) and b.get("section"):
                        cur_y += Mm(3)   # extra breathing room before section
                        _rect(slide,
                              l=ML, t=cur_y + Mm(2),
                              w=Mm(3), h=Mm(5),
                              fill=BT.PRIMARY_500_HEX, radius_mm=1.0)
                        _txb(slide, b["section"],
                             l=ML + Mm(5), t=cur_y, w=CW - Mm(5), h=BULLET_H,
                             sz=12, bold=True, color=BT.PRIMARY_500_HEX)
                        cur_y += BULLET_H + Mm(3)
                        continue   # skip standard ROW_GAP
                    # ── plain string bullet ───────────────────────────────────
                    if isinstance(b, str):
                        _txb(slide, f"• {b}", l=ML, t=cur_y, w=CW, h=BULLET_H,
                             sz=15, color=BT.NEUTRAL_700_HEX)
                    # ── pill + text ───────────────────────────────────────────
                    else:
                        pill_text = b.get("pill", "")
                        body_text = b.get("text", "")
                        style     = b.get("style", "primary-soft")
                        _ps       = _PILL_STYLES.get(style, _PILL_STYLES["primary-soft"])
                        text_x    = ML
                        if pill_text:
                            pw     = _pill(slide, pill_text, l=ML, t=cur_y + Mm(1),
                                          bg=_ps[0], fg=_ps[1], font_sz=9)
                            text_x = ML + pw + Mm(3)
                        if body_text:
                            _txb(slide, body_text, l=text_x, t=cur_y,
                                 w=CW - (text_x - ML), h=BULLET_H,
                                 sz=14, color=BT.NEUTRAL_700_HEX, ls_pt=20)
                    cur_y += BULLET_H + ROW_GAP
        elif body_text:
            _txb(slide, body_text,
                 l=ML, t=content_top, w=CW, h=content_h,
                 sz=16, color=BT.NEUTRAL_700_HEX, ls_pt=26)

        if note:
            _note_t = CONTENT_Y + CONTENT_H - CALLOUT_H - Mm(4)
            _callout(slide, note, l=ML, t=_note_t, w=CW, style=note_style)

        return slide

    # ── Two-Column Slide ──────────────────────────────────────────────────────

    def add_two_col_slide(self,
                          title: str,
                          left_content: str,
                          right_content: str,
                          left_title: str = "",
                          right_title: str = "",
                          subtitle: str = "",
                          label: str = "",
                          title_deco=None,
                          note: str = "",
                          note_style: str = "note"):
        """Two-column body slide with optional column subtitles."""
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        top    = CONTENT_Y + Mm(5)
        height = CONTENT_H - Mm(5)

        for i, (col_title, content) in enumerate([
            (left_title,  left_content),
            (right_title, right_content),
        ]):
            left_x = ML + i * (C2_W + C2_GAP)
            if col_title:
                # Plain bold sub-header — no pill background
                _txb(slide, col_title, l=left_x, t=top, w=C2_W, h=Mm(9),
                     sz=13, bold=True, color=BT.NEUTRAL_900_HEX)
                body_top = top + Mm(12)
            else:
                body_top = top
            _render_content_with_arrows(
                slide, content,
                x=left_x, y=body_top, w=C2_W, h=height - Mm(12),
                sz=14, color=BT.NEUTRAL_700_HEX, ls_pt=22,
                arrow_style="primary", arrow_size_mm=5.5)

        # Vertical divider
        _rect(slide,
              l=ML + C2_W + C2_GAP // 2,
              t=top + Mm(4),
              w=Mm(0.4),
              h=height - Mm(8),
              fill=BT.NEUTRAL_200_HEX)

        if note:
            _note_t = CONTENT_Y + CONTENT_H - CALLOUT_H - Mm(4)
            _callout(slide, note, l=ML, t=_note_t, w=CW, style=note_style)

        return slide

    # ── Two-Column Pill Tags ──────────────────────────────────────────────────

    def add_two_col_pills(
        self,
        title: str,
        subtitle: str = "",
        label: str = "",
        title_deco=None,
        left_title: str = "",
        left_pills: list = None,
        left_desc: str = "",
        right_title: str = "",
        right_pills: list = None,
        right_desc: str = "",
    ):
        """
        Two-column callout layout with pill/badge tags.

        Each column is a callout-style bordered card (no left strip).
        Card identity comes from the colored border + background, matching the
        callout design in add_about_slide.

        left_pills / right_pills : list[str] — each item rendered as a pill badge
        left_desc / right_desc   : str       — optional footer text (11pt NEUTRAL_400)

        Left column uses PRIMARY palette (green pills + green border).
        Right column uses SECONDARY palette (yellow-green pills + secondary border).
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        PAD_LEFT   = Mm(5)     # inner left padding (no strip, so tighter)
        PAD_X      = Mm(5)     # right inner padding
        PAD_TOP    = Mm(5)     # top inner padding
        PAD_BOT    = Mm(4)     # bottom inner padding
        TITLE_H    = Mm(12)    # column title box height
        PILL_H     = Mm(8.5)   # pill height
        PILL_GAP_H = Mm(2.5)   # horizontal gap between pills
        PILL_GAP_V = Mm(2.5)   # vertical gap between pill rows
        PILL_PX    = 4.6       # mm per char (CJK-mix heuristic at 10pt)
        PILL_PAD_X = Mm(3.5)   # horizontal text padding inside pill

        cy     = CONTENT_Y + Mm(4)
        card_h = CONTENT_H - Mm(4)

        _COLS = [
            (left_title,  left_pills  or [], left_desc,
             BT.PRIMARY_500_HEX,   BT.PRIMARY_100_HEX,
             BT.PRIMARY_500_HEX,   BT.WHITE_HEX),
            (right_title, right_pills or [], right_desc,
             BT.SECONDARY_500_HEX, BT.SECONDARY_100_HEX,
             BT.SECONDARY_500_HEX, BT.NEUTRAL_900_HEX),
        ]

        for col_idx, (col_title, pills, col_desc,
                      accent, card_bg, pill_bg, pill_fg) in enumerate(_COLS):
            cx         = ML + col_idx * (C2_W + C2_GAP)
            content_w  = C2_W - PAD_LEFT - PAD_X   # usable pill-flow width

            # Colored card — no border (only white-bg cards get a gray border)
            _card(slide, l=cx, t=cy, w=C2_W, h=card_h,
                  bg=card_bg)

            # Column title (accent color, 13pt bold)
            _txb(slide, col_title,
                 l=cx + PAD_LEFT, t=cy + PAD_TOP,
                 w=content_w, h=TITLE_H,
                 sz=13, bold=True, color=accent)

            # Pill flow
            row_x = cx + PAD_LEFT
            row_y = cy + PAD_TOP + TITLE_H + Mm(2)

            for pill_text in pills:
                n      = len(pill_text)
                pill_w = min(Mm(max(16, n * PILL_PX + 7)), content_w)

                # Wrap when current pill would overflow the column
                if row_x + pill_w > cx + C2_W - PAD_X + Mm(0.5):
                    row_x  = cx + PAD_LEFT
                    row_y += PILL_H + PILL_GAP_V

                _rect(slide, l=row_x, t=row_y, w=pill_w, h=PILL_H,
                      fill=pill_bg, radius_mm=BT.RADIUS_PILL_MM)
                _txb(slide, pill_text,
                     l=row_x + PILL_PAD_X, t=row_y + Mm(1.3),
                     w=pill_w - PILL_PAD_X * 2, h=Mm(7),
                     sz=10, bold=False, color=pill_fg)

                row_x += pill_w + PILL_GAP_H

            # Optional description at card bottom
            if col_desc:
                desc_h = Mm(18)
                desc_y = cy + card_h - PAD_BOT - desc_h
                _rect(slide, l=cx + PAD_LEFT, t=desc_y - Mm(3),
                      w=content_w, h=Mm(0.3), fill=BT.NEUTRAL_200_HEX)
                _txb(slide, col_desc,
                     l=cx + PAD_LEFT, t=desc_y,
                     w=content_w, h=desc_h,
                     sz=11, color=BT.NEUTRAL_400_HEX, ls_pt=17)

        return slide

    # ── Three-Cards Slide ─────────────────────────────────────────────────────

    def add_three_cards(self,
                        title: str,
                        cards: List[Dict[str, str]],
                        subtitle: str = "",
                        label: str = "",
                        title_deco=None):
        """
        Three-column feature cards.
        cards: [{"title": "...", "body": "...", "tag": "(optional)"}]
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        card_top = CONTENT_Y + Mm(6)
        card_h   = CONTENT_H - Mm(6)

        card_colors   = [BT.PRIMARY_100_HEX, BT.NEUTRAL_100_HEX, BT.SECONDARY_100_HEX]
        accent_colors = [BT.PRIMARY_500_HEX, BT.SUCCESS_HEX, BT.SECONDARY_500_HEX]

        for i, c in enumerate(cards[:3]):
            x          = ML + i * (C3_W + C3_GAP)
            # Three-card is always ≤ 3 items, which is below BT.MIN_CARDS_FOR_DARK.
            # Dark accent is structurally excluded here regardless of the config toggle.
            is_dark   = False
            is_danger = c.get("danger", False)
            if is_danger:
                bg, tag_color, txt_color = BT.CARD_DANGER_BG, BT.DANGER_HEX, BT.NEUTRAL_900_HEX
                body_color = BT.NEUTRAL_700_HEX
            else:
                bg        = card_colors[i % 3]
                tag_color = accent_colors[i % 3]
                txt_color = BT.NEUTRAL_900_HEX
                body_color = BT.NEUTRAL_700_HEX
            _card(slide, l=x, t=card_top, w=C3_W, h=card_h, bg=bg)

            layout = c.get("layout", "vertical_stack")
            if layout != "vertical_stack":
                _render_card_inner(slide, layout, x, card_top, C3_W, card_h,
                                   c, tag_color, txt_color=txt_color,
                                   body_color=body_color, pad_s=Mm(5), pad_t=Mm(6))
                continue

            inner_l = x + Mm(5)
            inner_w = C3_W - Mm(10)
            y_off   = card_top + Mm(6)

            tag = c.get("tag", "")
            if tag:
                _tc = BT.DANGER_HEX if is_danger else tag_color
                _txb(slide, tag, l=inner_l, t=y_off, w=inner_w, h=Mm(6),
                     sz=9, bold=True, color=_tc)
                y_off += Mm(7)

            # Card title
            _txb(slide, c.get("title", ""),
                 l=inner_l, t=y_off, w=inner_w, h=Mm(24),
                 sz=16, bold=True, color=txt_color)
            y_off += Mm(25)

            # Accent bar — inside card, below title
            _rect(slide, l=inner_l, t=y_off, w=Mm(28), h=Mm(1),
                  fill=tag_color)
            y_off += Mm(5)

            # Card body — 8mm bottom padding consistent with horizontal padding
            body_tb = _txb(slide, c.get("body", ""),
                           l=inner_l, t=y_off, w=inner_w,
                           h=card_h - (y_off - card_top) - Mm(8),
                           sz=13, color=body_color, ls_pt=18)
            body_tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        return slide

    # ── Six-Cards Slide (2 × 3 grid) ─────────────────────────────────────────

    def add_six_cards(self,
                      title: str,
                      cards: List[Dict[str, str]],
                      subtitle: str = "",
                      label: str = "",
                      title_deco=None,
                      intro_text: str = "",
                      intro_label: str = "",
                      intro_flow: list = None):
        """
        2-row × 3-column compact cards grid.
        cards: [{"title": "...", "body": "...", "tag": "(optional)"}] — up to 6
        Column index drives accent colour so each column pair shares a theme.

        intro_text / intro_label / intro_flow: when exactly 5 cards are given, the top-left slot
        becomes a summary text block instead of a card → balanced 2+3 layout:
          Row 0: [intro_text_slot]  [card0]  [card1]
          Row 1: [card2]  [card3]  [card4]
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        ROW_GAP = Mm(4)
        PAD_TOP = Mm(4)
        card_h  = (CONTENT_H - PAD_TOP - ROW_GAP) // 2

        # ── Canonical 6-slot card palette (semantic rules apply everywhere) ──────
        # Slot 1 gray  (#F2F3F5):  safe/supplementary → value text ONLY SUCCESS_HEX
        # Slot 3 orange (#FFF1DF): high-risk/warning  → bg+accent match P10 template
        # dark card   (NEUTRAL_900 #0E1216):   ★★★ 突出/亮点/强调 → SECONDARY_500 + white text
        # danger card (CARD_DANGER_BG #FFF2F2): ★ 慎用，危险/风险  → DANGER_HEX accent（浅色底）
        card_colors   = [BT.PRIMARY_100_HEX, BT.NEUTRAL_100_HEX, BT.SECONDARY_100_HEX,
                         BT.CARD_ORANGE_BG,           BT.CARD_TEAL_BG,           BT.CARD_PURPLE_BG]
        accent_colors = [BT.PRIMARY_500_HEX, BT.SUCCESS_HEX, BT.SECONDARY_500_HEX,
                         BT.WARNING_HEX,      BT.TEAL_HEX,           BT.PURPLE_HEX]

        # Dark card budget driven by brand_config.json via BT constants:
        #   enabled=False → 0 dark cards always
        #   enabled=True  → 0 if n < MIN_CARDS_FOR_DARK, else MAX_DARK_PER_SLIDE
        # Extra dark requests are silently downgraded to their slot default color.
        _n_cards     = len(cards[:6])
        _dark_budget = (
            BT.MAX_DARK_PER_SLIDE if _n_cards >= BT.MIN_CARDS_FOR_DARK else 0
        ) if BT.DARK_ACCENT_CARDS_ENABLED else 0
        _dark_used   = 0
        _cards_adj   = []
        for _c in cards[:6]:
            if _c.get("dark", False):
                if _dark_used < _dark_budget:
                    _cards_adj.append(_c)
                    _dark_used += 1
                else:
                    _c2 = {k: v for k, v in _c.items() if k != "dark"}
                    _cards_adj.append(_c2)
            else:
                _cards_adj.append(_c)

        # ── Intro text slot (5-card balanced layout) ──────────────────────────────
        # When exactly 5 cards + intro_text are provided, slot [row=0, col=0] becomes
        # a summary text block. Cards fill slots 1-5:
        #   Row 0: [intro_slot]  [card0]  [card1]
        #   Row 1: [card2]  [card3]  [card4]    → 2+3 balanced vs awkward bare 3+2
        _use_intro   = (bool(intro_text) or bool(intro_flow)) and len(_cards_adj) == 5
        _slot_offset = 1 if _use_intro else 0

        if _use_intro:
            ts_x  = ML
            ts_y  = CONTENT_Y + PAD_TOP
            _pad  = Mm(5)
            _iw   = C3_W - 2 * _pad
            cur_y = ts_y + _pad
            if intro_label:
                _txb(slide, intro_label, l=ts_x + _pad, t=cur_y,
                     w=_iw, h=Mm(6),
                     sz=8, bold=True, color=BT.PRIMARY_500_HEX)
                cur_y += Mm(8)
            if intro_text:
                _txb(slide, intro_text,
                     l=ts_x + _pad, t=cur_y,
                     w=_iw, h=Mm(20),
                     sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=17)
                cur_y += Mm(22)
            if intro_flow:
                _flow_pills(slide, intro_flow,
                            x0=ts_x + _pad, y0=cur_y, max_w=_iw,
                            font_sz=9)

        for i, c in enumerate(_cards_adj):
            slot      = i + _slot_offset
            col       = slot % 3
            row       = slot // 3
            x         = ML + col * (C3_W + C3_GAP)
            y         = CONTENT_Y + PAD_TOP + row * (card_h + ROW_GAP)
            is_dark   = c.get("dark",   False)
            is_danger = c.get("danger", False)
            if is_dark:
                bg, acc, txt_color = BT.NEUTRAL_900_HEX, BT.SECONDARY_500_HEX, BT.WHITE_HEX
                body_color = BT.NEUTRAL_400_HEX
            elif is_danger:
                bg, acc, txt_color = BT.CARD_DANGER_BG, BT.DANGER_HEX, BT.NEUTRAL_900_HEX
                body_color = BT.NEUTRAL_700_HEX
            else:
                bg         = card_colors[slot % 6]
                acc        = accent_colors[slot % 6]
                txt_color  = BT.NEUTRAL_900_HEX
                body_color = BT.NEUTRAL_700_HEX

            _card(slide, l=x, t=y, w=C3_W, h=card_h, bg=bg)

            layout = c.get("layout", "vertical_stack")
            if layout != "vertical_stack":
                _render_card_inner(slide, layout, x, y, C3_W, card_h,
                                   c, acc, txt_color=txt_color,
                                   body_color=body_color, pad_s=Mm(4), pad_t=Mm(5))
                continue

            _seq_badge(slide, x=x + C3_W - Mm(13), y=y + Mm(4), seq=i + 1, color=acc)

            inner_l = x + Mm(4)
            inner_w = C3_W - Mm(8)
            y_off   = y + Mm(5)

            tag = c.get("tag", "")
            if tag:
                if is_dark:
                    _tc = BT.SECONDARY_500_HEX
                elif is_danger:
                    _tc = BT.DANGER_HEX
                else:
                    _tc = acc
                _txb(slide, tag, l=inner_l, t=y_off, w=inner_w, h=Mm(6),
                     sz=8, bold=True, color=_tc)
                y_off += Mm(7)

            _txb(slide, c.get("title", ""),
                 l=inner_l, t=y_off, w=inner_w, h=Mm(16),
                 sz=14, bold=True, color=txt_color)
            y_off += Mm(17)

            # Accent bar — inside card, below title
            _rect(slide, l=inner_l, t=y_off, w=Mm(22), h=Mm(0.8), fill=acc)
            y_off += Mm(4)

            body_h = card_h - (y_off - y) - Mm(8)
            if body_h > Mm(8):
                body_tb = _txb(slide, c.get("body", ""),
                               l=inner_l, t=y_off, w=inner_w, h=body_h,
                               sz=11, color=body_color, ls_pt=16)
                body_tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        return slide

    # ── Big Stats Slide ───────────────────────────────────────────────────────

    def add_big_stats(self,
                      title: str,
                      stats: List[Tuple[str, str, str]],
                      subtitle: str = "",
                      label: str = "",
                      colorful=None,
                      title_deco=None,
                      note: str = "",
                      note_style: str = "note",
                      card_layout: str = "vertical_stack"):
        """
        Stats card grid.
        colorful=None (auto): ≤4 items → 2-col colorful; 5+ items → 2-row neutral
        colorful=True:  force colorful cards (best for short/consistent content, ≤6)
        colorful=False: force neutral white+gray border (for long/variable content)
        Neutral mode uses BT.EXTENDED_PALETTE (11 colors) for value text so 8 items
        get 8 distinct colors (green→blue→red→purple→orange→teal→yel-green→suc-green).
        stats: [("98%", "准确率", "核算自动化"), ...]  → (value, label, desc)

        Card semantic system (applies to all card-based slides):
          green  #EAFAF5: standard/feature      → PRIMARY_500 text
          gray   #F2F3F5: safe/supplementary    → SUCCESS_HEX text ONLY (#5CC13C)
          ygreen #F8FBE7: innovation/opportunity → SECONDARY_500 text
          orange #FFF1DF: high-risk/caution     → WARNING_HEX text (per P10 template)
          teal   #E0F7FA: expansion/eco         → #3CC5CF text
          purple #F0E8FF: ★★ 补充/战略/特殊    → #8255E1 text（颜色不够时使用）
          dark   #0E1216: ★★★ 突出/亮点/强调  → SECONDARY_500 (#C8E13C) + white label
          danger #FFF2F2: ★  慎用/危险/风险   → DANGER_HEX (#F12D2D)（确实危险才用）
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        stats = list(stats[:8])
        n     = len(stats)

        # Unified palette — slot index maps bg ↔ value-color as a pair
        CARD_BGS = [BT.PRIMARY_100_HEX, BT.NEUTRAL_100_HEX, BT.SECONDARY_100_HEX,     BT.CARD_ORANGE_BG,
                    BT.CARD_TEAL_BG,           BT.CARD_PURPLE_BG,
                    BT.PRIMARY_100_HEX, BT.NEUTRAL_100_HEX]
        VAL_COLS = [BT.PRIMARY_500_HEX, BT.SUCCESS_HEX, BT.SECONDARY_500_HEX, BT.WARNING_HEX,
                    BT.TEAL_HEX,           BT.PURPLE_HEX,
                    BT.PRIMARY_500_HEX, BT.SUCCESS_HEX]

        # colorful=True only makes visual sense for ≤6 items
        if colorful is None:
            use_color = (n <= 4)
        else:
            use_color = bool(colorful) and (n <= 6)

        # Reserve space for the callout block so cards never overlap it.
        _note_reserve = (CALLOUT_H + Mm(6)) if note else Mm(0)

        if use_color and n <= 4:
            # ── 2-column colorful layout (≤4 items) — content-adaptive row height ──
            # Val and label heights "hug" their content (1-line tight sizing).
            # Each row's height = max needed across both cards, hard-clamped to keep
            # cards within the content zone (never overlapping footer/logo).
            # Desc is bottom-aligned so text sits flush above the card's bottom margin.
            rows   = (n + 1) // 2
            gap    = Mm(8)
            card_w = (CW - gap) // 2
            _inner_w_mm = (card_w - Mm(12)) / 36000

            # Tight top-block geometry:  top_pad + val + gap + label + gap = above_desc
            _VAL_H      = Mm(20)   # 44pt bold, 1 line: ~17.5mm + margin
            _LABEL_H    = Mm(8)    # 14pt bold, 1 line: ~5mm + margin
            _ABOVE_DESC = Mm(8) + _VAL_H + Mm(2) + _LABEL_H + Mm(3)   # = Mm(41)
            _PAD_BOT    = Mm(10)

            # Hard max: cards must not reach the callout or footer zone.
            _avail_rows = SLIDE_H - FOOTER_H - (CONTENT_Y + Mm(6)) - _note_reserve
            _MAX_ROW_H  = int((_avail_rows - (rows - 1) * gap) // rows)

            if card_layout == "horizontal_split":
                # desc sits in the LEFT zone (~62% of inner width after split + gap)
                _hs_inner  = card_w - Mm(12)
                _hs_right  = int(_hs_inner * 0.38)
                _hs_left_w = (_hs_inner - _hs_right - Mm(4)) / 36000  # in mm
                _HS_ABOVE  = Mm(8) + Mm(10) + Mm(2)   # pad_t + title(1-line) + gap = 20mm
                _MIN_CARD_H = Mm(48)
                def _required_card_h(desc_text):
                    if desc_text:
                        return _HS_ABOVE + Mm(_est_text_h_mm(desc_text, 12, _hs_left_w, ls_pt=17)) + _PAD_BOT
                    return _HS_ABOVE + Mm(8) + _PAD_BOT
            else:
                _MIN_CARD_H = Mm(54)
                def _required_card_h(desc_text):
                    if desc_text:
                        return _ABOVE_DESC + Mm(_est_text_h_mm(desc_text, 12, _inner_w_mm, ls_pt=18)) + _PAD_BOT
                    return _ABOVE_DESC + Mm(8) + _PAD_BOT

            row_h = []
            for r in range(rows):
                row_items = stats[r * 2 : r * 2 + 2]
                if not note:
                    # No callout: expand cards to fill available height
                    row_h.append(int(_MAX_ROW_H))
                else:
                    rh = max((_required_card_h(d) for _, _, d in row_items), default=_MIN_CARD_H)
                    row_h.append(int(max(_MIN_CARD_H, min(rh, _MAX_ROW_H))))

            row_y = [CONTENT_Y + Mm(6)]
            for r in range(1, rows):
                row_y.append(row_y[r - 1] + row_h[r - 1] + gap)

            for i, (val, label, desc) in enumerate(stats):
                col = i % 2
                row = i // 2
                x   = ML + col * (card_w + gap)
                y   = row_y[row]
                ch  = row_h[row]
                vc  = VAL_COLS[i % len(VAL_COLS)]

                _card(slide, l=x, t=y, w=card_w, h=ch,
                      bg=CARD_BGS[i % len(CARD_BGS)])

                if card_layout == "horizontal_split":
                    _render_card_inner(
                        slide, "horizontal_split", x, y, card_w, ch,
                        {"metric": val, "title": label, "body": desc,
                         "metric_sz": 44, "title_sz": 14, "body_sz": 12},
                        vc, body_color=BT.NEUTRAL_400_HEX,
                        pad_s=Mm(6), pad_t=Mm(8))
                else:
                    # Val — tight hug height
                    _txb(slide, val,
                         l=x + Mm(6), t=y + Mm(8),
                         w=card_w - Mm(12), h=_VAL_H,
                         sz=44, bold=True, align=PP_ALIGN.LEFT, color=vc)
                    # Label — immediately after val with 2mm gap
                    _txb(slide, label,
                         l=x + Mm(6), t=y + Mm(8) + _VAL_H + Mm(2),
                         w=card_w - Mm(12), h=_LABEL_H,
                         sz=14, bold=True, color=BT.NEUTRAL_900_HEX)
                    if desc:
                        # Desc spans remaining space; bottom-aligned so text sits
                        # flush above the card's bottom margin regardless of length.
                        _txb(slide, desc,
                             l=x + Mm(6), t=y + _ABOVE_DESC,
                             w=card_w - Mm(12), h=ch - _ABOVE_DESC - _PAD_BOT,
                             sz=12, color=BT.NEUTRAL_400_HEX, ls_pt=18,
                             valign=MSO_ANCHOR.BOTTOM)

        elif use_color:
            # ── 2-ROW colorful layout (5-6 items) — content-adaptive row height ──
            top_n   = n // 2
            bot_n   = n - top_n
            row_gap = Mm(5)
            col_gap = Mm(4)
            # Tight top-block geometry — val and label hug their single line.
            # Hard max ensures cards cannot reach footer/logo zone.
            _VAL_H  = Mm(18)   # 40pt bold 1-line (~14mm + margin)
            _LBL_H  = Mm(8)    # 13pt bold 1-line (~5mm + margin)
            _ABOVE  = Mm(8) + _VAL_H + Mm(2) + _LBL_H + Mm(3)   # = Mm(39)
            _BOT    = Mm(10)
            _MIN_H  = Mm(50)
            _avail  = SLIDE_H - FOOTER_H - (CONTENT_Y + Mm(3)) - _note_reserve
            _MAX_H  = int((_avail - row_gap) // 2)

            def _color_row_h(items, ncols):
                cw = (CW - (ncols - 1) * col_gap) // ncols
                iw = (cw - Mm(12)) / 36000
                return int(min(_MAX_H, max(_MIN_H, max(
                    (_ABOVE + Mm(_est_text_h_mm(d, 11, iw, ls_pt=16)) + _BOT
                     if d else _ABOVE + _BOT)
                    for _, _, d in items))))

            top_ch = _color_row_h(stats[:top_n], top_n)
            bot_ch = _color_row_h(stats[top_n:], bot_n)
            top_y  = CONTENT_Y + Mm(3)
            bot_y  = top_y + top_ch + row_gap

            def _render_color_row(items, start_idx, y, ncols, card_h):
                card_w = (CW - (ncols - 1) * col_gap) // ncols
                for j, (val, label, desc) in enumerate(items):
                    i  = start_idx + j
                    x  = ML + j * (card_w + col_gap)
                    vc = VAL_COLS[i % len(VAL_COLS)]
                    _card(slide, l=x, t=y, w=card_w, h=card_h,
                          bg=CARD_BGS[i % len(CARD_BGS)])
                    if not _is_numeric_val(val):
                        _seq_badge(slide, x=x + card_w - Mm(13), y=y + Mm(4),
                                   seq=i + 1, color=vc)
                    _txb(slide, val,
                         l=x + Mm(6), t=y + Mm(8),
                         w=card_w - Mm(12), h=_VAL_H,
                         sz=40, bold=True, align=PP_ALIGN.LEFT, color=vc)
                    _txb(slide, label,
                         l=x + Mm(6), t=y + Mm(8) + _VAL_H + Mm(2),
                         w=card_w - Mm(12), h=_LBL_H,
                         sz=13, bold=True, color=BT.NEUTRAL_900_HEX)
                    if desc:
                        _txb(slide, desc,
                             l=x + Mm(6), t=y + _ABOVE,
                             w=card_w - Mm(12), h=card_h - _ABOVE - _BOT,
                             sz=11, color=BT.NEUTRAL_400_HEX, ls_pt=16,
                             valign=MSO_ANCHOR.BOTTOM)

            _render_color_row(stats[:top_n], 0,     top_y, top_n, top_ch)
            _render_color_row(stats[top_n:], top_n, bot_y, bot_n, bot_ch)

        else:
            # ── 2-ROW neutral layout (white+gray border) — content-adaptive ──
            # Value text uses BT.EXTENDED_PALETTE (11 colors, max contrast).
            top_n   = n // 2
            bot_n   = n - top_n
            row_gap = Mm(5)
            col_gap = Mm(4)
            _VAL_H  = Mm(18)
            _LBL_H  = Mm(8)
            _ABOVE  = Mm(8) + _VAL_H + Mm(2) + _LBL_H + Mm(3)   # = Mm(39)
            _BOT    = Mm(10)
            _MIN_H  = Mm(50)
            _avail  = SLIDE_H - FOOTER_H - (CONTENT_Y + Mm(3)) - _note_reserve
            _MAX_H  = int((_avail - row_gap) // 2)

            def _neutral_row_h(items, ncols):
                cw = (CW - (ncols - 1) * col_gap) // ncols
                iw = (cw - Mm(12)) / 36000
                return int(min(_MAX_H, max(_MIN_H, max(
                    (_ABOVE + Mm(_est_text_h_mm(d, 11, iw, ls_pt=16)) + _BOT
                     if d else _ABOVE + _BOT)
                    for _, _, d in items))))

            top_ch = _neutral_row_h(stats[:top_n], top_n)
            bot_ch = _neutral_row_h(stats[top_n:], bot_n)
            top_y  = CONTENT_Y + Mm(3)
            bot_y  = top_y + top_ch + row_gap

            def _render_neutral_row(items, start_idx, y, ncols, card_h):
                card_w = (CW - (ncols - 1) * col_gap) // ncols
                for j, (val, label, desc) in enumerate(items):
                    i  = start_idx + j
                    x  = ML + j * (card_w + col_gap)
                    vc = BT.EXTENDED_PALETTE[i % len(BT.EXTENDED_PALETTE)]
                    _card(slide, l=x, t=y, w=card_w, h=card_h,
                          bg=BT.WHITE_HEX, border=BT.BORDER_DEFAULT_HEX)
                    if not _is_numeric_val(val):
                        _seq_badge(slide, x=x + card_w - Mm(13), y=y + Mm(4),
                                   seq=i + 1, color=vc)
                    _txb(slide, val,
                         l=x + Mm(6), t=y + Mm(8),
                         w=card_w - Mm(12), h=_VAL_H,
                         sz=40, bold=True, align=PP_ALIGN.LEFT, color=vc)
                    _txb(slide, label,
                         l=x + Mm(6), t=y + Mm(8) + _VAL_H + Mm(2),
                         w=card_w - Mm(12), h=_LBL_H,
                         sz=13, bold=True, color=BT.NEUTRAL_900_HEX)
                    if desc:
                        _txb(slide, desc,
                             l=x + Mm(6), t=y + _ABOVE,
                             w=card_w - Mm(12), h=card_h - _ABOVE - _BOT,
                             sz=11, color=BT.NEUTRAL_400_HEX, ls_pt=16,
                             valign=MSO_ANCHOR.BOTTOM)

            _render_neutral_row(stats[:top_n], 0,     top_y, top_n, top_ch)
            _render_neutral_row(stats[top_n:], top_n, bot_y, bot_n, bot_ch)

        if note:
            _note_t = CONTENT_Y + CONTENT_H - CALLOUT_H - Mm(4)
            _callout(slide, note, l=ML, t=_note_t, w=CW, style=note_style)

        return slide

    # ── Timeline Slide ────────────────────────────────────────────────────────

    def add_timeline(self,
                     title: str,
                     milestones: List[Tuple[str, str, str]],
                     subtitle: str = "",
                     label: str = "",
                     title_deco=None,
                     note: str = "",
                     note_style: str = "note"):
        """
        Adaptive timeline — no hard cap on item count.
        - 1–6  items → single horizontal row (full content height)
        - 7–12 items → two-row wrap: top row = n//2 items, bottom = remainder
          Both rows share the same visual language; split is as balanced as possible.
        milestones: [("period", "title", "description"), ...]
        """
        n = len(milestones)
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)
        if n == 0:
            return slide

        if n <= 6:
            # Single row — full content area
            self._draw_timeline_row(
                slide, milestones,
                row_top=CONTENT_Y + Mm(2),
                row_h=CONTENT_H - Mm(4),
            )
        else:
            # Two-row wrap
            # Row 1 spine extends to slide right edge;
            # Row 2 spine extends from slide left edge.
            # This creates a visual "snake" showing timeline continuity.
            n1      = n // 2          # top row (floor → fewer or equal)
            ROW_GAP = Mm(7)
            PAD_V   = Mm(2)
            row_h   = (CONTENT_H - ROW_GAP - PAD_V * 2) // 2

            row1_top = CONTENT_Y + PAD_V
            row2_top = row1_top + row_h + ROW_GAP

            self._draw_timeline_row(slide, milestones[:n1],
                                    row1_top, row_h, extend_right=True)
            self._draw_timeline_row(slide, milestones[n1:],
                                    row2_top, row_h, extend_left=True)

        if note:
            _note_t = CONTENT_Y + CONTENT_H - CALLOUT_H - Mm(4)
            _callout(slide, note, l=ML, t=_note_t, w=CW, style=note_style)

        return slide

    def _draw_timeline_row(self,
                           slide,
                           items: List[Tuple[str, str, str]],
                           row_top: int,
                           row_h: int,
                           extend_left: bool = False,
                           extend_right: bool = False):
        """
        Draw one horizontal timeline row inside the given vertical bounds.
        Font sizes and dot radius scale automatically with item density.

        items:        [("period", "title", "description"), ...]
        row_top:      top-y of the row area (EMU)
        row_h:        total height of the row area (EMU)
        extend_left:  stretch spine from slide left edge (x=0) instead of
                      stopping at the first dot — used for the bottom wrap row
        extend_right: stretch spine to slide right edge (x=SLIDE_W) instead of
                      stopping at the last dot — used for the top wrap row
        """
        n = len(items)
        if n == 0:
            return

        shrink   = max(0, n - 3)
        sz_per   = max(9,  12 - shrink)
        sz_title = max(10, 14 - shrink)
        sz_desc  = max(9,  12 - shrink)
        ls_desc  = max(13, 17 - shrink)
        dot_r    = Mm(3.5) if n <= 2 else (Mm(3) if n <= 4 else Mm(2.5))

        item_w  = (CW - Mm(4)) // n
        # Spine sits at ~30 % from the top of the row area so period labels
        # have room above and descriptions have the majority of space below.
        spine_y = row_top + int(row_h * 0.30)
        dot_y   = spine_y - dot_r

        # Dot centre positions
        cx_first = ML + item_w // 2
        cx_last  = ML + item_w // 2 + (n - 1) * item_w

        # Spine: extend to slide edges when wrapping, otherwise dot-to-dot
        spine_l = 0        if extend_left  else cx_first
        spine_r = SLIDE_W  if extend_right else cx_last

        _rect(slide,
              l=spine_l,
              t=spine_y - Mm(0.7),
              w=spine_r - spine_l,
              h=Mm(1.4),
              fill=BT.PRIMARY_500_HEX)

        ACCENT = [BT.PRIMARY_500_HEX, BT.SUCCESS_HEX,
                  BT.SECONDARY_500_HEX, BT.SUCCESS_HEX,
                  BT.PRIMARY_500_HEX, BT.SUCCESS_HEX]

        for i, (period, m_title, desc) in enumerate(items):
            cx = ML + item_w // 2 + i * item_w
            ac = ACCENT[i % len(ACCENT)]

            # White halo so dot pops over spine
            _rect(slide,
                  l=cx - dot_r - Mm(1), t=dot_y - Mm(1),
                  w=(dot_r + Mm(1)) * 2, h=(dot_r + Mm(1)) * 2,
                  fill=BT.WHITE_HEX, radius_mm=BT.RADIUS_PILL_MM)
            _rect(slide,
                  l=cx - dot_r, t=dot_y,
                  w=dot_r * 2, h=dot_r * 2,
                  fill=ac, radius_mm=BT.RADIUS_PILL_MM)

            # Period label — pill badge, width fitted to content, centered over dot.
            # Badge sits just above the dot with a small gap.
            _BADGE_H = Mm(6.5)
            # Content-aware badge width: estimate from character count + padding,
            # cap at (item_w - 6mm) so it never overflows the column.
            # Hardcoded 28mm was too small for long mixed Latin/CJK period strings.
            _CH_W    = Mm(sz_per * 0.36)                          # mixed-script per-char estimate
            _BADGE_W = min(
                item_w - Mm(6),                                    # never wider than the column
                max(Mm(20), int(len(period) * _CH_W + Mm(7)))     # text width + pill padding
            )
            _BADGE_L = max(int(ML), cx - _BADGE_W // 2)           # centered; never left of margin
            _BADGE_T = dot_y - Mm(2) - _BADGE_H
            _rect(slide, l=_BADGE_L, t=_BADGE_T, w=_BADGE_W, h=_BADGE_H,
                  fill=ac, radius_mm=BT.RADIUS_PILL_MM)
            _txb(slide, period,
                 l=_BADGE_L, t=_BADGE_T,
                 w=_BADGE_W, h=_BADGE_H,
                 sz=sz_per, bold=True, color=BT.WHITE_HEX,
                 align=PP_ALIGN.CENTER, wrap=False,
                 valign=MSO_ANCHOR.MIDDLE)

            # Title — below dot
            title_top = spine_y + dot_r + Mm(3)
            title_h   = Mm(14)
            _txb(slide, m_title,
                 l=cx - item_w // 2 + Mm(2), t=title_top,
                 w=item_w - Mm(4), h=title_h,
                 sz=sz_title, bold=True, color=BT.NEUTRAL_900_HEX,
                 align=PP_ALIGN.CENTER)

            # Description — fills remaining row height
            if desc:
                desc_top = title_top + title_h + Mm(2)
                desc_h   = row_top + row_h - desc_top - Mm(2)
                if desc_h > Mm(6):
                    _txb(slide, desc,
                         l=cx - item_w // 2 + Mm(2), t=desc_top,
                         w=item_w - Mm(4), h=desc_h,
                         sz=sz_desc, color=BT.NEUTRAL_400_HEX,
                         align=PP_ALIGN.CENTER, ls_pt=ls_desc)

    # ── Quote Slide ───────────────────────────────────────────────────────────

    def add_quote(self,
                  quote_text: str,
                  author: str = "",
                  role: str = "",
                  subtitle: str = ""):
        """
        Dark quote / testimonial slide.
        Large quote mark, body text white, author in gradient.
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.NEUTRAL_900_HEX)

        # Big quotation mark
        _txb(slide, "“",
             l=ML, t=Mm(28), w=Mm(30), h=Mm(26),
             sz=56, bold=True, color=BT.PRIMARY_500_HEX)

        # Quote body
        _txb(slide, quote_text,
             l=ML + Mm(2), t=Mm(54),
             w=CW - Mm(4), h=Mm(52),
             sz=22, bold=False, color=BT.WHITE_HEX,
             align=PP_ALIGN.LEFT, ls_pt=33)

        # Thin accent line
        _rect(slide, l=ML, t=Mm(112), w=Mm(32), h=Mm(1.5),
              fill=BT.SECONDARY_500_HEX)

        # Author — solid primary color (gradient reserved for covers/section/closing)
        if author:
            _txb(slide, author,
                 l=ML, t=Mm(117), w=Mm(180), h=Mm(12),
                 sz=16, bold=True, color=BT.PRIMARY_500_HEX)

        # Role
        if role:
            _txb(slide, role,
                 l=ML, t=Mm(131), w=Mm(180), h=Mm(10),
                 sz=12, color=BT.NEUTRAL_400_HEX)

        # Optional subtitle label (top-right)
        if subtitle:
            _txb(slide, subtitle,
                 l=SLIDE_W - ML - Mm(80), t=Mm(12),
                 w=Mm(80), h=Mm(10),
                 sz=10, color=BT.NEUTRAL_400_HEX, align=PP_ALIGN.RIGHT)

        # Logo reversed
        _add_logo_h(slide, right_mm=10, bottom_mm=4, h_mm=7)

        # Bottom bar
        _rect(slide, l=0, t=SLIDE_H - Mm(3.5), w=SLIDE_W, h=Mm(3.5),
              fill=BT.PRIMARY_500_HEX)
        return slide

    # ── Table Slide ───────────────────────────────────────────────────────────

    def add_table_slide(self,
                        title: str,
                        headers: List[str],
                        rows: List[List[str]],
                        subtitle: str = "",
                        note: str = "",
                        label: str = "",
                        title_deco=None):
        """
        Data table slide.
        Header row: #3EC99E bg + white text.
        Zebra rows: alternate #F2F3F5 / white.
        """
        from pptx.util import Pt
        from pptx.oxml.ns import qn

        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        n_cols = len(headers)
        n_rows = len(rows)
        if n_cols == 0:
            return slide

        t_top  = int(CONTENT_Y + Mm(5))
        t_h    = int(CONTENT_H - Mm(8) - (CALLOUT_H + Mm(5) if note else 0))
        row_h  = min(int(t_h // (n_rows + 1)), int(Mm(12)))

        table  = slide.shapes.add_table(
            n_rows + 1, n_cols,
            int(ML), t_top,
            int(CW), row_h * (n_rows + 1)
        ).table

        # Header row
        for j, hdr in enumerate(headers):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(BT.PRIMARY_500_HEX)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = hdr
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = _rgb(BT.WHITE_HEX)
            _set_run_fonts(run)
            p.alignment = PP_ALIGN.LEFT

        # Data rows
        for i, row in enumerate(rows):
            bg = BT.NEUTRAL_100_HEX if i % 2 == 0 else BT.WHITE_HEX
            for j, val in enumerate(row[:n_cols]):
                cell = table.cell(i + 1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(bg)
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(val)
                run.font.size = Pt(11)
                run.font.color.rgb = _rgb(
                    BT.NEUTRAL_900_HEX if j == 0 else BT.NEUTRAL_700_HEX
                )
                _set_run_fonts(run)
                p.alignment = PP_ALIGN.LEFT

        # Note — callout block below table
        if note:
            _callout(slide, note,
                     l=ML, t=t_top + row_h * (n_rows + 1) + Mm(4),
                     w=CW, style="note")

        return slide

    # ── Text + Image Layouts ──────────────────────────────────────────────────

    def add_text_image_right(self,
                             title: str,
                             body_text: str,
                             image_path: Optional[str] = None,
                             subtitle: str = "",
                             label: str = "",
                             title_deco=None):
        """Left 60% text, right 40% image placeholder."""
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        text_w  = int(CW * 0.58)
        img_w   = int(CW - text_w - Mm(6))
        img_l   = ML + text_w + Mm(6)

        _txb(slide, body_text,
             l=ML, t=CONTENT_Y + Mm(6), w=text_w,
             h=CONTENT_H - Mm(6), sz=14, ls_pt=22)

        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, img_l,
                int(CONTENT_Y + Mm(6)), img_w,
                int(CONTENT_H - Mm(8)))
        else:
            _card(slide, img_l, int(CONTENT_Y + Mm(6)),
                  img_w, int(CONTENT_H - Mm(8)),
                  bg=BT.PRIMARY_100_HEX, radius_mm=BT.RADIUS_LG_MM)
        return slide

    def add_image_left_text(self,
                            title: str,
                            body_text: str,
                            image_path: Optional[str] = None,
                            subtitle: str = "",
                            label: str = "",
                            title_deco=None):
        """Left 40% image, right 60% text."""
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        img_w  = int(CW * 0.40)
        text_w = int(CW - img_w - Mm(6))
        text_l = ML + img_w + Mm(6)

        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, int(ML),
                int(CONTENT_Y + Mm(6)), img_w,
                int(CONTENT_H - Mm(8)))
        else:
            _card(slide, int(ML), int(CONTENT_Y + Mm(6)),
                  img_w, int(CONTENT_H - Mm(8)),
                  radius_mm=BT.RADIUS_LG_MM,
                  bg=BT.PRIMARY_100_HEX)

        _txb(slide, body_text,
             l=text_l, t=CONTENT_Y + Mm(6), w=text_w,
             h=CONTENT_H - Mm(6), sz=14, ls_pt=22)
        return slide

    # ── Table of Contents ─────────────────────────────────────────────────────

    def add_toc(self,
                title: str,
                chapters: List[Dict[str, str]],
                label: str = "TABLE OF CONTENTS",
                description: str = "",
                title_deco=None):
        """
        2×3 chapter card grid TOC.
        chapters: [{"num":"01","title":"关于我们","subtitle":"公司简介·愿景","state":"done|current|upcoming"}]

        state values:
          "done"     — green num, light card bg
          "current"  — dark card (#0E1216), green num, white text
          "upcoming" — dark num, light card, neutral arrow
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, label=label, title_deco=title_deco)
        _footer(slide)

        if description:
            _txb(slide, description, l=Mm(208), t=Mm(11), w=Mm(109), h=Mm(18),
                 sz=9, color=BT.NEUTRAL_700_HEX, wrap=True)

        CARD_W  = Mm(93.5)
        CARD_H  = Mm(56.4)
        COL_GAP = Mm(7.9)
        ROW_GAP = Mm(5.3)
        GRID_Y  = CONTENT_Y + Mm(4)

        # Arrow button: 6.3×6.3mm, right-aligned inside card (master P2 measurement)
        ARR_SZ  = Mm(6.3)
        ARR_L   = CARD_W - ARR_SZ - Mm(5.6)   # 81.4mm from card left
        ARR_T   = Mm(14.6)                     # from card top

        for i, ch in enumerate(chapters[:6]):
            col   = i % 3
            row   = i // 3
            cx    = ML + col * (CARD_W + COL_GAP)
            cy    = GRID_Y + row * (CARD_H + ROW_GAP)
            state = ch.get("state", "done")

            if state == "current":
                card_bg = BT.NEUTRAL_900_HEX
                num_c   = BT.PRIMARY_500_HEX
                ttl_c   = BT.WHITE_HEX
                sub_c   = BT.NEUTRAL_400_HEX
                arr_bg  = BT.PRIMARY_500_HEX
                arr_c   = BT.WHITE_HEX
                line_c  = BT.PRIMARY_500_HEX
            elif state == "upcoming":
                # Static PPT/PDF has no interactive state — use same green as "done"
                # so all non-active chapters read with full visual weight.
                card_bg = BT.BG_PAGE_HEX
                num_c   = BT.PRIMARY_500_HEX
                ttl_c   = BT.NEUTRAL_900_HEX
                sub_c   = BT.NEUTRAL_700_HEX
                arr_bg  = BT.PRIMARY_100_HEX
                arr_c   = BT.PRIMARY_500_HEX
                line_c  = BT.PRIMARY_500_HEX
            else:  # done
                card_bg = BT.BG_PAGE_HEX
                num_c   = BT.PRIMARY_500_HEX
                ttl_c   = BT.NEUTRAL_900_HEX
                sub_c   = BT.NEUTRAL_700_HEX
                arr_bg  = BT.PRIMARY_100_HEX
                arr_c   = BT.PRIMARY_500_HEX
                line_c  = BT.PRIMARY_500_HEX

            # Card background (no border — visual hierarchy via number + arrow)
            _rect(slide, l=cx, t=cy, w=CARD_W, h=CARD_H, fill=card_bg,
                  radius_mm=BT.RADIUS_SM_MM)

            # Large chapter number — top-left
            _txb(slide, ch.get("num", ""),
                 l=cx + Mm(5.6), t=cy + Mm(5.6), w=Mm(20), h=Mm(15.5),
                 sz=44, bold=True, color=num_c, wrap=False)

            # Arrow button — top-right (pill circle + "→" text)
            _rect(slide, l=cx + ARR_L, t=cy + ARR_T, w=ARR_SZ, h=ARR_SZ,
                  fill=arr_bg, radius_mm=BT.RADIUS_PILL_MM)
            _txb(slide, "→",
                 l=cx + ARR_L, t=cy + ARR_T + Mm(0.6), w=ARR_SZ, h=ARR_SZ,
                 sz=8, bold=True, color=arr_c, align=PP_ALIGN.CENTER)

            # Title
            _txb(slide, ch.get("title", ""),
                 l=cx + Mm(5.6), t=cy + Mm(27.5), w=CARD_W - Mm(10), h=Mm(7.8),
                 sz=15, bold=True, color=ttl_c)

            # Thin accent line
            _rect(slide, l=cx + Mm(5.6), t=cy + Mm(37), w=Mm(5.6), h=Mm(0.5),
                  fill=line_c)

            # Subtitle
            sub = ch.get("subtitle", "")
            if sub:
                _txb(slide, sub,
                     l=cx + Mm(5.6), t=cy + Mm(40.7), w=CARD_W - Mm(10), h=Mm(10),
                     sz=8, color=sub_c, wrap=True)

        return slide

    # ── Dark Cover with Background Image (Slide 4 style) ──────────────────────

    def add_cover_image(self,
                        title: str,
                        subtitle: str = "",
                        tagline: str = "",
                        date_label: str = "",
                        meta_pairs: Optional[List[Tuple[str, str]]] = None,
                        bg_image_path: Optional[str] = None):
        """
        Full-bleed dark cover with optional background image overlay.
        Mirrors the photo-cover variant in the master template.

        meta_pairs: up to 3 (value, label) pairs shown as bottom info strip.
                    e.g. [("2026·03","PRESENTATION VERSION"),
                          ("www.example.com","OFFICIAL WEBSITE"),
                          ("Enterprise AI","CORE FOCUS")]
        bg_image_path: optional path to a full-slide background photo.
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.NEUTRAL_900_HEX)

        if bg_image_path and os.path.exists(bg_image_path):
            slide.shapes.add_picture(
                bg_image_path, 0, 0, int(SLIDE_W), int(SLIDE_H))

        # Stacked reverse logo — top left
        _add_logo_stacked(slide, reverse=True, l_mm=25, t_mm=13, h_mm=26)

        if tagline:
            _txb(slide, tagline,
                 l=SLIDE_W - Mm(100), t=Mm(13), w=Mm(96), h=Mm(5),
                 sz=7, color=BT.NEUTRAL_400_HEX, align=PP_ALIGN.RIGHT)

        # Section tag
        if subtitle:
            _txb(slide, subtitle,
                 l=Mm(28), t=Mm(56), w=Mm(200), h=Mm(6),
                 sz=8, color=BT.PRIMARY_500_HEX)

        # Main title (white, 48pt, multiline)
        _txb(slide, title,
             l=Mm(28), t=Mm(65), w=Mm(256), h=Mm(50),
             sz=48, bold=True, color=BT.WHITE_HEX, ls_pt=56)

        if date_label:
            _txb(slide, date_label,
                 l=Mm(28), t=Mm(124), w=Mm(200), h=Mm(8),
                 sz=14, color=BT.NEUTRAL_200_HEX)

        # Bottom meta strip
        if meta_pairs:
            strip_y = Mm(164)
            col_w   = (SLIDE_W - Mm(56)) // len(meta_pairs)
            for j, (val, lbl) in enumerate(meta_pairs[:3]):
                bx = Mm(28) + j * col_w
                _txb(slide, val,
                     l=bx, t=strip_y, w=col_w - Mm(4), h=Mm(8),
                     sz=11, color=BT.WHITE_HEX)
                _txb(slide, lbl,
                     l=bx, t=strip_y + Mm(8), w=col_w - Mm(4), h=Mm(6),
                     sz=9, color=BT.NEUTRAL_400_HEX)

        # Bottom bar
        _rect(slide, l=0, t=SLIDE_H - Mm(3.5), w=SLIDE_W, h=Mm(3.5),
              fill=BT.PRIMARY_500_HEX)
        return slide

    # ── Rich Chapter Divider (Slide 6 style) ──────────────────────────────────

    def add_divider_rich(self,
                         chapter_num: str,
                         chapter_title: str,
                         subtitle: str = "",
                         chapter_items: Optional[List[Tuple[str, str]]] = None,
                         current_item: int = 0,
                         bg_image_path: Optional[str] = None):
        """
        Dark divider with oversized decorative chapter number and right sidebar.

        chapter_num:   "01"  (shown at 260pt as teal background accent)
        chapter_title: main title (60pt, white)
        subtitle:      optional subtitle below the title
        chapter_items: [("01","关于我们"), ("02","公司简介"), ...]
                       shown in the right sidebar; current_item is highlighted white
        current_item:  0-based index of the active entry in chapter_items
        bg_image_path: optional full-slide background image overlay
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.NEUTRAL_900_HEX)

        if bg_image_path and os.path.exists(bg_image_path):
            slide.shapes.add_picture(
                bg_image_path, 0, 0, int(SLIDE_W), int(SLIDE_H))

        # Oversized decorative number — serves as background texture.
        # alpha=8 keeps it barely visible so it doesn't compete with foreground text.
        _txb(slide, chapter_num,
             l=Mm(28), t=Mm(38), w=Mm(160), h=Mm(83),
             sz=260, bold=True, color=BT.PRIMARY_500_HEX, alpha=8)

        # "CHAPTER · XX" small label
        _txb(slide, f"CHAPTER · {chapter_num}",
             l=Mm(28), t=Mm(84.7), w=Mm(106), h=Mm(6),
             sz=9, bold=True, color=BT.PRIMARY_500_HEX)

        # Chapter title
        _txb(slide, chapter_title,
             l=Mm(28), t=Mm(93), w=Mm(212), h=Mm(32),
             sz=60, bold=True, color=BT.WHITE_HEX, ls_pt=68)

        if subtitle:
            _txb(slide, subtitle,
                 l=Mm(28.9), t=Mm(124), w=Mm(212), h=Mm(10),
                 sz=14, color=BT.NEUTRAL_200_HEX)

        # Thin accent line
        _rect(slide, l=Mm(28), t=Mm(140), w=Mm(21), h=Mm(0.7),
              fill=BT.WHITE_HEX)

        # Right sidebar — chapter list
        SIDEBAR_X = Mm(232.8)
        sidebar_w = SLIDE_W - SIDEBAR_X - MR

        if chapter_items:
            _txb(slide, "IN THIS CHAPTER",
                 l=SIDEBAR_X, t=Mm(93), w=sidebar_w, h=Mm(7),
                 sz=7, color=BT.NEUTRAL_400_HEX)
            _rect(slide, l=SIDEBAR_X, t=Mm(101), w=sidebar_w, h=Mm(0.2),
                  fill=BT.WHITE_HEX)
            for j, (num, name) in enumerate(chapter_items):
                item_y  = Mm(103.5) + j * Mm(6.4)
                is_curr = (j == current_item)
                # Active: secondary accent (#C8E13C) + bold for clear highlight.
                # Inactive: neutral-400 (medium grey) — clearly receded.
                _txb(slide, f"{num}  {name}",
                     l=SIDEBAR_X, t=item_y, w=sidebar_w, h=Mm(6),
                     sz=9, bold=is_curr,
                     color=BT.SECONDARY_500_HEX if is_curr else BT.NEUTRAL_400_HEX)

        _footer(slide)
        return slide

    # ── About / Intro Slide (Slide 13/14 style) ───────────────────────────────

    def add_about_slide(self,
                        title: str,
                        body_text: str,
                        label: str = "",
                        callout_items: Optional[List[Dict[str, str]]] = None,
                        right_panel: Optional[Dict[str, Any]] = None,
                        title_deco=None):
        """
        Left text + right dark panel layout (company intro / team page).

        body_text:     main paragraph on the left
        label:         small section label above title (e.g. "01 · ABOUT US")
        callout_items: compact info boxes below body_text.
                       [{"label":"OFFICIAL WEBSITE","value":"https://..."}]
        right_panel:   dark panel on the right with accent items.
                       {"title_label":"OUR VISION","items":[{"label":"...","text":"...","accent":"#3EC99E"}]}
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, label=label, title_deco=title_deco)
        _footer(slide)

        LEFT_W = Mm(138)

        # Body text
        BODY_T = Mm(50)
        BODY_H = Mm(30)
        _txb(slide, body_text,
             l=ML, t=BODY_T, w=LEFT_W, h=BODY_H,
             sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=17, wrap=True)

        # Callout cards — gradient fill + primary border + label pill
        if callout_items:
            ITEM_W = LEFT_W // 2 - Mm(3)
            ITEM_H = Mm(24)
            GAP_X  = Mm(6)
            _GRAD  = [
                (0,   BT.PRIMARY_100_HEX),   # #EAFAF5  light teal
                (50,  "#EDFBEE"),              # light fresh green (success tint)
                (100, BT.SECONDARY_100_HEX),  # #F8FBE7  light yellow-green
            ]
            for k, item in enumerate(callout_items[:2]):
                bx = ML + k * (ITEM_W + GAP_X)
                by = BODY_T + BODY_H + Mm(8)  # 8mm gap below body text
                # Gradient rounded card with primary border
                card = _rect(slide, l=bx, t=by, w=ITEM_W, h=ITEM_H,
                             fill=BT.PRIMARY_100_HEX,
                             line=BT.PRIMARY_500_HEX, lw_pt=0.4,
                             radius_mm=BT.RADIUS_MD_MM)
                _apply_shape_gradient(card, _GRAD, angle_deg=135)
                # Label pill
                _rect(slide, l=bx + Mm(5), t=by + Mm(5), w=Mm(24), h=Mm(6.5),
                      fill=BT.PRIMARY_500_HEX, radius_mm=BT.RADIUS_PILL_MM)
                _txb(slide, item.get("label", ""),
                     l=bx + Mm(5), t=by + Mm(5), w=Mm(24), h=Mm(6.5),
                     sz=7, bold=True, color=BT.WHITE_HEX,
                     align=PP_ALIGN.CENTER, wrap=False)
                # Value text
                _txb(slide, item.get("value", ""),
                     l=bx + Mm(5), t=by + Mm(13.5), w=ITEM_W - Mm(10), h=Mm(7),
                     sz=10, bold=True, color=BT.NEUTRAL_900_HEX)

        # Right dark panel — height is derived from item count so top == bottom padding
        PANEL_X   = Mm(172.8)
        PANEL_W   = SLIDE_W - PANEL_X - MR
        PANEL_T   = Mm(49.4)
        PANEL_PAD = Mm(9)       # equal top / bottom inner padding
        item_h    = Mm(15.5)    # row stride (dot + label + text)
        ITEM_TEXT_H = Mm(13.5)  # label(5mm) + gap(0.5mm) + text-box(8mm)

        _items = right_panel.get("items", [])[:3] if right_panel else []
        n      = len(_items)
        # content bottom (rel to panel top) = first-item offset + (n-1)*stride + single-item height
        content_h = PANEL_PAD + (max(n - 1, 0) * item_h) + ITEM_TEXT_H
        PANEL_H   = content_h + PANEL_PAD

        _rect(slide, l=PANEL_X, t=PANEL_T, w=PANEL_W, h=PANEL_H,
              fill=BT.NEUTRAL_900_HEX, radius_mm=BT.RADIUS_MD_MM)

        if right_panel and _items:
            _ICON_SZ   = Mm(12)    # outer halo circle diameter
            _INNER_SZ  = Mm(9)     # inner solid circle (icon bg)
            _ICON_IMG  = Mm(5.5)   # PNG icon size
            _ICON_L    = PANEL_X + Mm(8)
            for k, it in enumerate(_items):
                iy     = PANEL_T + k * item_h + PANEL_PAD
                accent = it.get("accent", BT.PRIMARY_500_HEX)
                icon   = it.get("icon", "")

                # Outer halo: white semi-transparent circle (frosted glow on dark)
                _halo_offset = (_ICON_SZ - _INNER_SZ) / 2
                _rect(slide, l=_ICON_L, t=iy, w=_ICON_SZ, h=_ICON_SZ,
                      fill=BT.WHITE_HEX, radius_mm=BT.RADIUS_PILL_MM, fill_alpha=12)

                # Inner circle: accent color, fully opaque — clear icon background
                _rect(slide, l=_ICON_L + _halo_offset, t=iy + _halo_offset,
                      w=_INNER_SZ, h=_INNER_SZ,
                      fill=accent, radius_mm=BT.RADIUS_PILL_MM)

                # Icon PNG — lookup same as module_grid
                if icon and _BRAND_ROOT:
                    _stem = os.path.splitext(icon)[0]
                    _candidates = [
                        os.path.join(_BRAND_ROOT, "assets", "icons", "png", _stem + ".png"),
                        os.path.join(_BRAND_ROOT, "assets", "icons", _stem + ".png"),
                    ]
                    _icon_path = next((c for c in _candidates if os.path.exists(c)), None)
                    if _icon_path:
                        _il = _ICON_L + (_ICON_SZ - _ICON_IMG) / 2
                        _it = iy + (_ICON_SZ - _ICON_IMG) / 2
                        slide.shapes.add_picture(_icon_path, int(_il), int(_it),
                                                 width=int(_ICON_IMG), height=int(_ICON_IMG))

                # Text column — offset right of icon circle
                _TEXT_L = _ICON_L + _ICON_SZ + Mm(4)
                _TEXT_W = PANEL_W - (_TEXT_L - PANEL_X) - Mm(4)
                # Item label
                _txb(slide, it.get("label", ""),
                     l=_TEXT_L, t=iy + Mm(0.5), w=_TEXT_W, h=Mm(5),
                     sz=7, bold=True, color=accent)
                # Item text
                _txb(slide, it.get("text", ""),
                     l=_TEXT_L, t=iy + Mm(5.5), w=_TEXT_W, h=Mm(8),
                     sz=13, bold=True, color=BT.WHITE_HEX, ls_pt=17)

        # Bottom cards row (brand tone / core focus chips)
        return slide

    # ── Feature Module Grid (Slide 15 style) ──────────────────────────────────

    def add_module_grid(self,
                        title: str,
                        modules: List[Dict],
                        subtitle: str = "",
                        label: str = "",
                        title_deco=None,
                        intro_text: str = "",
                        intro_label: str = "",
                        intro_flow: list = None):
        """
        Feature module grid, up to 8 modules. Default: 4-per-row (2×4).

        modules: [{
            "num":      "01",          # sequence number shown top-right
            "title":    "智能检查",
            "en":       "SMART INSPECTION",
            "bullets":  ["AI 一键创建检查表单", "..."],
            "accent":   BT.SUCCESS_HEX,       # EN label color  (optional, defaults cycle)
            "icon_bg":  BT.PRIMARY_100_HEX,    # icon square bg  (optional, defaults cycle)
            "featured": False,                  # True → dark bg, white text (flagship)
        }]

        intro_text / intro_label: when 5 or 7 modules are given, slot [0,0] becomes a
        summary text block and the column count adapts for a balanced layout:
          5 mods → 3-col  (row0: text+m0+m1,  row1: m2+m3+m4)    → 2+3
          7 mods → 4-col  (row0: text+m0+m1+m2, row1: m3+m4+m5+m6) → 3+4
        """
        # Paired semantically: each slot (icon_bg, accent) comes from the same color family.
        # icon_bg is always a *_100 / *_BG light color — never a *_500 saturated color.
        # Slot order cycles through 6 distinct palettes then wraps for decks with 7-8 modules.
        _ICON_DEFAULTS = [
            BT.PRIMARY_100_HEX,    BT.CARD_ORANGE_BG,   BT.NEUTRAL_100_HEX,   BT.SECONDARY_100_HEX,
            BT.CARD_TEAL_BG,       BT.CARD_PURPLE_BG,   BT.PRIMARY_100_HEX,   BT.CARD_ORANGE_BG,
        ]
        _ACCENT_DEFAULTS = [
            BT.PRIMARY_500_HEX,    BT.WARNING_HEX,      BT.SUCCESS_HEX,        BT.SECONDARY_500_HEX,
            BT.TEAL_HEX,           BT.PURPLE_HEX,       BT.PRIMARY_500_HEX,    BT.WARNING_HEX,
        ]
        # When caller sets accent but not icon_bg, derive icon_bg from this map
        # so the circular badge and the en label always share the same color family.
        _ACCENT_TO_ICON_BG = {
            BT.PRIMARY_500_HEX:   BT.PRIMARY_100_HEX,
            BT.SUCCESS_HEX:       BT.NEUTRAL_100_HEX,
            BT.SECONDARY_500_HEX: BT.SECONDARY_100_HEX,
            BT.WARNING_HEX:       BT.CARD_ORANGE_BG,
            BT.TEAL_HEX:          BT.CARD_TEAL_BG,
            BT.PURPLE_HEX:        BT.CARD_PURPLE_BG,
        }

        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        CELL_H  = Mm(60)
        COL_GAP = Mm(3.5)
        ROW_GAP = Mm(3.5)
        GRID_Y  = CONTENT_Y + Mm(3)

        # Featured (dark) budget — same policy as card layouts, driven by BT constants.
        _n_mods       = len(modules[:8])
        _feat_budget  = (
            BT.MAX_DARK_PER_SLIDE if _n_mods >= BT.MIN_CARDS_FOR_DARK else 0
        ) if BT.DARK_ACCENT_CARDS_ENABLED else 0
        _feat_used    = 0
        _modules_adj  = []
        for _m in modules[:8]:
            if _m.get("featured", False):
                if _feat_used < _feat_budget:
                    _modules_adj.append(_m)
                    _feat_used += 1
                else:
                    _m2 = {k: v for k, v in _m.items() if k != "featured"}
                    _modules_adj.append(_m2)
            else:
                _modules_adj.append(_m)

        # ── Column count + cell width: adapts when intro_text is provided ─────────
        # 5 mods + intro_text → 3-col (row0: text+m0+m1,    row1: m2+m3+m4)   2+3
        # 7 mods + intro_text → 4-col (row0: text+m0+m1+m2, row1: m3+m4+m5+m6) 3+4
        # all other counts    → 4-col standard grid
        _use_intro   = (bool(intro_text) or bool(intro_flow)) and (_n_mods in (5, 7))
        _slot_offset = 1 if _use_intro else 0
        if _use_intro and _n_mods == 5:
            COLS   = 3
            CELL_W = (CW - 2 * COL_GAP) // 3
        else:
            COLS   = 4
            CELL_W = Mm(74.1)

        if _use_intro:
            ts_x  = ML
            ts_y  = GRID_Y
            _pad  = Mm(5)
            _iw   = CELL_W - 2 * _pad
            cur_y = ts_y + _pad
            if intro_label:
                _txb(slide, intro_label, l=ts_x + _pad, t=cur_y,
                     w=_iw, h=Mm(5),
                     sz=7, bold=True, color=BT.PRIMARY_500_HEX)
                cur_y += Mm(7)
            if intro_text:
                _txb(slide, intro_text,
                     l=ts_x + _pad, t=cur_y,
                     w=_iw, h=Mm(18),
                     sz=10, color=BT.NEUTRAL_700_HEX, ls_pt=16)
                cur_y += Mm(20)
            if intro_flow:
                _flow_pills(slide, intro_flow,
                            x0=ts_x + _pad, y0=cur_y, max_w=_iw,
                            font_sz=8)

        for i, mod in enumerate(_modules_adj):
            slot     = i + _slot_offset
            col      = slot % COLS
            row      = slot // COLS
            cx       = ML + col * (CELL_W + COL_GAP)
            cy       = GRID_Y + row * (CELL_H + ROW_GAP)
            featured = mod.get("featured", False)
            accent   = mod.get("accent",  _ACCENT_DEFAULTS[slot % len(_ACCENT_DEFAULTS)])
            if "icon_bg" in mod:
                icon_bg = mod["icon_bg"]
            elif "accent" in mod:
                # auto-derive matching light bg from the custom accent so icon and en label share color family
                icon_bg = _ACCENT_TO_ICON_BG.get(mod["accent"], _ICON_DEFAULTS[slot % len(_ICON_DEFAULTS)])
            else:
                icon_bg = _ICON_DEFAULTS[slot % len(_ICON_DEFAULTS)]
            cell_bg  = BT.NEUTRAL_900_HEX if featured else BT.WHITE_HEX
            ttl_c    = BT.WHITE_HEX if featured else BT.NEUTRAL_900_HEX
            bul_c    = BT.NEUTRAL_200_HEX if featured else BT.NEUTRAL_700_HEX
            num_c    = accent if featured else BT.NEUTRAL_400_HEX

            # Use _card() so rounded corners (RADIUS_SM_MM) and border are applied.
            # White cells get a subtle border so they're visible against the slide bg.
            _card(slide, l=cx, t=cy, w=CELL_W, h=CELL_H, bg=cell_bg,
                  border=None if featured else BT.BORDER_DEFAULT_HEX,
                  radius_mm=BT.RADIUS_SM_MM)

            # Icon background — circular badge (pill radius on a square → perfect circle).
            _IBGL = cx + Mm(4.9)
            _IBGT = cy + Mm(5.0)
            _IBGW = Mm(8.5)
            _IBGH = Mm(8.5)
            _rect(slide, l=_IBGL, t=_IBGT, w=_IBGW, h=_IBGH,
                  fill=icon_bg, radius_mm=BT.RADIUS_PILL_MM)

            # Icon content — three modes:
            #   1. PNG/SVG file path (absolute): rendered as picture
            #   2. Short string / emoji: rendered as centered text without font
            #      override so the OS emoji/symbol font fallback activates
            #   3. Empty / None: shows the coloured background only (placeholder)
            #
            # Lookup order for non-absolute filenames (stem = name without extension):
            #   1. assets/icons/png/{stem}.png  ← pre-rasterized, used for PPTX
            #   2. assets/icons/svg/{stem}.svg  ← canonical (cannot embed in PPTX directly)
            #   3. assets/icons/{name}          ← legacy flat path
            #   4. Absolute / relative path as-is
            # Run assets/icons/render.mjs to regenerate PNGs from SVGs.
            icon_spec = mod.get("icon", "")
            if icon_spec:
                _icon_path = None
                # Try icon lookup candidates in priority order
                if _BRAND_ROOT:
                    _stem = os.path.splitext(icon_spec)[0]   # "sparkles" from "sparkles.png"
                    _candidates = [
                        os.path.join(_BRAND_ROOT, "assets", "icons", "png", _stem + ".png"),
                        os.path.join(_BRAND_ROOT, "assets", "icons", "svg", _stem + ".svg"),
                        os.path.join(_BRAND_ROOT, "assets", "icons", icon_spec),
                    ]
                    for _c in _candidates:
                        if os.path.exists(_c):
                            _icon_path = _c
                            break
                if not _icon_path and os.path.exists(icon_spec):
                    _icon_path = icon_spec

                if _icon_path:
                    # File-based icon: centered within the circular background.
                    # Inset = (bg_size - icon_size) / 2 on each axis → perfect center.
                    _ISZ = Mm(5.5)
                    _IL  = _IBGL + (_IBGW - _ISZ) / 2
                    _IT  = _IBGT + (_IBGH - _ISZ) / 2
                    slide.shapes.add_picture(
                        _icon_path, int(_IL), int(_IT),
                        width=int(_ISZ), height=int(_ISZ))
                else:
                    # Inline emoji / unicode character — do NOT set fonts so the
                    # system's emoji/symbol fallback can activate.
                    _ib = slide.shapes.add_textbox(
                        int(_IBGL), int(_IBGT),
                        int(_IBGW), int(_IBGH))
                    _itf = _ib.text_frame
                    _itf.word_wrap = False
                    _ip = _itf.paragraphs[0]
                    _ip.alignment = PP_ALIGN.CENTER
                    _irun = _ip.add_run()
                    _irun.text = icon_spec
                    _irun.font.size = Pt(12)
                    # Colour only for non-emoji (colour is ignored on system emoji)
                    _irun.font.color.rgb = _rgb(
                        BT.WHITE_HEX if featured else BT.NEUTRAL_700_HEX
                    )

            # Sequence number — top right of cell, single line, no wrap
            num_str = mod.get("num", f"{i+1:02d}")
            _txb(slide, num_str,
                 l=cx + CELL_W - Mm(14), t=cy + Mm(4.5), w=Mm(13), h=Mm(6),
                 sz=9, bold=True, color=num_c, align=PP_ALIGN.RIGHT, wrap=False)

            # Module title
            _txb(slide, mod.get("title", ""),
                 l=cx + Mm(4.9), t=cy + Mm(16.6), w=CELL_W - Mm(9), h=Mm(7),
                 sz=11, bold=True, color=ttl_c)

            # EN subtitle — small colored label below module title
            en = mod.get("en", "")
            if en:
                _txb(slide, en,
                     l=cx + Mm(4.9), t=cy + Mm(24),
                     w=int(CELL_W - Mm(9.8)), h=Mm(5.5),
                     sz=6, bold=True,
                     color=BT.WHITE_HEX if featured else accent)

            # Bullet points
            bullets = mod.get("bullets", [])
            if bullets:
                body = "\n".join(f"· {b}" for b in bullets)
                _txb(slide, body,
                     l=cx + Mm(4.9), t=cy + Mm(29.5),
                     w=CELL_W - Mm(9), h=CELL_H - Mm(33),
                     sz=7, color=bul_c, ls_pt=11, wrap=True)

        return slide

    # ── Closing Slide ─────────────────────────────────────────────────────────

    def add_closing(self,
                    slogan: str = "谢谢",
                    slogan_label: str = "LET'S BUILD THE FUTURE TOGETHER",
                    slogan_sub: str = "",
                    slogan_parts=None,
                    message: str = "",
                    contact: str = ""):
        """
        Closing slide — fully code-generated (dark bg + right 2×2 contact cards).

          slogan_label  small caps label above slogan
          slogan        main text — gradient when slogan_parts=None
          slogan_parts  [(text, hex_color), ...] for partial highlight; '\\n' = line break
          slogan_sub    supporting paragraph below accent line
        """
        if message and not slogan:
            slogan = message

        slide = self._new_slide()
        _set_slide_bg(slide, BT.NEUTRAL_900_HEX)

        # Logo (top-left)
        _add_logo_stacked(slide, reverse=True, l_mm=21, t_mm=11, h_mm=28)

        # Top-right label
        _txb(slide, "THANK YOU · CONTACT US",
             l=Mm(220), t=Mm(13), w=Mm(98), h=Mm(5),
             sz=8, color=BT.NEUTRAL_400_HEX, align=PP_ALIGN.RIGHT)

        # ── Left slogan area ──────────────────────────────────────────
        _SY = Mm(46)

        if slogan_label:
            _txb(slide, slogan_label,
                 l=ML, t=_SY, w=Mm(170), h=Mm(6),
                 sz=9, bold=True, color=BT.PRIMARY_500_HEX)
            _SY += Mm(9)

        _SLOGAN_T = _SY
        _SLOGAN_H = Mm(52)

        if slogan_parts:
            _txb_runs(slide, slogan_parts,
                      l=ML, t=_SLOGAN_T, w=Mm(170), h=_SLOGAN_H,
                      sz=42, bold=True, ls_pt=52)
        else:
            _txb_gradient(slide, slogan,
                          l=ML, t=_SLOGAN_T, w=Mm(170), h=_SLOGAN_H,
                          sz=42, bold=True, align=PP_ALIGN.LEFT,
                          stops=TITLE_GRADIENT, ls_pt=52)

        _SY = _SLOGAN_T + _SLOGAN_H + Mm(4)

        _rect(slide, l=ML, t=_SY, w=Mm(22), h=Mm(0.8), fill=BT.PRIMARY_500_HEX)
        _SY += Mm(6)

        if slogan_sub:
            _txb(slide, slogan_sub,
                 l=ML, t=_SY, w=Mm(155), h=Mm(20),
                 sz=10, color=BT.NEUTRAL_400_HEX, ls_pt=15, wrap=True)

        # ── Right 2×2 contact card grid ───────────────────────────────
        GRID_X = Mm(208.1)
        GRID_Y = Mm(46.2)
        CARD_W = Mm(52.9)
        CARD_H = Mm(56.4)
        H_GAP  = Mm(3.5)
        V_GAP  = Mm(3.6)

        _CARDS = [
            {"col": 0, "row": 0,
             "accent": BT.PRIMARY_500_HEX,   "cat": "SALES",
             "title": "销售对接",   "value": BT.COMPANY_PHONE,
             "sub": "电话 / 微信号", "icon": "monitor-smartphone"},
            {"col": 1, "row": 0,
             "accent": BT.SECONDARY_500_HEX, "cat": "EMAIL",
             "title": "合作邮箱",   "value": BT.COMPANY_EMAIL,
             "sub": "24 小时内回复", "icon": "mail"},
            {"col": 0, "row": 1,
             "accent": BT.PRIMARY_500_HEX,   "cat": "WEBSITE",
             "title": "公司官网",   "value": BT.COMPANY_WEBSITE,
             "sub": "访问了解更多",  "icon": "shield-check"},
            {"col": 1, "row": 1,
             "accent": BT.SECONDARY_500_HEX, "cat": "WECHAT",
             "title": "官方公众号",  "value": "未来方舟智能科技",
             "sub": "关注公众号",   "icon": "message-circle"},
        ]

        for cd in _CARDS:
            cx = GRID_X + cd["col"] * (CARD_W + H_GAP)
            cy = GRID_Y + cd["row"] * (CARD_H + V_GAP)
            ac = cd["accent"]

            _rect(slide, l=cx, t=cy, w=CARD_W, h=CARD_H,
                  fill=BT.WHITE_HEX, radius_mm=BT.RADIUS_SM_MM)

            DOT_S = Mm(6.3)
            DOT_L = cx + Mm(5.6)
            DOT_T = cy + Mm(5.0)
            _rect(slide, l=DOT_L, t=DOT_T, w=DOT_S, h=DOT_S,
                  fill=ac, radius_mm=BT.RADIUS_PILL_MM)

            _icon = cd.get("icon", "")
            if _icon and _BRAND_ROOT:
                _ip = os.path.join(_BRAND_ROOT, "assets", "icons", "png",
                                   os.path.splitext(_icon)[0] + ".png")
                if os.path.exists(_ip):
                    _ISZ = Mm(3.2)
                    slide.shapes.add_picture(
                        _ip,
                        int(DOT_L + (DOT_S - _ISZ) / 2),
                        int(DOT_T + (DOT_S - _ISZ) / 2),
                        width=int(_ISZ), height=int(_ISZ))

            _txb(slide, cd["cat"],
                 l=DOT_L + DOT_S + Mm(2), t=DOT_T + Mm(1.5),
                 w=CARD_W - DOT_S - Mm(9), h=Mm(4),
                 sz=7, bold=True, color=ac)

            _txb(slide, cd["title"],
                 l=DOT_L, t=cy + Mm(14),
                 w=CARD_W - Mm(7), h=Mm(6),
                 sz=10, bold=True, color=BT.NEUTRAL_900_HEX)

            SUB_PAD = Mm(4)
            SUB_L   = cx + SUB_PAD
            SUB_T   = cy + Mm(22)
            SUB_W   = CARD_W - 2 * SUB_PAD
            SUB_H   = CARD_H - SUB_T + cy - SUB_PAD
            _rect(slide, l=SUB_L, t=SUB_T, w=SUB_W, h=SUB_H,
                  fill=BT.NEUTRAL_100_HEX, radius_mm=BT.RADIUS_SM_MM)

            _txb(slide, cd["value"],
                 l=SUB_L + Mm(3), t=SUB_T + Mm(3),
                 w=SUB_W - Mm(6), h=SUB_H - Mm(10),
                 sz=9, bold=True, color=BT.NEUTRAL_900_HEX, wrap=True)

            if cd.get("sub"):
                _txb(slide, cd["sub"],
                     l=SUB_L + Mm(3), t=SUB_T + SUB_H - Mm(5.5),
                     w=SUB_W - Mm(6), h=Mm(4.5),
                     sz=7, color=BT.NEUTRAL_400_HEX)

        # Footer + bottom bar
        _txb(slide, f"© 2026 {BT.BRAND_NAME_EN} · {BT.BRAND_FULL_CN}",
             l=ML, t=SLIDE_H - Mm(11), w=Mm(260), h=Mm(5),
             sz=7, color=BT.NEUTRAL_400_HEX)
        _rect(slide, l=0, t=SLIDE_H - Mm(3.5), w=SLIDE_W, h=Mm(3.5),
              fill=BT.PRIMARY_500_HEX)
        return slide

    # ── Decorative Arrow Helper ───────────────────────────────────────────────

    def add_deco_arrow(self, slide, arrow_type: str,
                       l_mm: float, t_mm: float,
                       w_mm: float = 55.0,
                       series: str = "a"):
        """Place an arrow decoration linking text and image areas.

        arrow_type: "arrow1" — upward-right, for TOP-aligned text+image layouts
                    "arrow2" — downward-right, for BOTTOM-aligned text+image layouts
        l_mm, t_mm: top-left corner of the arrow (mm from slide top-left)
        w_mm:       desired width; height auto-calculated from image ratio
        series:     "a" (white bg) | "b" (gray/image bg) | "c" (dark bg)

        Typical usage:
            slide = prs.add_text_image_right(...)
            prs.add_deco_arrow(slide, "arrow1", l_mm=130, t_mm=70, w_mm=50)
        """
        return _deco_img(slide, arrow_type, l_mm, t_mm, w_mm, series=series)

    # ── Four-Column Card Row ──────────────────────────────────────────────────

    def add_four_cards(self,
                       title: str,
                       cards: List[Dict[str, str]],
                       subtitle: str = "",
                       label: str = "",
                       title_deco=None):
        """
        Single-row four-column card layout — for exactly 4 parallel items.

        Use this instead of add_six_cards when you have 4 items:
        add_six_cards with 4 items produces a 3+1 orphan layout.
        add_four_cards gives 4 equal-width tall cards in one balanced row.

        Each card shares the visual language of add_three_cards:
          tag (optional eyebrow) → title → accent bar → body text

        Color palette (cycles across 4 slots):
          0: PRIMARY_100  / PRIMARY_500    (green)
          1: NEUTRAL_100  / SUCCESS        (gray-green)
          2: SECONDARY_100/ SECONDARY_500  (lime)
          3: CARD_ORANGE_BG/ WARNING       (orange)

        cards: [{"title": "...", "body": "...", "tag": "(optional)"}]
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        n        = min(len(cards), 4)
        C4_GAP   = Mm(4)
        C4_W     = (CW - 3 * C4_GAP) // 4
        card_top = CONTENT_Y + Mm(5)
        card_h   = CONTENT_H - Mm(5)

        card_colors   = [BT.PRIMARY_100_HEX, BT.NEUTRAL_100_HEX,
                         BT.SECONDARY_100_HEX, BT.CARD_ORANGE_BG]
        accent_colors = [BT.PRIMARY_500_HEX, BT.SUCCESS_HEX,
                         BT.SECONDARY_500_HEX, BT.WARNING_HEX]

        for i, c in enumerate(cards[:4]):
            x         = ML + i * (C4_W + C4_GAP)
            bg        = card_colors[i % 4]
            acc       = accent_colors[i % 4]
            _card(slide, l=x, t=card_top, w=C4_W, h=card_h, bg=bg)

            layout = c.get("layout", "vertical_stack")
            if layout != "vertical_stack":
                _render_card_inner(slide, layout, x, card_top, C4_W, card_h,
                                   c, acc, pad_s=Mm(4), pad_t=Mm(5))
                continue

            inner_l = x + Mm(4)
            inner_w = C4_W - Mm(8)
            y_off   = card_top + Mm(5)

            tag = c.get("tag", "")
            if tag:
                _txb(slide, tag, l=inner_l, t=y_off, w=inner_w, h=Mm(6),
                     sz=8, bold=True, color=acc)
                y_off += Mm(7)

            _txb(slide, c.get("title", ""),
                 l=inner_l, t=y_off, w=inner_w, h=Mm(22),
                 sz=14, bold=True, color=BT.NEUTRAL_900_HEX)
            y_off += Mm(23)

            _rect(slide, l=inner_l, t=y_off, w=Mm(20), h=Mm(1), fill=acc)
            y_off += Mm(5)

            body_tb = _txb(slide, c.get("body", ""),
                           l=inner_l, t=y_off, w=inner_w,
                           h=card_h - (y_off - card_top) - Mm(6),
                           sz=12, color=BT.NEUTRAL_700_HEX, ls_pt=17)
            body_tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        return slide

    # ── Numbered Rows (problem / reason / defect list) ───────────────────────

    def add_numbered_rows(self,
                          title: str,
                          items: list,
                          subtitle: str = "",
                          label: str = "",
                          title_deco=None,
                          note: str = "",
                          note_style: str = "note"):
        """
        3-item numbered horizontal rows — for problem/reason/defect lists.

        Each item renders as a full-width split-panel card:
          Left  (~28%): accent-colored bg — large number badge + item title
          Right (~72%): light-colored bg  — body text (auto-shrinks to fit)

        items: [{"num": "01", "title": "...", "body": "...", "tag": "(optional eyebrow)"}]

        Color palette per row (0→2 cycles):
          0: PRIMARY_500 left / PRIMARY_100 right
          1: SUCCESS     left / NEUTRAL_100 right
          2: SECONDARY_500 left / SECONDARY_100 right  (dark text variant)
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        n = min(len(items), 3)
        if n == 0:
            return slide

        _note_reserve = CALLOUT_H + Mm(5) if note else 0
        TOP_PAD  = Mm(4)
        ROW_GAP  = Mm(3)
        row_h    = int((CONTENT_H - TOP_PAD - _note_reserve - ROW_GAP * (n - 1)) / n)

        LEFT_W   = int(CW * 0.28)
        CARD_GAP = Mm(3)
        RIGHT_W  = CW - LEFT_W - CARD_GAP

        # (accent_bg, light_bg, left_text_color)
        _PALETTES = [
            (BT.PRIMARY_500_HEX,   BT.PRIMARY_100_HEX,   BT.WHITE_HEX),
            (BT.SUCCESS_HEX,       BT.NEUTRAL_100_HEX,   BT.WHITE_HEX),
            (BT.SECONDARY_500_HEX, BT.SECONDARY_100_HEX, BT.NEUTRAL_900_HEX),
        ]

        for i, item in enumerate(items[:3]):
            acc_bg, light_bg, acc_txt = _PALETTES[i % 3]
            y = CONTENT_Y + TOP_PAD + i * (row_h + ROW_GAP)

            # Left accent card
            _card(slide, l=ML, t=y, w=LEFT_W, h=row_h, bg=acc_bg)

            # Right light card
            _card(slide, l=ML + LEFT_W + CARD_GAP, t=y, w=RIGHT_W, h=row_h, bg=light_bg)

            # Number badge pill (top-left of left panel)
            num_text  = item.get("num", f"{i + 1:02d}")
            BADGE_H   = Mm(7)
            BADGE_W   = Mm(14)
            badge_bg  = BT.WHITE_HEX if acc_txt == BT.WHITE_HEX else acc_bg
            badge_fg  = acc_bg       if acc_txt == BT.WHITE_HEX else BT.WHITE_HEX
            _rect(slide, l=ML + Mm(5), t=y + Mm(5),
                  w=BADGE_W, h=BADGE_H,
                  fill=badge_bg, radius_mm=BT.RADIUS_PILL_MM)
            _txb(slide, num_text,
                 l=ML + Mm(5), t=y + Mm(5),
                 w=BADGE_W, h=BADGE_H,
                 sz=9, bold=True, color=badge_fg,
                 align=PP_ALIGN.CENTER, wrap=False)

            # Optional tag eyebrow
            tag = item.get("tag", "")
            y_inner = y + Mm(14)
            if tag:
                _txb(slide, tag, l=ML + Mm(5), t=y_inner,
                     w=LEFT_W - Mm(10), h=Mm(5),
                     sz=7, bold=True,
                     color=BT.WHITE_HEX if acc_txt == BT.WHITE_HEX else BT.NEUTRAL_700_HEX)
                y_inner += Mm(6)

            # Item title in left panel
            item_title_h = row_h - (y_inner - y) - Mm(5)
            tb = _txb(slide, item.get("title", ""),
                      l=ML + Mm(5), t=y_inner,
                      w=LEFT_W - Mm(10), h=item_title_h,
                      sz=14, bold=True, color=acc_txt, ls_pt=20)
            tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            # Body text in right panel (arrows auto-replaced with brand PNG)
            _render_content_with_arrows(
                slide, item.get("body", ""),
                x=ML + LEFT_W + CARD_GAP + Mm(5),
                y=y + Mm(5),
                w=RIGHT_W - Mm(10),
                h=row_h - Mm(10),
                sz=13, color=BT.NEUTRAL_700_HEX, ls_pt=18,
                arrow_style="primary", arrow_size_mm=4.5)

        if note:
            _note_t = CONTENT_Y + CONTENT_H - CALLOUT_H - Mm(4)
            _callout(slide, note, l=ML, t=_note_t, w=CW, style=note_style)

        return slide

    # ── Accent Rows (horizontal data / spec comparison) ───────────────────────

    def add_accent_rows(self,
                        title: str,
                        items: list,
                        subtitle: str = "",
                        label: str = "",
                        title_deco=None,
                        note: str = "",
                        note_style: str = "note"):
        """
        Horizontal spec-comparison rows with left accent bar — for data/ledger/type lists.

        Each item is a full-width rounded card with:
          • Narrow 2.5 mm accent bar on the inside-left edge
          • Header row: bold title (left) + optional badge pill (right-aligned)
          • Body area: body text with auto-size

        items: [{
            "tag":   "ACCOUNT TYPE",    # small accent-colored label above title
            "title": "账本名称",        # bold 14pt title
            "body":  "描述文字...",     # body text (auto-shrinks to fit)
            "badge": "confidence=1.0",  # right-aligned badge pill (optional)
        }]

        Color palette per row (0→2):
          0: PRIMARY_500   / PRIMARY_100
          1: TEAL          / CARD_TEAL_BG
          2: SUCCESS       / NEUTRAL_100
        """
        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        n = min(len(items), 3)
        if n == 0:
            return slide

        _note_reserve = CALLOUT_H + Mm(5) if note else 0
        TOP_PAD  = Mm(4)
        ROW_GAP  = Mm(3)
        row_h    = int((CONTENT_H - TOP_PAD - _note_reserve - ROW_GAP * (n - 1)) / n)

        STRIPE_W = Mm(2.5)
        CONTENT_L = ML + STRIPE_W + Mm(5)   # text starts after stripe + gap

        _PALETTES = [
            (BT.PRIMARY_500_HEX,   BT.PRIMARY_100_HEX),
            (BT.TEAL_HEX,          BT.CARD_TEAL_BG),
            (BT.SUCCESS_HEX,       BT.NEUTRAL_100_HEX),
        ]

        for i, item in enumerate(items[:3]):
            acc, light_bg = _PALETTES[i % 3]
            y = CONTENT_Y + TOP_PAD + i * (row_h + ROW_GAP)
            inner_w = CW - STRIPE_W - Mm(15)   # usable text width

            # Card background (full-width, rounded)
            _card(slide, l=ML, t=y, w=CW, h=row_h, bg=light_bg)

            # Left accent bar (inset inside card, avoids covering rounded corners)
            _rect(slide, l=ML + Mm(3), t=y + Mm(3),
                  w=STRIPE_W, h=row_h - Mm(6),
                  fill=acc, radius_mm=BT.RADIUS_PILL_MM)

            # Optional tag label (small, accent-colored)
            tag = item.get("tag", "")
            y_inner = y + Mm(5)
            if tag:
                _txb(slide, tag, l=CONTENT_L, t=y_inner,
                     w=Mm(80), h=Mm(5.5),
                     sz=7, bold=True, color=acc)
                y_inner += Mm(6.5)

            # Title + optional right-aligned badge on same baseline
            title_text = item.get("title", "")
            _txb(slide, title_text,
                 l=CONTENT_L, t=y_inner,
                 w=inner_w - Mm(40), h=Mm(8),
                 sz=14, bold=True, color=BT.NEUTRAL_900_HEX)

            badge_text = item.get("badge", "")
            if badge_text:
                _pill(slide, badge_text,
                      l=ML + CW - Mm(5) - int(len(badge_text) * Mm(9 * 0.44) + 2 * Mm(3.5)),
                      t=y_inner + Mm(0.5),
                      bg=acc, font_sz=8, h=Mm(7))

            y_inner += Mm(10)

            # Body text
            body_h = row_h - (y_inner - y) - Mm(5)
            if body_h > Mm(8):
                _render_content_with_arrows(
                    slide, item.get("body", ""),
                    x=CONTENT_L, y=y_inner,
                    w=inner_w, h=body_h,
                    sz=12, color=BT.NEUTRAL_700_HEX, ls_pt=17,
                    arrow_style="primary", arrow_size_mm=4.5)

        if note:
            _note_t = CONTENT_Y + CONTENT_H - CALLOUT_H - Mm(4)
            _callout(slide, note, l=ML, t=_note_t, w=CW, style=note_style)

        return slide

    # ── Strip List ────────────────────────────────────────────────────────────

    def add_strip_list(self,
                       title: str,
                       items: list,
                       strip_style: str = "accent",
                       subtitle: str = "",
                       label: str = "",
                       title_deco=None,
                       note: str = "",
                       note_style: str = "note"):
        """
        Present items as full-width horizontal strip cards stacked vertically.

        strip_style:
          "accent"    — light bg + inset vertical accent stripe, title bold, body below,
                        optional pill tag floats right
          "numbered"  — white bg + grey border, circle number badge on left, title + body
          "tag_left"  — light bg, pill tag anchored left, title + body to the right

        Item dict keys:
          "title"  : str | bilingual dict  (required)
          "body"   : str | bilingual dict  (optional description)
          "tag"    : str                   (pill label; accent → floats right,
                                           tag_left → anchors left as category pill)
          "color"  : semantic key          ("primary","neutral","secondary","warning",…)
          "num"    : explicit badge label override for numbered style (default "01","02",…)
        """
        from scripts.components.base import resolve_card_color
        from scripts.i18n import T

        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        _note_reserve = CALLOUT_H + Mm(6) if note else 0
        content_top   = CONTENT_Y + Mm(4)
        available_h   = CONTENT_H - Mm(4) - _note_reserve

        n = max(1, len(items))
        STRIP_GAP = Mm(4)
        strip_h = min(Mm(22), max(Mm(13), int((available_h - (n - 1) * STRIP_GAP) // n)))

        cur_y = content_top
        for idx, item in enumerate(items):
            item_title = T(item.get("title", ""))
            item_body  = T(item.get("body",  ""))
            item_tag   = T(item.get("tag",   ""))
            color_key  = item.get("color", "primary")
            num_label  = item.get("num", f"{idx + 1:02d}")
            bg, acc, txt = resolve_card_color(color_key)

            if strip_style == "accent":
                # ── Card bg ───────────────────────────────────────────────────
                _card(slide, l=ML, t=cur_y, w=CW, h=strip_h, bg=bg)
                # Inset vertical accent stripe (stays within card's rounded area)
                _rect(slide, l=ML + Mm(2), t=cur_y + Mm(2),
                      w=Mm(3), h=strip_h - Mm(4), fill=acc, radius_mm=0.5)

                PAD_L     = Mm(9)    # stripe(2+3) + gap(4)
                PAD_T     = Mm(4)
                PAD_R     = Mm(6)
                content_w = CW - PAD_L - PAD_R

                # Optional tag pill — float right, vertically centred
                tag_reserve = 0
                if item_tag:
                    est_tag_w   = int(len(item_tag) * Mm(8 * 0.44) + 2 * Mm(3.5))
                    tag_t       = cur_y + (strip_h - Mm(8)) // 2
                    tag_l       = ML + CW - PAD_R - est_tag_w
                    _pill(slide, item_tag, l=tag_l, t=tag_t,
                          bg=acc, fg=_pill_fg(acc), font_sz=8, h=Mm(8))
                    tag_reserve = est_tag_w + Mm(4)

                title_h = Mm(8)
                body_h  = strip_h - PAD_T - title_h - Mm(2) - Mm(3)

                if item_title:
                    _txb(slide, item_title,
                         l=ML + PAD_L, t=cur_y + PAD_T,
                         w=content_w - tag_reserve, h=title_h,
                         sz=13, bold=True, color=txt)
                if item_body and body_h > Mm(4):
                    tb = _txb(slide, item_body,
                              l=ML + PAD_L, t=cur_y + PAD_T + title_h + Mm(2),
                              w=content_w, h=body_h,
                              sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=16)
                    try:
                        tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    except Exception:
                        pass

            elif strip_style == "numbered":
                # ── White card with border ────────────────────────────────────
                _card(slide, l=ML, t=cur_y, w=CW, h=strip_h,
                      bg=BT.WHITE_HEX, border=BT.NEUTRAL_200_HEX)
                # Circle badge
                BADGE_D = Mm(12)
                badge_l = ML + Mm(5)
                badge_t = cur_y + (strip_h - BADGE_D) // 2
                _rect(slide, l=badge_l, t=badge_t, w=BADGE_D, h=BADGE_D,
                      fill=acc, radius_mm=BT.RADIUS_PILL_MM)
                _txb(slide, num_label,
                     l=badge_l, t=badge_t + Mm(1.5), w=BADGE_D, h=BADGE_D - Mm(3),
                     sz=10, bold=True, color=BT.WHITE_HEX,
                     align=PP_ALIGN.CENTER, wrap=False)

                CONTENT_L = badge_l + BADGE_D + Mm(5)
                content_w = ML + CW - CONTENT_L - Mm(5)
                title_h   = Mm(7)
                body_h    = strip_h - Mm(4) - title_h - Mm(2) - Mm(3)
                y_title   = cur_y + (Mm(3) if item_body else int((strip_h - title_h) // 2))

                if item_title:
                    _txb(slide, item_title,
                         l=CONTENT_L, t=y_title, w=content_w, h=title_h,
                         sz=13, bold=True, color=BT.NEUTRAL_900_HEX)
                if item_body and body_h > Mm(4):
                    tb = _txb(slide, item_body,
                              l=CONTENT_L, t=y_title + title_h + Mm(2),
                              w=content_w, h=body_h,
                              sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=16)
                    try:
                        tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    except Exception:
                        pass

            elif strip_style == "tag_left":
                # ── Light bg card with left pill tag ─────────────────────────
                _card(slide, l=ML, t=cur_y, w=CW, h=strip_h, bg=bg)
                TAG_PAD   = Mm(5)
                tag_label = item_tag or (item_title[:10] if item_title else "")
                est_tag_w = int(len(tag_label) * Mm(9 * 0.44) + 2 * Mm(4))
                tag_t     = cur_y + (strip_h - Mm(8)) // 2
                _pill(slide, tag_label, l=ML + TAG_PAD, t=tag_t,
                      bg=acc, fg=_pill_fg(acc), font_sz=9, h=Mm(8))

                CONTENT_L = ML + TAG_PAD + est_tag_w + Mm(5)
                content_w = ML + CW - CONTENT_L - Mm(5)
                title_h   = Mm(7)
                body_h    = strip_h - Mm(4) - title_h - Mm(2) - Mm(3)
                y_title   = cur_y + (Mm(3) if item_body else int((strip_h - title_h) // 2))

                if item_title:
                    _txb(slide, item_title,
                         l=CONTENT_L, t=y_title, w=content_w, h=title_h,
                         sz=13, bold=True, color=txt)
                if item_body and body_h > Mm(4):
                    tb = _txb(slide, item_body,
                              l=CONTENT_L, t=y_title + title_h + Mm(2),
                              w=content_w, h=body_h,
                              sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=16)
                    try:
                        tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    except Exception:
                        pass

            cur_y += strip_h + STRIP_GAP

        if note:
            _note_t = CONTENT_Y + CONTENT_H - CALLOUT_H - Mm(4)
            _callout(slide, note, l=ML, t=_note_t, w=CW, style=note_style)

        return slide

    def add_two_panel_strips(self,
                              title: str,
                              left_label: str,
                              left_items: list,
                              right_label: str,
                              right_items: list,
                              left_style: str = "accent",
                              right_style: str = "numbered",
                              subtitle: str = "",
                              label: str = "",
                              title_deco=None,
                              note: str = "",
                              note_style: str = "note"):
        """
        Single slide with two side-by-side panel containers, each holding
        stacked strip cards.  Ideal for presenting two related lists on
        one page (e.g. methods + workflow steps).

        left_style / right_style: same values as add_strip_list
            "accent"    — coloured bg + inset accent stripe, title + body + optional tag
            "numbered"  — numbered circle badge, title (+ body if strip tall enough)
            "tag_left"  — left pill tag + title + body

        Item dict keys: same as add_strip_list (title, body, tag, color, num).
        """
        from scripts.components.base import resolve_card_color
        from scripts.i18n import T

        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label, title_deco=title_deco)
        _footer(slide)

        _note_reserve = CALLOUT_H + Mm(6) if note else 0

        # ── Panel layout constants ────────────────────────────────────────────
        PANEL_PAD = Mm(5)
        LABEL_H   = Mm(8)
        LABEL_GAP = Mm(4)

        panel_t = CONTENT_Y + Mm(4)
        panel_h = CONTENT_H - Mm(8) - _note_reserve
        # Height available for strips inside each panel (after label + padding)
        strips_avail = panel_h - PANEL_PAD - LABEL_H - LABEL_GAP - PANEL_PAD

        # ── Calculate per-strip heights ───────────────────────────────────────
        n_left    = max(1, len(left_items))
        L_GAP     = Mm(3)
        l_strip_h = int((strips_avail - (n_left - 1) * L_GAP) // n_left)

        n_right   = max(1, len(right_items))
        R_GAP     = Mm(3) if n_right <= 3 else Mm(2)
        r_strip_h = int((strips_avail - (n_right - 1) * R_GAP) // n_right)

        # ── Render each panel ─────────────────────────────────────────────────
        for side in ("left", "right"):
            is_left = (side == "left")
            panel_l  = ML if is_left else ML + C2_W + C2_GAP
            items    = left_items  if is_left else right_items
            slabel   = left_label  if is_left else right_label
            style    = left_style  if is_left else right_style
            n        = n_left      if is_left else n_right
            strip_h  = l_strip_h  if is_left else r_strip_h
            gap      = L_GAP      if is_left else R_GAP

            # Left panel: white bg + subtle border; right panel: NEUTRAL_100 tray bg
            if is_left:
                _card(slide, l=panel_l, t=panel_t, w=C2_W, h=panel_h,
                      bg=BT.WHITE_HEX, border=BT.NEUTRAL_200_HEX)
            else:
                _card(slide, l=panel_l, t=panel_t, w=C2_W, h=panel_h,
                      bg=BT.NEUTRAL_100_HEX, border=None)

            # Section label — plain bold text, no pill
            _txb(slide, slabel,
                 l=panel_l + PANEL_PAD,
                 t=panel_t + PANEL_PAD,
                 w=C2_W - 2 * PANEL_PAD, h=LABEL_H,
                 sz=12, bold=True, color=BT.NEUTRAL_900_HEX,
                 wrap=False)

            strip_w = C2_W - 2 * PANEL_PAD
            cur_y   = panel_t + PANEL_PAD + LABEL_H + LABEL_GAP

            for idx, item in enumerate(items):
                item_title = T(item.get("title", ""))
                item_body  = T(item.get("body",  ""))
                item_tag   = T(item.get("tag",   ""))
                color_key  = item.get("color", "primary")
                num_label  = item.get("num", f"{idx + 1:02d}")
                bg_col, acc, txt = resolve_card_color(color_key)
                strip_l = panel_l + PANEL_PAD

                if style == "accent":
                    _card(slide, l=strip_l, t=cur_y, w=strip_w, h=strip_h, bg=bg_col)
                    _rect(slide, l=strip_l + Mm(2), t=cur_y + Mm(2),
                          w=Mm(3), h=strip_h - Mm(4), fill=acc, radius_mm=0.5)

                    PAD_L     = Mm(9)
                    PAD_T     = Mm(4)
                    PAD_R     = Mm(5)
                    content_w = strip_w - PAD_L - PAD_R

                    tag_reserve = 0
                    if item_tag:
                        est_tag_w  = int(len(item_tag) * Mm(8 * 0.44) + 2 * Mm(3.5))
                        tag_t      = cur_y + (strip_h - Mm(8)) // 2
                        tag_l      = strip_l + strip_w - PAD_R - est_tag_w
                        _pill(slide, item_tag, l=tag_l, t=tag_t,
                              bg=acc, fg=_pill_fg(acc), font_sz=8, h=Mm(8))
                        tag_reserve = est_tag_w + Mm(4)

                    title_h = Mm(8)
                    body_h  = strip_h - PAD_T - title_h - Mm(2) - Mm(3)

                    if item_title:
                        _txb(slide, item_title,
                             l=strip_l + PAD_L, t=cur_y + PAD_T,
                             w=content_w - tag_reserve, h=title_h,
                             sz=13, bold=True, color=txt)
                    if item_body and body_h > Mm(4):
                        tb = _txb(slide, item_body,
                                  l=strip_l + PAD_L, t=cur_y + PAD_T + title_h + Mm(2),
                                  w=content_w, h=body_h,
                                  sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=16)
                        try:
                            tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        except Exception:
                            pass

                elif style == "numbered":
                    # White card on NEUTRAL_100 panel bg → clear visual contrast
                    _card(slide, l=strip_l, t=cur_y, w=strip_w, h=strip_h,
                          bg=BT.WHITE_HEX, border=BT.NEUTRAL_200_HEX)
                    BADGE_D = Mm(10)
                    badge_l = strip_l + Mm(4)
                    badge_t = cur_y + (strip_h - BADGE_D) // 2
                    _rect(slide, l=badge_l, t=badge_t, w=BADGE_D, h=BADGE_D,
                          fill=BT.PRIMARY_500_HEX, radius_mm=BT.RADIUS_PILL_MM)
                    _txb(slide, num_label,
                         l=badge_l, t=badge_t + Mm(1.2), w=BADGE_D, h=BADGE_D - Mm(2.4),
                         sz=9, bold=True, color=BT.WHITE_HEX,
                         align=PP_ALIGN.CENTER, wrap=False)

                    CONTENT_L = badge_l + BADGE_D + Mm(4)
                    content_w = strip_l + strip_w - CONTENT_L - Mm(4)
                    title_h   = Mm(7)
                    body_h    = strip_h - Mm(4) - title_h - Mm(2) - Mm(3)
                    y_title   = (cur_y + int((strip_h - title_h) // 2)
                                 if not item_body or body_h < Mm(4)
                                 else cur_y + Mm(3))

                    if item_title:
                        _txb(slide, item_title,
                             l=CONTENT_L, t=y_title, w=content_w, h=title_h,
                             sz=12, bold=True, color=BT.NEUTRAL_900_HEX)
                    if item_body and body_h > Mm(4):
                        tb = _txb(slide, item_body,
                                  l=CONTENT_L, t=y_title + title_h + Mm(2),
                                  w=content_w, h=body_h,
                                  sz=10, color=BT.NEUTRAL_700_HEX, ls_pt=15)
                        try:
                            tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        except Exception:
                            pass

                elif style == "tag_left":
                    _card(slide, l=strip_l, t=cur_y, w=strip_w, h=strip_h, bg=bg_col)
                    TAG_PAD   = Mm(4)
                    tag_label = item_tag or (item_title[:10] if item_title else "")
                    est_tag_w = int(len(tag_label) * Mm(9 * 0.44) + 2 * Mm(4))
                    tag_t     = cur_y + (strip_h - Mm(8)) // 2
                    _pill(slide, tag_label, l=strip_l + TAG_PAD, t=tag_t,
                          bg=acc, fg=_pill_fg(acc), font_sz=9, h=Mm(8))

                    CONTENT_L = strip_l + TAG_PAD + est_tag_w + Mm(4)
                    content_w = strip_l + strip_w - CONTENT_L - Mm(4)
                    title_h   = Mm(7)
                    body_h    = strip_h - Mm(4) - title_h - Mm(2) - Mm(3)
                    y_title   = (cur_y + int((strip_h - title_h) // 2)
                                 if not item_body or body_h < Mm(4)
                                 else cur_y + Mm(3))

                    if item_title:
                        _txb(slide, item_title,
                             l=CONTENT_L, t=y_title, w=content_w, h=title_h,
                             sz=12, bold=True, color=txt)
                    if item_body and body_h > Mm(4):
                        tb = _txb(slide, item_body,
                                  l=CONTENT_L, t=y_title + title_h + Mm(2),
                                  w=content_w, h=body_h,
                                  sz=10, color=BT.NEUTRAL_700_HEX, ls_pt=15)
                        try:
                            tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        except Exception:
                            pass

                cur_y += strip_h + gap

        if note:
            _note_t = CONTENT_Y + CONTENT_H - CALLOUT_H - Mm(4)
            _callout(slide, note, l=ML, t=_note_t, w=CW, style=note_style)

        return slide

    # ── Grid-first slide (add_slide) ─────────────────────────────────────────

    def add_slide(self,
                  title,
                  label: str = "",
                  subtitle: str = "",
                  title_deco=None,
                  grid: str = "full",
                  slots: list = None,
                  note: str = "",
                  note_style: str = "note",
                  # ── backwards compat (component-first path) ───────────────
                  layout=None,
                  slide_label: str = ""):
        """
        Grid-first slide builder.

        Grid mode (new — preferred):
            grid       — preset name: "full" | "2col" | "3col" | "2row" | "3row"
                         "top_2col" | "2col_bot" | "top_3col"
                         "left_2row" | "2row_right"
            slots      — list of container dicts, one per grid slot,
                         in reading order (left→right, top→bottom).
                         Each dict must have "type" plus type-specific keys.
            note       — optional bottom callout text (auto height, outside grid)
            note_style — "note" | "tip" | "warning" | "danger"

        Component-first mode (legacy, backwards compat):
            layout     — a BaseLayout instance (Grid / Stack / Split)

        Container types currently available in slots:
            "text"        — paragraph + bullet list
            "placeholder" — dashed box for drafting
            "blank"       — empty spacer
            More types added in step 3 (cards, strips, stats, flow, image, tags).
        """
        # ── Legacy component-first path ───────────────────────────────────────
        if layout is not None:
            from pptx.util import Mm as _Mm
            slide = self._new_slide()
            _set_slide_bg(slide, BT.WHITE_HEX)
            _header(slide, title, subtitle=subtitle, label=label,
                    title_deco=title_deco)
            _footer(slide, layout_label=slide_label)
            note_reserve = CALLOUT_H + _Mm(6) if note else 0
            cy = CONTENT_Y + _Mm(4)
            ch = CONTENT_H - _Mm(4) - note_reserve
            layout.render_pptx(slide, int(ML), int(cy), int(CW), int(ch))
            if note:
                nt = CONTENT_Y + CONTENT_H - CALLOUT_H - _Mm(4)
                if hasattr(note, "render_pptx"):
                    note.render_pptx(slide, int(ML), int(nt), int(CW),
                                     int(CALLOUT_H))
                else:
                    _callout(slide, str(note), l=ML, t=nt, w=CW,
                             style=note_style)
            return slide

        # ── Grid-first path ───────────────────────────────────────────────────
        if slots is None:
            slots = []

        if grid not in GRID_PRESETS:
            raise ValueError(
                f"Unknown grid preset '{grid}'. "
                f"Available: {sorted(GRID_PRESETS)}"
            )
        expected = GRID_SLOT_COUNT[grid]
        if len(slots) != expected:
            raise ValueError(
                f"Grid '{grid}' expects {expected} slot(s), got {len(slots)}. "
                f"Slots are in reading order: left→right, top→bottom."
            )

        slide = self._new_slide()
        _set_slide_bg(slide, BT.WHITE_HEX)
        _header(slide, title, subtitle=subtitle, label=label,
                title_deco=title_deco)
        _footer(slide)

        note_reserve = CALLOUT_H + Mm(6) if note else 0
        ax = ML
        ay = CONTENT_Y + Mm(4)
        aw = CW
        ah = CONTENT_H - Mm(4) - note_reserve

        slot_rects = GRID_PRESETS[grid](ax, ay, aw, ah)
        for slot_dict, (sl, st, sw, sh) in zip(slots, slot_rects):
            _render_container(slide, slot_dict, sl, st, sw, sh)

        if note:
            note_t = CONTENT_Y + CONTENT_H - CALLOUT_H - Mm(4)
            _callout(slide, note, l=ML, t=note_t, w=CW, style=note_style)

        return slide

    # ── Save ─────────────────────────────────────────────────────────────────

    def save(self, output_path: str):
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        self._prs.save(output_path)
        print(f"✓ Saved: {output_path}")

    @property
    def presentation(self):
        return self._prs
