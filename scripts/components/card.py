"""Card component — renders a single branded card for Grid/Stack layouts."""
from __future__ import annotations
from dataclasses import dataclass, field
from pptx.util import Mm
from scripts.components.base import BaseComponent, resolve_card_color
import scripts.brand_tokens as BT


@dataclass
class Card(BaseComponent):
    title: object = ""    # str | {"cn": ..., "en": ...}
    body:  object = ""
    tag:   object = ""
    color: str    = "primary"
    layout: str   = "vertical_stack"
    # horizontal_split
    metric: object = ""
    # icon_left
    icon: object = ""
    # quote_card
    quote:       object = ""
    attribution: object = ""
    # override
    dark: bool = False

    def _colors(self):
        if self.dark:
            return BT.NEUTRAL_900_HEX, BT.SECONDARY_500_HEX, BT.WHITE_HEX
        return resolve_card_color(self.color)

    def render_pptx(self, slide, x: int, y: int, w: int, h: int) -> None:
        from scripts.pptx_builder import _card, _render_card_inner, _txb, _rect
        from scripts.i18n import T, is_bilingual, get_cn, get_en, \
            CN_COLOR, EN_COLOR, EN_TITLE_COLOR, EN_SZ_RATIO
        try:
            from pptx.util import MSO_AUTO_SIZE
        except ImportError:
            from pptx.enum.text import MSO_AUTO_SIZE

        bg, acc, txt = self._colors()
        body_color_default = (BT.NEUTRAL_700_HEX if bg != BT.NEUTRAL_900_HEX
                              else BT.NEUTRAL_100_HEX)
        border = None if bg != BT.WHITE_HEX else BT.BORDER_DEFAULT_HEX
        _card(slide, l=x, t=y, w=w, h=h, bg=bg, border=border)

        if self.layout != "vertical_stack":
            data = {
                "title":       T(self.title),
                "body":        T(self.body),
                "tag":         T(self.tag),
                "metric":      T(self.metric),
                "icon":        T(self.icon),
                "quote":       T(self.quote),
                "attribution": T(self.attribution),
            }
            _render_card_inner(slide, self.layout, x, y, w, h, data, acc,
                               txt_color=txt, body_color=body_color_default,
                               pad_s=Mm(5), pad_t=Mm(5))
            return

        # ── vertical_stack ────────────────────────────────────────────────────
        PAD_S = Mm(5)
        PAD_T = Mm(5)
        PAD_B = Mm(6)
        inner_l = x + PAD_S
        inner_w = w - 2 * PAD_S
        available = h - PAD_T - PAD_B

        has_tag  = bool(T(self.tag))
        tag_h    = Mm(6)  if has_tag  else 0
        tag_gap  = Mm(3)  if has_tag  else 0
        bar_h    = Mm(1)
        bar_gap  = Mm(5)

        after_tag    = available - tag_h - tag_gap
        title_budget = min(Mm(20), max(Mm(8), int(after_tag * 0.30)))
        bar_space    = bar_h + bar_gap if after_tag - title_budget > Mm(12) else 0
        body_h       = after_tag - title_budget - bar_space

        # ── bilingual: split body_h between CN and EN ─────────────────────────
        bilingual    = is_bilingual()
        cn_body_text = get_cn(self.body) if bilingual else T(self.body)
        en_body_text = get_en(self.body) if bilingual else ""
        has_en_body  = bilingual and bool(en_body_text)

        if has_en_body and body_h > Mm(8):
            cn_body_h = int(body_h * 0.58)
            en_body_h = body_h - cn_body_h - Mm(1)
        else:
            cn_body_h = body_h
            en_body_h = 0

        y_off = y + PAD_T

        # tag (CN only — short label, not worth duplicating)
        if has_tag:
            _txb(slide, T(self.tag), l=inner_l, t=y_off, w=inner_w, h=tag_h,
                 sz=9, bold=True, color=acc)
            y_off += tag_h + tag_gap

        # title: CN (+ EN subtitle in bilingual)
        cn_title = get_cn(self.title) if bilingual else T(self.title)
        en_title = get_en(self.title) if bilingual else ""
        has_en_title = bilingual and bool(en_title)

        if has_en_title:
            cn_title_h = max(Mm(8), int(title_budget * 0.60))
            en_title_h = title_budget - cn_title_h
        else:
            cn_title_h = title_budget
            en_title_h = 0

        if cn_title:
            tb = _txb(slide, cn_title, l=inner_l, t=y_off, w=inner_w,
                      h=cn_title_h, sz=15, bold=True, color=txt)
            try:
                tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
            y_off += cn_title_h

        if has_en_title and en_title_h > Mm(4):
            tb = _txb(slide, en_title, l=inner_l, t=y_off, w=inner_w,
                      h=en_title_h, sz=10, color=EN_TITLE_COLOR,
                      en_font="Inter", cn_font="Inter")
            try:
                tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
            y_off += en_title_h

        if bar_space:
            _rect(slide, l=inner_l, t=y_off, w=Mm(24), h=bar_h, fill=acc)
            y_off += bar_space

        # body: CN paragraph
        cn_color = CN_COLOR if bilingual else body_color_default
        if cn_body_text and cn_body_h > Mm(4):
            tb = _txb(slide, cn_body_text, l=inner_l, t=y_off, w=inner_w,
                      h=cn_body_h, sz=12, color=cn_color, ls_pt=17)
            try:
                tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
            y_off += cn_body_h

        # body: EN paragraph (bilingual only)
        if has_en_body and en_body_h > Mm(4):
            en_sz = max(8, int(12 * EN_SZ_RATIO))
            tb = _txb(slide, en_body_text, l=inner_l, t=y_off + Mm(1), w=inner_w,
                      h=en_body_h, sz=en_sz, color=EN_COLOR,
                      ls_pt=int(en_sz * 1.5),
                      en_font="Inter", cn_font="Inter")
            try:
                tb.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass

    def render_html(self) -> str:
        from scripts.i18n import T, is_bilingual, get_cn, get_en, EN_COLOR
        bg, acc, txt = self._colors()
        cls = f"card card--{self.color}"

        if self.layout == "horizontal_split":
            return (
                f'<div class="{cls}" style="background:{bg}">'
                f'<div class="card__left">'
                f'<div class="card__tag" style="color:{acc}">{T(self.tag)}</div>'
                f'<div class="card__title" style="color:{txt}">{T(self.title)}</div>'
                f'<div class="card__body">{T(self.body)}</div>'
                f'</div>'
                f'<div class="card__metric" style="color:{acc}">{T(self.metric)}</div>'
                f'</div>'
            )
        if self.layout == "icon_left":
            return (
                f'<div class="{cls}" style="background:{bg}">'
                f'<div class="card__icon-badge" style="background:{acc}">'
                f'{T(self.icon)[:2].upper()}</div>'
                f'<div class="card__content">'
                f'<div class="card__tag" style="color:{acc}">{T(self.tag)}</div>'
                f'<div class="card__title" style="color:{txt}">{T(self.title)}</div>'
                f'<div class="card__body">{T(self.body)}</div>'
                f'</div></div>'
            )
        if self.layout == "quote_card":
            return (
                f'<div class="{cls} card--quote" style="background:{bg}">'
                f'<div class="card__quote-mark" style="color:{acc}">"</div>'
                f'<div class="card__quote">{T(self.quote) or T(self.body)}</div>'
                f'<div class="card__attribution">— {T(self.attribution)}</div>'
                f'</div>'
            )
        # vertical_stack
        bar = f'<div class="card__bar" style="background:{acc}"></div>' if T(self.title) else ""
        en_body = get_en(self.body) if is_bilingual() else ""
        en_html = (
            f'<div class="card__body card__body--en" '
            f'style="color:{EN_COLOR};font-family:Inter,sans-serif;font-size:0.85em">'
            f'{en_body}</div>'
        ) if en_body else ""
        return (
            f'<div class="{cls}" style="background:{bg}">'
            f'<div class="card__tag" style="color:{acc}">{T(self.tag)}</div>'
            f'<div class="card__title" style="color:{txt}">{T(self.title)}</div>'
            f'{bar}'
            f'<div class="card__body">{T(self.body)}</div>'
            f'{en_html}'
            f'</div>'
        )
