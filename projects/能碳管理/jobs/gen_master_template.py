#!/usr/bin/env python3
"""
gen_master_template.py

未来方舟品牌 PPT 母版展示文件生成脚本。
输出: projects/能碳管理/docs/brand_ppt_master.pptx

包含 12 种版式的完整示范幻灯片，配以品牌装饰性手绘元素。
AI 可读取此文件理解版式结构，后续生成 PPT 时以此为参照。

版式目录（幻灯片顺序）:
  01  封面页       Cover
  02  章节分隔     Section
  03  标准内容     Title + Content
  04  两列内容     Two Columns
  05  三列卡片     Three Cards
  06  左文右图     Text + Image Right
  07  左图右文     Image Left + Text
  08  大数据指标   Big Stats
  09  水平时间轴   Timeline
  10  引用证言     Quote
  11  表格页       Table
  12  结尾致谢     Closing
"""

import sys, os
from PIL import Image as PILImage

_BRAND   = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..'))
_PROJECT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
sys.path.insert(0, _BRAND)

_DOCS   = os.path.join(_PROJECT, "docs")
_ASSETS = os.path.join(_BRAND, "assets")
_DECO   = os.path.join(_ASSETS, "装饰性元素")

os.makedirs(_DOCS, exist_ok=True)
os.makedirs(os.path.join(_DOCS, "img"), exist_ok=True)

from pptx import Presentation
from pptx.util import Inches, Pt, Mm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


# ─── 尺寸常量 ────────────────────────────────────────────────────────────────
W = Inches(13.33)   # 幻灯片宽  338.6mm
H = Inches(7.5)     # 幻灯片高  190.5mm

ML = Mm(21)         # 内容左边距
MR = Mm(21)         # 内容右边距
CW = W - ML - MR    # 内容宽度

HEADER_H   = Mm(36)         # 页眉区高（顶绿线+标题+副标题+分隔线）
FOOTER_H   = Mm(13)         # 页脚区高
CONTENT_Y  = HEADER_H       # 内容区起始 y
CONTENT_H  = H - CONTENT_Y - FOOTER_H  # 内容区可用高度

# 三列/两列卡片间距
C3_GAP = Mm(5)
C3_W   = (CW - 2 * C3_GAP) / 3
C2_GAP = Mm(6)
C2_W   = (CW - C2_GAP) / 2


# ─── 品牌色 ──────────────────────────────────────────────────────────────────
P5 = "#3EC99E"   # 主绿
P1 = "#EAFAF5"   # 主绿浅
S5 = "#C8E13C"   # 辅色黄绿
N9 = "#0E1216"   # 中性900（近黑）
N7 = "#3D444A"   # 中性700（正文）
N4 = "#8A9199"   # 中性400（说明）
N2 = "#D0D5DD"   # 中性200（边框）
N1 = "#F2F3F5"   # 中性100（背景）
WH = "#FFFFFF"   # 白

OK = "#53AF36"   # 成功绿
WA = "#F3B021"   # 警告橙
ER = "#F12D2D"   # 危险红
PU = "#8255E1"   # 补充紫

FONT_CN = "PingFang SC"
FONT_EN = "Inter"


# ─── 资源路径 ─────────────────────────────────────────────────────────────────
LOGO_H_P = os.path.join(_ASSETS, "logo-horizontal-primary.png")
LOGO_H_R = os.path.join(_ASSETS, "logo-horizontal-reverse.png")
LOGO_S_P = os.path.join(_ASSETS, "logo-stacked-primary.png")
LOGO_S_R = os.path.join(_ASSETS, "logo-stacked-reverse.png")

DECO = {
    "A_ul": os.path.join(_DECO, "A下划@4x.png"),    # 主绿下划线
    "A_ci": os.path.join(_DECO, "A画圈@4x.png"),    # 主绿画圈
    "A_a1": os.path.join(_DECO, "A箭头1@4x.png"),   # 主绿箭头→↑
    "A_a2": os.path.join(_DECO, "A箭头2@4x.png"),   # 主绿箭头→↓
    "B_ul": os.path.join(_DECO, "B下划@4x.png"),    # 深绿下划线
    "B_ci": os.path.join(_DECO, "B画圈@4x.png"),    # 深绿画圈
    "B_a1": os.path.join(_DECO, "B箭头1@4x.png"),   # 深绿箭头→↑
    "B_a2": os.path.join(_DECO, "B箭头2@4x.png"),   # 深绿箭头→↓
    "C_ul": os.path.join(_DECO, "C下划@4x.png"),    # 黄绿下划线
    "C_ci": os.path.join(_DECO, "C画圈@4x.png"),    # 黄绿画圈
    "C_a1": os.path.join(_DECO, "C箭头1@4x.png"),   # 黄绿箭头→↑
    "C_a2": os.path.join(_DECO, "C箭头2@4x.png"),   # 黄绿箭头→↓
}


# ─── 底层 Helper ──────────────────────────────────────────────────────────────

def _rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def rect(slide, l, t, w, h, fill=None, line=None, lw=0.75, rounded=False):
    sid = 5 if rounded else 1
    s = slide.shapes.add_shape(sid, int(l), int(t), int(w), int(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = _rgb(line)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def txb(slide, text, l, t, w, h,
        sz=16, bold=False, color=N7,
        align=PP_ALIGN.LEFT, wrap=True,
        ls_pt=None, name=""):
    """添加文本框，自动设置中英文双字体。"""
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if ls_pt:
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        etree.SubElement(lnSpc, qn("a:spcPts"),
                         attrib={"val": str(int(ls_pt * 100))})
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = FONT_EN
    rPr = run._r.get_or_add_rPr()
    for tag, face in [("a:latin", FONT_EN), ("a:ea", FONT_CN)]:
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", face)
    return tb


def txb_lines(slide, lines, l, t, w, h,
              sz=14, bold=False, color=N7,
              align=PP_ALIGN.LEFT, ls_pt=22):
    """
    添加多段落文本框。
    lines: list of str 或 list of (text, sz, bold, color)
    """
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, lsz, lb, lc = line, sz, bold, color
        else:
            txt, lsz, lb, lc = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if ls_pt:
            pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            etree.SubElement(lnSpc, qn("a:spcPts"),
                             attrib={"val": str(int(ls_pt * 100))})
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(lsz)
        run.font.bold = lb
        run.font.color.rgb = _rgb(lc)
        run.font.name = FONT_EN
        rPr = run._r.get_or_add_rPr()
        for tag, face in [("a:latin", FONT_EN), ("a:ea", FONT_CN)]:
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set("typeface", face)
    return tb


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(color)


def _img_aspect(path):
    img = PILImage.open(path)
    return img.width / img.height


def add_pic(slide, path, l, t, w, h=None):
    """添加图片，h=None 时按宽高比自动计算。"""
    if not (path and os.path.exists(path)):
        return None
    if h is None:
        r = _img_aspect(path)
        h = int(w / r)
    return slide.shapes.add_picture(path, int(l), int(t),
                                    width=int(w), height=int(h))


def logo_h(slide, reverse=False, r_mm=9, b_mm=3.5, h_mm=7):
    """水平 Logo，放右下角。"""
    p = LOGO_H_R if reverse else LOGO_H_P
    if not (p and os.path.exists(p)):
        return None
    lh = Mm(h_mm)
    lw = int(lh * _img_aspect(p))
    return slide.shapes.add_picture(
        p, int(W - Mm(r_mm) - lw), int(H - Mm(b_mm) - lh),
        width=lw, height=lh)


def logo_s(slide, reverse=True, l_mm=12, t_mm=10, h_mm=26):
    """堆叠 Logo，用于封面/章节页。"""
    p = LOGO_S_R if reverse else LOGO_S_P
    if not (p and os.path.exists(p)):
        return None
    lh = Mm(h_mm)
    lw = int(lh * _img_aspect(p))
    return slide.shapes.add_picture(
        p, Mm(l_mm), Mm(t_mm), width=lw, height=lh)


def deco(slide, key, l, t, w, h=None):
    """添加装饰元素。"""
    return add_pic(slide, DECO.get(key, ""), l, t, w, h)


def _new(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ─── 复合组件：页眉 / 页脚 ────────────────────────────────────────────────────

def header(slide, title, subtitle="", label=""):
    """
    通用页眉：顶部 2mm 绿线 + 标题 + 副标题 + 底部细分隔。
    返回内容区起始 y（CONTENT_Y）。
    """
    # 顶部绿线
    rect(slide, l=0, t=0, w=W, h=Mm(2.5), fill=P5)
    # 眉标（可选小分类标签）
    if label:
        txb(slide, label,
            l=ML, t=Mm(4.5),
            w=Mm(80), h=Mm(6),
            sz=9, color=P5, name="header_label")
    # 主标题
    txb(slide, title,
        l=ML, t=Mm(5.5 if not label else 11),
        w=CW * 0.82, h=Mm(18),
        sz=26, bold=True, color=N9, name="header_title")
    # 副标题
    if subtitle:
        txb(slide, subtitle,
            l=ML, t=Mm(24 if not label else 27),
            w=CW * 0.80, h=Mm(9),
            sz=12, color=N4, name="header_subtitle")
    # 分隔线
    rect(slide, l=ML, t=Mm(33.5), w=CW, h=Mm(0.4), fill=N2)
    return CONTENT_Y


def footer(slide, layout_label=""):
    """通用页脚：细线 + 水平 Logo + 版式标签（灰色右下角，AI 识别用）。"""
    rect(slide, l=0, t=H - Mm(13), w=W, h=Mm(0.4), fill=N2)
    logo_h(slide, reverse=False, r_mm=9, b_mm=3, h_mm=7)
    if layout_label:
        txb(slide, layout_label,
            l=W - Mm(72), t=H - Mm(12),
            w=Mm(60), h=Mm(10),
            sz=8, color=N2, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, bg=N1, accent=None, accent_side="top",
         accent_h_mm=1.2, accent_w_mm=3):
    """
    卡片底板：圆角矩形背景 + 可选顶部/左侧强调线。
    accent_side: "top" | "left"
    """
    r = rect(slide, l=l, t=t, w=w, h=h, fill=bg, rounded=False)
    if accent:
        if accent_side == "top":
            rect(slide, l=l, t=t,
                 w=w, h=Mm(accent_h_mm), fill=accent)
        elif accent_side == "left":
            rect(slide, l=l, t=t,
                 w=Mm(accent_w_mm), h=h, fill=accent)
    return r


# ─── 版式 01：封面 ────────────────────────────────────────────────────────────
def layout_01_cover(prs,
                    title="ArkSus® AI+ GEMP\n能碳一体化平台",
                    subtitle="为头部品牌方上游供应商打造的合规一体机",
                    eyebrow="能碳管理平台  ·  产品介绍",
                    meta="2026年06月"):
    """
    版式 01 Cover：
    - 深色背景 #0E1216
    - 左侧 4mm 主绿竖线
    - stacked-reverse Logo 左上角
    - 眉标 / 主标题（48pt Bold 白）/ 副标题 / 元信息
    - A系列下划线装饰于标题下
    - A箭头2 右下角点缀
    - 底部 3.5mm 主绿横线

    AI 使用说明：
      · title      主标题，最多两行，约 14 字/行
      · subtitle   副标题，一句话描述
      · eyebrow    左上标签（项目名/部门）
      · meta       右下元信息（日期/版本）
    """
    slide = _new(prs)
    set_bg(slide, N9)

    # 左侧竖线装饰条
    rect(slide, l=0, t=Mm(25), w=Mm(4), h=Mm(105), fill=P5)

    # Logo stacked-reverse 左上
    logo_s(slide, reverse=True, l_mm=12, t_mm=10, h_mm=24)

    # 眉标
    txb(slide, eyebrow,
        l=Mm(18), t=Mm(50),
        w=Mm(220), h=Mm(9),
        sz=11, color=N4, name="cover_eyebrow")

    # 水平分隔线（标题上方）
    rect(slide, l=Mm(18), t=Mm(60), w=W - Mm(36), h=Mm(0.5), fill=N7)

    # 主标题
    tb = slide.shapes.add_textbox(int(Mm(18)), int(Mm(63)),
                                  int(Mm(225)), int(Mm(52)))
    tb.name = "cover_title"
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(title.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        etree.SubElement(lnSpc, qn("a:spcPts"), attrib={"val": "5400"})
        run = p.add_run()
        run.text = line
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = _rgb(WH)
        run.font.name = FONT_EN
        rPr = run._r.get_or_add_rPr()
        for tag, face in [("a:latin", FONT_EN), ("a:ea", FONT_CN)]:
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set("typeface", face)

    # A系列下划线装饰（标题第一行下）
    deco(slide, "A_ul", l=Mm(18), t=Mm(107), w=Mm(130))

    # 副标题
    txb(slide, subtitle,
        l=Mm(18), t=Mm(114),
        w=Mm(220), h=Mm(12),
        sz=14, color=N4, ls_pt=21, name="cover_subtitle")

    # 元信息
    txb(slide, meta,
        l=Mm(18), t=Mm(130),
        w=Mm(100), h=Mm(9),
        sz=11, color=N4, name="cover_meta")

    # 装饰：右下角 A箭头2（向右下）
    deco(slide, "A_a2",
         l=W - Mm(62), t=H - Mm(60), w=Mm(52))

    # 底部绿横线
    rect(slide, l=0, t=H - Mm(3.5), w=W, h=Mm(3.5), fill=P5)

    # 版式标识（AI读取）
    txb(slide, "LAYOUT: 01_Cover",
        l=W - Mm(70), t=H - Mm(10),
        w=Mm(60), h=Mm(8),
        sz=7, color=N7, align=PP_ALIGN.RIGHT)
    return slide


# ─── 版式 02：章节分隔 ───────────────────────────────────────────────────────
def layout_02_section(prs,
                      num="01",
                      title="平台核心能力",
                      subtitle="Core Capabilities"):
    """
    版式 02 Section：
    - 主绿背景 #3EC99E
    - 章节编号（大号，半透明白）
    - 章节标题（白色 36pt Bold）
    - 英文副标题（白色 16pt）
    - A画圈装饰环绕章节编号
    - stacked-reverse Logo 右下角

    AI 使用说明：
      · num        章节编号，如 "01" "02"
      · title      章节标题
      · subtitle   英文副标题（可选）
    """
    slide = _new(prs)
    set_bg(slide, P5)

    # 大号半透明章节编号（背景装饰层）
    txb(slide, num,
        l=W * 0.52, t=Mm(38),
        w=Mm(110), h=Mm(80),
        sz=120, bold=True, color=WH,
        align=PP_ALIGN.LEFT, name="section_bg_num")
    # 让大号数字半透明（通过设置颜色近似处理）
    # 使用 XML 设置透明度
    _set_text_alpha(slide.shapes[-1], alpha_pct=15)

    # A画圈装饰（围绕章节编号）
    deco(slide, "A_ci",
         l=W * 0.50, t=Mm(28), w=Mm(85))

    # 左侧内容区
    # 章节编号（前景）
    txb(slide, num,
        l=ML, t=Mm(62),
        w=Mm(50), h=Mm(25),
        sz=32, bold=True, color=WH, name="section_num")

    # 章节标题
    txb(slide, title,
        l=ML, t=Mm(88),
        w=Mm(240), h=Mm(35),
        sz=38, bold=True, color=WH,
        align=PP_ALIGN.LEFT, ls_pt=48, name="section_title")

    # 英文副标题
    if subtitle:
        txb(slide, subtitle,
            l=ML, t=Mm(124),
            w=Mm(200), h=Mm(12),
            sz=16, color=WH, name="section_subtitle")

    # 水平细线（标题上方）
    rect(slide, l=ML, t=Mm(85), w=Mm(60), h=Mm(1), fill=WH)

    # Logo stacked-reverse 右下
    _place_logo_s_reverse_bottom_right(slide, h_mm=20, r_mm=14, b_mm=12)

    # 版式标识
    txb(slide, "LAYOUT: 02_Section",
        l=W - Mm(70), t=H - Mm(10),
        w=Mm(60), h=Mm(8),
        sz=7, color=WH, align=PP_ALIGN.RIGHT)
    return slide


def _set_text_alpha(shape, alpha_pct=20):
    """设置文本框文字透明度（通过 XML）。"""
    try:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                rPr = run._r.get_or_add_rPr()
                solidFill = rPr.find(qn("a:solidFill"))
                if solidFill is None:
                    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
                srgb = solidFill.find(qn("a:srgbClr"))
                if srgb is None:
                    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
                srgb.set("val", "FFFFFF")
                alpha_el = srgb.find(qn("a:alpha"))
                if alpha_el is None:
                    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
                alpha_el.set("val", str(alpha_pct * 1000))
    except Exception:
        pass


def _place_logo_s_reverse_bottom_right(slide, h_mm=20, r_mm=14, b_mm=12):
    """在右下角放置 stacked-reverse Logo。"""
    p = LOGO_S_R
    if not (p and os.path.exists(p)):
        return
    lh = Mm(h_mm)
    lw = int(lh * _img_aspect(p))
    slide.shapes.add_picture(
        p,
        int(W - Mm(r_mm) - lw),
        int(H - Mm(b_mm) - lh),
        width=lw, height=lh)


# ─── 版式 03：标准内容页 ──────────────────────────────────────────────────────
def layout_03_title_content(prs,
                             title="平台功能架构",
                             subtitle="一横四纵 · 六大行业全链覆盖",
                             label="产品架构"):
    """
    版式 03 Title + Content（单列内容）：
    - 白色背景
    - 通用页眉（顶绿线+标题+副标题+分隔线）
    - 宽内容区，适合纯文字、大图、图文混排
    - 通用页脚

    AI 使用说明：
      · title    主标题（26pt Bold）
      · subtitle 副标题（12pt 灰）
      · label    眉标标签（9pt 绿）
      · content  内容区自由放置
    """
    slide = _new(prs)
    set_bg(slide, WH)
    header(slide, title, subtitle, label)

    # 示例内容：架构层次说明（竖排）
    # 顶部：三个横向功能模块卡片
    card_h = Mm(28)
    cards = [
        (P5,  "数据采集层", "IoT / ERP / MES 多源接入，支持 OPC-UA、MQTT、API"),
        (P5,  "核算计算层", "Scope 1/2/3 自动核算，ISO 14064 & GHG Protocol"),
        (S5,  "报告输出层", "报告中心一键生成，Limited Assurance 自动归档"),
    ]
    for i, (ac, ttl, desc) in enumerate(cards):
        cx = ML + i * (C3_W + C3_GAP)
        cy = CONTENT_Y + Mm(3)
        card(slide, l=cx, t=cy, w=C3_W, h=card_h, bg=N1, accent=ac, accent_side="top")
        txb(slide, ttl, l=cx + Mm(4), t=cy + Mm(4),
            w=C3_W - Mm(8), h=Mm(10), sz=14, bold=True, color=N9)
        txb(slide, desc, l=cx + Mm(4), t=cy + Mm(14),
            w=C3_W - Mm(8), h=Mm(14), sz=10, color=N7, ls_pt=15)

    # 中间：横向流程箭头示意
    flow_y = CONTENT_Y + Mm(36)
    rect(slide, l=ML, t=flow_y + Mm(6), w=CW, h=Mm(1.5), fill=N2)
    steps = ["排放源识别", "数据采集", "自动核算", "内部审核", "核查报告"]
    sw = CW / len(steps)
    for i, s in enumerate(steps):
        cx = ML + i * sw
        rect(slide, l=cx + Mm(2), t=flow_y + Mm(2), w=sw - Mm(4), h=Mm(11),
             fill=P1, rounded=True)
        txb(slide, s, l=cx + Mm(3), t=flow_y + Mm(3),
            w=sw - Mm(6), h=Mm(8), sz=11, bold=True, color=P5,
            align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            txb(slide, "→", l=cx + sw - Mm(4), t=flow_y + Mm(3),
                w=Mm(8), h=Mm(8), sz=11, color=P5, align=PP_ALIGN.CENTER)

    # 下方：大图占位区域
    rect(slide, l=ML, t=flow_y + Mm(18),
         w=CW, h=CONTENT_H - Mm(36) - Mm(18),
         fill=N1, line=N2, lw=0.5)
    txb(slide, "[ 图片 / 截图区域 · 建议比例 16:5 ]",
        l=ML + Mm(10), t=flow_y + Mm(30),
        w=CW - Mm(20), h=Mm(14),
        sz=12, color=N4, align=PP_ALIGN.CENTER)
    # B箭头2 指向图片区
    deco(slide, "B_a2",
         l=ML + CW * 0.72, t=flow_y + Mm(14), w=Mm(22))

    footer(slide, "LAYOUT: 03_Title+Content")
    return slide


# ─── 版式 04：两列内容 ───────────────────────────────────────────────────────
def layout_04_two_col(prs,
                      title="合规难点 · 与平台解法",
                      subtitle="每一条都对应一个工厂真实损失敞口",
                      label="痛点解法"):
    """
    版式 04 Two Columns：
    - 左右两列，各有子标题 + 内容
    - 中间 0.5mm 分隔线
    - 可用于：对比/方案/痛点-解法/before-after

    AI 使用说明：
      · left_title / left_items    左列子标题与内容列表
      · right_title / right_items  右列子标题与内容列表
      · 每列宽约 145mm，内容高约 130mm
    """
    slide = _new(prs)
    set_bg(slide, WH)
    header(slide, title, subtitle, label)

    cy = CONTENT_Y + Mm(4)
    ch = CONTENT_H - Mm(4)
    gx = ML + C2_W + C2_GAP / 2  # 中线 x

    # 中间分隔线
    rect(slide, l=gx - Mm(0.25), t=cy, w=Mm(0.5), h=ch, fill=N2)

    left_title = "挑战"
    right_title = "解法"
    left_items = [
        ("01", "SCAR 触发", "品牌方发现 Scope 3 缺漏，单点触发整线停产风险，当年订单直接清零。"),
        ("02", "BOM 数据错配", "BOM 因子库 vintage 不一致，PCF 报告每季度返工，核查机构要求重复举证。"),
        ("03", "能源管理手动", "峰谷用电靠人工调度，PV 弃光 18%，节费机会日复一日流失。"),
    ]
    right_items = [
        ("01", "Scope 3 自动盯防", "平台 80% 实质性排放源自动识别，距截止日倒计时提醒，一次提交三件套。"),
        ("02", "因子 vintage 留痕", "Tier-2 PCF 追溯 CRM，ISO 14040/14044/14067 注脚自动生成。"),
        ("03", "AI 多能源调度", "24h 联合寻优峰谷+PV+储能+工序，节费按月入账，自动派工。"),
    ]

    for col_idx, (ct, items) in enumerate([(left_title, left_items),
                                            (right_title, right_items)]):
        cx = ML if col_idx == 0 else ML + C2_W + C2_GAP
        # 列标题
        txb(slide, ct, l=cx, t=cy,
            w=C2_W, h=Mm(11), sz=13, bold=True, color=P5)
        rect(slide, l=cx, t=cy + Mm(11), w=Mm(24), h=Mm(1), fill=P5)

        item_y = cy + Mm(16)
        for num_icon, item_title, item_desc in items:
            # 序号 badge
            b = rect(slide, l=cx, t=item_y, w=Mm(9), h=Mm(9), fill=P1)
            txb(slide, "·", l=cx + Mm(2), t=item_y + Mm(1),
                w=Mm(6), h=Mm(7), sz=14, bold=True, color=P5,
                align=PP_ALIGN.CENTER)
            # 项目标题 + 描述
            txb(slide, item_title, l=cx + Mm(12), t=item_y,
                w=C2_W - Mm(13), h=Mm(9), sz=13, bold=True, color=N9)
            txb(slide, item_desc, l=cx + Mm(12), t=item_y + Mm(10),
                w=C2_W - Mm(13), h=Mm(18), sz=11, color=N7, ls_pt=17)
            item_y += Mm(32)

    footer(slide, "LAYOUT: 04_Two-Columns")
    return slide


# ─── 版式 05：三列卡片 ───────────────────────────────────────────────────────
def layout_05_three_cards(prs,
                           title="四大可量化回报",
                           subtitle="Payback < 6 个月 · 工厂首年数据",
                           label="价值主张"):
    """
    版式 05 Three Cards：
    - 三列等宽卡片，每列独立卡片背景
    - 卡片顶部可选强调色线（主绿/辅色/成功绿等）
    - 适合：功能对比/价值卡片/产品模块/团队/步骤

    AI 使用说明：
      · 每列卡片宽约 96mm，高约 130mm
      · 卡片顶部强调色线：cards[i] 中 accent 字段指定颜色
      · 每卡片含：编号 badge + 标题 + 正文（最多 5 行）
    """
    slide = _new(prs)
    set_bg(slide, WH)
    header(slide, title, subtitle, label)

    cards_data = [
        (P5,  "01", "保住大品牌订单",
         "4 步向导确保 PER/清洁能源/限用物质一次过审",
         ["FY 审计一次通过", "SCAR 归零", "AVL 名录稳定"]),
        (S5,  "02", "年省电费 + 减排",
         "AI 多能源调度 + 能效诊断，峰段电费、PV 弃光按回利润表",
         ["峰段占比 47%→32%", "PV 弃光 18%→4%", "年节费 ¥ 420 万+"]),
        (OK,  "03", "绿色工厂溢价",
         "GB/T 36132 五星评级，年度议价直接折算单价 3–5% 上浮",
         ["三星→四星升级", "长协优先资格", "绿色信贷利率下浮"]),
    ]

    cy = CONTENT_Y + Mm(4)
    card_h = CONTENT_H - Mm(6)

    for i, (ac, num, ttl, desc, bullets) in enumerate(cards_data):
        cx = ML + i * (C3_W + C3_GAP)
        card(slide, l=cx, t=cy, w=C3_W, h=card_h, bg=N1,
             accent=ac, accent_side="top", accent_h_mm=2.5)
        # 序号
        txb(slide, num, l=cx + Mm(4), t=cy + Mm(6),
            w=Mm(22), h=Mm(14), sz=22, bold=True, color=_hex_lighten(ac))
        # 标题
        txb(slide, ttl, l=cx + Mm(4), t=cy + Mm(22),
            w=C3_W - Mm(8), h=Mm(14), sz=15, bold=True, color=N9)
        # 描述
        txb(slide, desc, l=cx + Mm(4), t=cy + Mm(38),
            w=C3_W - Mm(8), h=Mm(22), sz=11, color=N7, ls_pt=17)
        # 细分隔线
        rect(slide, l=cx + Mm(4), t=cy + Mm(62), w=C3_W - Mm(8), h=Mm(0.4), fill=N2)
        # Bullets
        bul_y = cy + Mm(65)
        for b in bullets:
            txb(slide, f"  ✓  {b}", l=cx + Mm(4), t=bul_y,
                w=C3_W - Mm(8), h=Mm(10), sz=11, color=N7)
            bul_y += Mm(11)

    # C系列画圈装饰在第三列标题附近
    cx3 = ML + 2 * (C3_W + C3_GAP)
    deco(slide, "C_ci", l=cx3 - Mm(8), t=cy + Mm(14), w=C3_W + Mm(16))

    footer(slide, "LAYOUT: 05_Three-Cards")
    return slide


def _hex_lighten(h, alpha=0.4):
    """将颜色按 alpha 混合白色，返回十六进制（近似淡化效果）。"""
    h = h.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r2 = int(r * alpha + 255 * (1 - alpha))
    g2 = int(g * alpha + 255 * (1 - alpha))
    b2 = int(b * alpha + 255 * (1 - alpha))
    return "#{:02X}{:02X}{:02X}".format(r2, g2, b2)


# ─── 版式 06：左文右图 ───────────────────────────────────────────────────────
def layout_06_text_image_right(prs,
                                title="能碳驾驶舱",
                                subtitle="一屏掌握全厂碳能双指标",
                                label="产品截图"):
    """
    版式 06 Text + Image Right：
    - 左侧文字区（约 54% 宽）
    - 右侧图片区（约 42% 宽，全高，带浅色背景占位）
    - 适合：功能说明 + 截图/示意图

    AI 使用说明：
      · 左侧文字区: l=ML, t=CONTENT_Y, w=约 182mm, h=CONTENT_H
      · 右侧图片区: l=约 220mm, t=0, w=约 120mm, h=全高
      · 图片路径替换右侧占位矩形中的 add_picture 调用
    """
    slide = _new(prs)
    set_bg(slide, WH)
    header(slide, title, subtitle, label)

    img_w = Mm(122)        # 右侧图片区宽
    txt_w = W - img_w - ML - Mm(8)  # 左侧文字区宽
    img_l = W - img_w

    # 右侧图片背景区（占位）
    rect(slide, l=img_l, t=0, w=img_w, h=H, fill=N1)
    txb(slide, "[ 截图 / 示意图 ]",
        l=img_l + Mm(10), t=H / 2 - Mm(8),
        w=img_w - Mm(20), h=Mm(14),
        sz=12, color=N4, align=PP_ALIGN.CENTER)

    # 左侧内容区
    cy = CONTENT_Y + Mm(4)

    # 功能亮点列表（3条带竖线的条目）
    items = [
        ("Scope 1/2/3 实时核算",
         "排放源自动识别，核算结果每日更新，距目标差距一目了然。"),
        ("能耗异常 AI 预警",
         "设备级 3D 热点地图，一屏定位告警/预警/能效偏差，响应时间 < 30 秒。"),
        ("合规进度仪表盘",
         "品牌方 FY 截止日倒计时，PER/清洁能源/限用物质���项合规状态全显示。"),
    ]
    for j, (item_t, item_d) in enumerate(items):
        iy = cy + j * Mm(38)
        # 左侧绿竖线
        rect(slide, l=ML, t=iy, w=Mm(3), h=Mm(28), fill=P5)
        txb(slide, item_t, l=ML + Mm(7), t=iy,
            w=txt_w - Mm(8), h=Mm(12), sz=14, bold=True, color=N9)
        txb(slide, item_d, l=ML + Mm(7), t=iy + Mm(13),
            w=txt_w - Mm(8), h=Mm(16), sz=11, color=N7, ls_pt=17)

    # 底部 CTA 区块
    cta_y = cy + Mm(118)
    rect(slide, l=ML, t=cta_y, w=txt_w, h=Mm(14), fill=P5, rounded=True)
    txb(slide, "→  预约功能演示",
        l=ML + Mm(6), t=cta_y + Mm(3),
        w=txt_w - Mm(8), h=Mm(9),
        sz=12, bold=True, color=WH)

    # B箭头1 指向 CTA
    deco(slide, "B_a1", l=ML + txt_w - Mm(30), t=cta_y - Mm(22), w=Mm(28))

    footer(slide, "LAYOUT: 06_Text+Image-Right")
    return slide


# ─── 版式 07：左图右文 ───────────────────────────────────────────────────────
def layout_07_image_left_text(prs,
                               title="绿色工厂评价模块",
                               subtitle="GB/T 36132-2025 五星评级自动辅助",
                               label="功能详解"):
    """
    版式 07 Image Left + Text：
    - 左侧图片区（约 42% 宽，从顶部延伸到底）
    - 右侧文字区
    - 版式 06 的镜像，适合交替使用保持视觉节奏

    AI 使用说明：
      · 左侧图片区: l=0, t=0, w=约 120mm, h=全高
      · 右侧文字区: l=约 128mm, t=CONTENT_Y, w=约 185mm, h=CONTENT_H
    """
    slide = _new(prs)
    set_bg(slide, WH)

    img_w = Mm(120)
    txt_l = img_w + Mm(8)
    txt_w = W - txt_l - MR

    # 左侧图片区（深色背景占位，保留品牌感）
    rect(slide, l=0, t=0, w=img_w, h=H, fill=N9)
    txb(slide, "[ 截图 / 示意图 ]",
        l=Mm(10), t=H / 2 - Mm(8),
        w=img_w - Mm(20), h=Mm(14),
        sz=12, color=N4, align=PP_ALIGN.CENTER)
    # Logo stacked-reverse 在图片区左下
    _place_logo_s_reverse_bottom_right_in_area(slide, img_w, h_mm=18, r_mm=10, b_mm=10)

    # 右侧页眉
    rect(slide, l=txt_l, t=0, w=W - txt_l, h=Mm(2.5), fill=P5)
    txb(slide, "功能详解",
        l=txt_l, t=Mm(4.5), w=Mm(60), h=Mm(6), sz=9, color=P5)
    txb(slide, title,
        l=txt_l, t=Mm(11), w=txt_w, h=Mm(16),
        sz=24, bold=True, color=N9)
    if subtitle:
        txb(slide, subtitle,
            l=txt_l, t=Mm(27), w=txt_w, h=Mm(8), sz=11, color=N4)
    rect(slide, l=txt_l, t=Mm(36), w=txt_w, h=Mm(0.4), fill=N2)

    # 右侧内容：评价维度列表
    dims = [
        (P5,   "能源消耗综合指标", "单位产值综合能耗 ≤ 基准值 80%，可再生能源占比 ≥ 20%"),
        (S5,   "污染物达标排放",   "工业废水/废气/固废 100% 达标，VOCs 在线监测数据完整"),
        (OK,   "温室气体排放",     "Scope 1+2 排放强度年降幅 ≥ 3%，CCER 抵消可选"),
        ("#3CC5CF", "用水效率",         "工业用水重复利用率 ≥ 90%，废水回用率持续提升"),
    ]
    dim_y = CONTENT_Y + Mm(5)
    for ac, dt, dd in dims:
        rect(slide, l=txt_l, t=dim_y, w=Mm(3.5), h=Mm(22), fill=ac)
        txb(slide, dt, l=txt_l + Mm(7), t=dim_y,
            w=txt_w - Mm(8), h=Mm(11), sz=13, bold=True, color=N9)
        txb(slide, dd, l=txt_l + Mm(7), t=dim_y + Mm(12),
            w=txt_w - Mm(8), h=Mm(11), sz=10, color=N7, ls_pt=15)
        dim_y += Mm(27)

    footer(slide, "LAYOUT: 07_Image-Left+Text")
    return slide


def _place_logo_s_reverse_bottom_right_in_area(slide, area_w, h_mm=18, r_mm=8, b_mm=8):
    p = LOGO_S_R
    if not (p and os.path.exists(p)):
        return
    lh = Mm(h_mm)
    lw = int(lh * _img_aspect(p))
    slide.shapes.add_picture(
        p, int(area_w - Mm(r_mm) - lw), int(H - Mm(b_mm) - lh),
        width=lw, height=lh)


# ─── 版式 08：大数据指标 ─────────────────────────────────────────────────────
def layout_08_big_stats(prs,
                         title="90 天硬指标成效",
                         subtitle="某品牌 Tier-1 电池模组厂 · 华东 · 年产值 ¥18 亿",
                         label="客户成效"):
    """
    版式 08 Big Stats（2×2 指标卡）：
    - 四个大数字指标卡，2行×2列布局
    - 每卡含：大数字（48pt）+ 单位 + 指标名 + 变化趋势
    - 适合：ROI 数据/指标成效/量化回报/关键参数

    AI 使用说明：
      · stats 列表：(before, after, unit, label, trend_color)
      · 2×2 网格，每卡宽约 (CW-gap)/2，高约 CONTENT_H/2
    """
    slide = _new(prs)
    set_bg(slide, WH)
    header(slide, title, subtitle, label)

    stats = [
        ("6 周",   "5 天",  "",    "核查周期",     P5,  "↓ 83%"),
        ("4 次/年", "0 次",  "",    "PCF 重做次数", OK,  "↓ 100%"),
        ("47%",    "32%",   "",    "峰段电费占比", WA,  "↓ 15pp"),
        ("¥",      "420万", "+",   "年节约电费",   S5,  "节费净增"),
    ]

    cy = CONTENT_Y + Mm(4)
    cg = Mm(6)
    sw = (CW - cg) / 2
    sh = (CONTENT_H - Mm(6) - cg) / 2

    for i, (before, after, unit, lbl, ac, trend) in enumerate(stats):
        row, col = divmod(i, 2)
        cx = ML + col * (sw + cg)
        sy = cy + row * (sh + cg)

        card(slide, l=cx, t=sy, w=sw, h=sh, bg=N1,
             accent=ac, accent_side="top", accent_h_mm=2.5)

        # Before（灰色删除线效果用颜色表达）
        txb(slide, f"原：{before}", l=cx + Mm(4), t=sy + Mm(6),
            w=sw - Mm(8), h=Mm(8), sz=11, color=N4)

        # After（大数字）
        txb(slide, after + unit,
            l=cx + Mm(4), t=sy + Mm(14),
            w=sw - Mm(8), h=Mm(26),
            sz=40, bold=True, color=ac, ls_pt=40)

        # 指标名
        txb(slide, lbl, l=cx + Mm(4), t=sy + sh - Mm(24),
            w=sw * 0.65, h=Mm(12), sz=13, bold=True, color=N9)

        # 趋势标签
        rect(slide, l=cx + sw - Mm(28), t=sy + sh - Mm(22),
             w=Mm(24), h=Mm(10), fill=_hex_lighten(ac, 0.3), rounded=True)
        txb(slide, trend,
            l=cx + sw - Mm(28), t=sy + sh - Mm(21),
            w=Mm(24), h=Mm(9), sz=10, bold=True, color=ac,
            align=PP_ALIGN.CENTER)

    # C画圈装饰在右下指标区
    deco(slide, "C_ci",
         l=ML + CW * 0.52, t=cy + sh + cg + Mm(2),
         w=sw + Mm(10))

    footer(slide, "LAYOUT: 08_Big-Stats")
    return slide


# ─── 版式 09：水平时间轴 ─────────────────────────────────────────────────────
def layout_09_timeline(prs,
                        title="合规路线图",
                        subtitle="欧盟四大核心法规 · 关键时间节点",
                        label="政策时间轴"):
    """
    版式 09 Timeline（水平时间轴）：
    - 水平主绿时间轴线
    - 节点：实心圆点 + 时间标签（上/下交替）+ 事件标题 + 描述
    - 适合：政策时间线/产品路线图/项目里程碑/流程步骤

    AI 使用说明：
      · events 列表：(year, title, desc, above)
      · above=True 时标签在轴线上方，False 在下方（交替排列）
      · 轴线 y = CONTENT_Y + CONTENT_H/2，上下各约 55mm 内容区
    """
    slide = _new(prs)
    set_bg(slide, WH)
    header(slide, title, subtitle, label)

    # 时间轴参数
    line_y = CONTENT_Y + CONTENT_H * 0.42
    events = [
        ("2023.10", "CBAM 过渡期", "开始强制季度申报，\n默认值可用 8%",         True),
        ("2025.02", "电池碳足迹声明", "EV 电池须强制\n碳足迹声明",              False),
        ("2026.01", "CBAM 正式征税", "按 EU ETS 对标\n价格缴纳证书",            True),
        ("2027.02", "电池 DBP",     "数字电池护照\n（DBP）强制实施",            False),
        ("2028+",   "CSRD 扩展",   "中国母公司\n纳入 CSRD 范围",               True),
    ]

    n = len(events)
    seg_w = CW / (n + 1)

    # 主轴线
    rect(slide, l=ML, t=line_y - Mm(1), w=CW, h=Mm(2), fill=P5)

    # 轴线起止箭头
    txb(slide, "▶", l=ML + CW - Mm(3), t=line_y - Mm(4),
        w=Mm(10), h=Mm(8), sz=10, color=P5)

    for i, (yr, ttl, desc, above) in enumerate(events):
        ex = ML + (i + 1) * seg_w

        # 节点圆点（用矩形模拟）
        dot_r = Mm(4)
        rect(slide, l=ex - dot_r, t=line_y - dot_r,
             w=dot_r * 2, h=dot_r * 2,
             fill=P5 if above else S5, rounded=True)
        # 内圈白点
        rect(slide, l=ex - Mm(1.5), t=line_y - Mm(1.5),
             w=Mm(3), h=Mm(3), fill=WH, rounded=True)

        # 竖连接线
        if above:
            rect(slide, l=ex - Mm(0.5), t=line_y - Mm(28),
                 w=Mm(1), h=Mm(24), fill=N2)
        else:
            rect(slide, l=ex - Mm(0.5), t=line_y + Mm(4),
                 w=Mm(1), h=Mm(24), fill=N2)

        # 时间标签
        ty = (line_y - Mm(40)) if above else (line_y + Mm(30))
        txb(slide, yr, l=ex - Mm(20), t=ty,
            w=Mm(40), h=Mm(8), sz=10, bold=True,
            color=P5 if above else S5,
            align=PP_ALIGN.CENTER)

        # 事件标题
        ty2 = (line_y - Mm(30)) if above else (line_y + Mm(40))
        txb(slide, ttl, l=ex - Mm(22), t=ty2,
            w=Mm(44), h=Mm(12), sz=12, bold=True, color=N9,
            align=PP_ALIGN.CENTER)

        # 事件描述
        ty3 = (line_y - Mm(18)) if above else (line_y + Mm(53))
        txb(slide, desc, l=ex - Mm(22), t=ty3,
            w=Mm(44), h=Mm(20), sz=10, color=N7,
            align=PP_ALIGN.CENTER, ls_pt=14)

    # 装饰：A箭头1 在时间轴右侧
    deco(slide, "A_a1", l=W - Mm(60), t=line_y - Mm(50), w=Mm(42))

    footer(slide, "LAYOUT: 09_Timeline")
    return slide


# ─── 版式 10：引用证言 ───────────────────────────────────────────────────────
def layout_10_quote(prs,
                    quote="平台把以前 6 周的核查工作压到 5 天，我们终于能把人留在车间而不是 Excel 里。",
                    author="—— 可持续总监，某品牌 Tier-1 电池模组厂",
                    context="华东 · 年产值 ¥18 亿 · 90 天落地"):
    """
    版式 10 Quote（引用/证言）：
    - 浅绿色背景或白色背景
    - 大引号装饰（主绿，超大字号）
    - 引用文字居中（22pt）
    - 来源作者右对齐（14pt 灰）
    - 适合：客户证言/领导寄语/数据声明/关键结论

    AI 使用说明：
      · quote   引用正文，建议 30–80 字
      · author  来源（姓名/职位/公司）
      · context 背景信息（可选）
    """
    slide = _new(prs)
    set_bg(slide, P1)

    # 背景装饰：左侧深绿渐变条
    rect(slide, l=0, t=0, w=Mm(8), h=H, fill=P5)

    # 大引号（背景层，主绿半透明）
    txb(slide, "“",   # " 左双引号
        l=Mm(28), t=Mm(22),
        w=Mm(80), h=Mm(80),
        sz=160, bold=True, color=P5,
        align=PP_ALIGN.LEFT)
    _set_text_alpha(slide.shapes[-1], alpha_pct=18)

    # 引用文字（前景）
    txb(slide, f"“{quote}”",
        l=Mm(36), t=H * 0.28,
        w=W - Mm(56), h=Mm(70),
        sz=20, bold=False, color=N9,
        align=PP_ALIGN.LEFT, ls_pt=32, name="quote_text")

    # A下划线装饰（引用文字下）
    deco(slide, "A_ul",
         l=Mm(36), t=H * 0.28 + Mm(74),
         w=Mm(140))

    # 作者
    txb(slide, author,
        l=Mm(36), t=H * 0.28 + Mm(82),
        w=W - Mm(56), h=Mm(12),
        sz=13, bold=True, color=N9,
        align=PP_ALIGN.LEFT, name="quote_author")

    # 背景
    if context:
        txb(slide, context,
            l=Mm(36), t=H * 0.28 + Mm(95),
            w=W - Mm(56), h=Mm(9),
            sz=11, color=N4,
            align=PP_ALIGN.LEFT)

    # Logo 水平 右下
    logo_h(slide, reverse=False, r_mm=9, b_mm=4, h_mm=7)

    # 版式标识
    txb(slide, "LAYOUT: 10_Quote",
        l=W - Mm(70), t=H - Mm(12),
        w=Mm(60), h=Mm(8),
        sz=7, color=N4, align=PP_ALIGN.RIGHT)
    return slide


# ─── 版式 11：表格页 ─────────────────────────────────────────────────────────
def layout_11_table(prs,
                    title="欧盟法规对比",
                    subtitle="四部核心法规要求对照",
                    label="合规矩阵"):
    """
    版式 11 Table：
    - 通用页眉
    - Header 行：主绿背景 + 白色文字
    - 奇数数据行：中性100背景
    - 偶数数据行：白色背景
    - 最右列可放置状态标签

    AI 使用说明：
      · 表格整体: l=ML, t=CONTENT_Y+Mm(4), w=CW
      · Header 行高 ~10mm，数据行高 ~12mm
      · 列宽按内容比例分配
    """
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn as _qn

    slide = _new(prs)
    set_bg(slide, WH)
    header(slide, title, subtitle, label)

    # 表格数据
    col_widths_pct = [0.14, 0.20, 0.22, 0.22, 0.22]
    col_headers = ["法规", "适用范围", "强制要求", "时间节点", "平台支持"]
    rows = [
        ["CBAM", "钢铁/铝/水泥/化肥/电力/氢",
         "季度提交嵌入排放报告，按 EU ETS 缴费",
         "2026.01 正式", "✓ 自动核算"],
        ["CSRD", "在欧营业额 >€150M 全球企业",
         "双重重要性分析 + Scope 1/2/3 有限保证",
         "2024 大型企业", "✓ 报告中心"],
        ["ESPR", "几乎所有在欧销售产品",
         "强制 DPP：材料/碳足迹/维修评分",
         "2027 电池率先", "✓ DPP 模块"],
        ["电池法规", "EV 电池/工业电池 >2kWh",
         "强制 DBP + Co/Li/Ni/Pb 再生含量披露",
         "2027.08 DBP", "✓ 碳护照"],
    ]

    # 用矩形+文本框模拟表格（python-pptx 原生 table 兼容性更好，但这里用手绘保持灵活性）
    # 改用 python-pptx 原生 table
    ty = int(CONTENT_Y + Mm(4))
    row_h = [Emu(int(Mm(10)))] + [Emu(int(Mm(11.5)))] * len(rows)
    col_w = [int(CW * p) for p in col_widths_pct]

    tbl = slide.shapes.add_table(
        len(rows) + 1, len(col_headers),
        int(ML), ty,
        int(CW), int(Mm(10) + Mm(11.5) * len(rows))
    ).table

    # 列宽
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = cw

    # 行高
    for ri, rh in enumerate(row_h):
        tbl.rows[ri].height = rh

    # Header 行
    for ci, ch in enumerate(col_headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(P5)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = ch
        run.font.size = _Pt(11)
        run.font.bold = True
        run.font.color.rgb = _rgb(WH)
        run.font.name = FONT_EN
        rPr = run._r.get_or_add_rPr()
        for tag, face in [("a:latin", FONT_EN), ("a:ea", FONT_CN)]:
            el = rPr.find(_qn(tag))
            if el is None:
                el = etree.SubElement(rPr, _qn(tag))
            el.set("typeface", face)

    # 数据行
    for ri, row in enumerate(rows):
        bg = N1 if ri % 2 == 0 else WH
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(bg)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci > 0 else PP_ALIGN.CENTER
            run = p.add_run()
            # 最后列（平台支持）用绿色粗体
            is_support = ci == len(col_headers) - 1
            run.text = val
            run.font.size = _Pt(10)
            run.font.bold = ci == 0 or is_support
            run.font.color.rgb = _rgb(P5 if is_support else (N9 if ci == 0 else N7))
            run.font.name = FONT_EN
            rPr = run._r.get_or_add_rPr()
            for tag, face in [("a:latin", FONT_EN), ("a:ea", FONT_CN)]:
                el = rPr.find(_qn(tag))
                if el is None:
                    el = etree.SubElement(rPr, _qn(tag))
                el.set("typeface", face)

    # 表格边框处理（通过 XML 移除所有边框，只保留行分隔）
    _clean_table_borders(tbl)

    footer(slide, "LAYOUT: 11_Table")
    return slide


def _clean_table_borders(tbl):
    """移除表格外框，保留浅灰行分隔线。"""
    try:
        from pptx.oxml.ns import qn as _qn
        for row in tbl.rows:
            for cell in row.cells:
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for side in ["lnL", "lnR", "lnT", "lnB"]:
                    ln = tcPr.find(_qn(f"a:{side}"))
                    if ln is not None:
                        tcPr.remove(ln)
    except Exception:
        pass


# ─── 版式 12：结尾致谢 ───────────────────────────────────────────────────────
def layout_12_closing(prs,
                      message="谢谢",
                      cta="预约产品演示 · 获取 ROI 测算",
                      contact="contact@futureark.ai  ·  400-000-0000",
                      tagline="Intelligence · Efficiency · Trust · Sustainability"):
    """
    版式 12 Closing：
    - 深色背景 #0E1216
    - 居中感谢语（大号白色）
    - CTA 区块（主绿背景）
    - Logo stacked-reverse 居中
    - 底部品牌 tagline + 主绿横线

    AI 使用说明：
      · message  主感谢语（48pt Bold 白）
      · cta      行动号召（绿色区块）
      · contact  联系方式
    """
    slide = _new(prs)
    set_bg(slide, N9)

    # 顶部细绿线
    rect(slide, l=0, t=0, w=W, h=Mm(2), fill=P5)

    # Logo stacked-reverse 居中
    p = LOGO_S_R
    if p and os.path.exists(p):
        lh = Mm(28)
        lw = int(lh * _img_aspect(p))
        slide.shapes.add_picture(p,
                                 int((W - lw) / 2), int(Mm(22)),
                                 width=lw, height=lh)

    # 主感谢语
    txb(slide, message,
        l=Mm(30), t=Mm(58),
        w=W - Mm(60), h=Mm(35),
        sz=52, bold=True, color=WH,
        align=PP_ALIGN.CENTER, name="closing_message")

    # A下划装饰
    deco(slide, "A_ul",
         l=(W - Mm(100)) / 2, t=Mm(90),
         w=Mm(100))

    # CTA 区块
    cta_y = Mm(98)
    cta_w = Mm(200)
    cta_l = (W - cta_w) / 2
    rect(slide, l=cta_l, t=cta_y, w=cta_w, h=Mm(16), fill=P5, rounded=True)
    txb(slide, f"→  {cta}",
        l=cta_l + Mm(10), t=cta_y + Mm(4),
        w=cta_w - Mm(20), h=Mm(10),
        sz=13, bold=True, color=WH,
        align=PP_ALIGN.CENTER)

    # 联系方式
    txb(slide, contact,
        l=Mm(30), t=Mm(120),
        w=W - Mm(60), h=Mm(10),
        sz=12, color=N4,
        align=PP_ALIGN.CENTER)

    # Tagline
    txb(slide, tagline,
        l=Mm(30), t=H - Mm(22),
        w=W - Mm(60), h=Mm(9),
        sz=10, color=N4,
        align=PP_ALIGN.CENTER)

    # 底部绿横线
    rect(slide, l=0, t=H - Mm(3.5), w=W, h=Mm(3.5), fill=P5)

    # 版式标识
    txb(slide, "LAYOUT: 12_Closing",
        l=W - Mm(70), t=H - Mm(10),
        w=Mm(60), h=Mm(8),
        sz=7, color=N7, align=PP_ALIGN.RIGHT)
    return slide


# ─── 版式索引页（首页说明） ───────────────────────────────────────────────────
def layout_00_index(prs):
    """版式目录页：列出所有 12 种版式供快速导览。"""
    slide = _new(prs)
    set_bg(slide, WH)
    rect(slide, l=0, t=0, w=W, h=Mm(2.5), fill=P5)

    txb(slide, "未来方舟 · 品牌 PPT 母版",
        l=ML, t=Mm(6), w=CW, h=Mm(18),
        sz=28, bold=True, color=N9)
    txb(slide, "Brand Presentation Master Template  ·  版式目录",
        l=ML, t=Mm(25), w=CW, h=Mm(9),
        sz=12, color=N4)
    rect(slide, l=ML, t=Mm(35), w=CW, h=Mm(0.4), fill=N2)

    entries = [
        ("01", "Cover",          "封面页",      "深色背景 · 大标题 · 装饰元素"),
        ("02", "Section",        "章节分隔",    "主绿背景 · 章节编号 · 画圈装饰"),
        ("03", "Title+Content",  "标准内容",    "白色 · 单列内容 · 图/文/流程"),
        ("04", "Two-Columns",    "两列内容",    "对比/痛点解法/Before-After"),
        ("05", "Three-Cards",    "三列卡片",    "功能/价值/步骤 卡片布局"),
        ("06", "Text+ImgRight",  "左文右图",    "功能说明 + 截图（60/40）"),
        ("07", "ImgLeft+Text",   "左图右文",    "版式 06 镜像，交替使用"),
        ("08", "Big-Stats",      "大数据指标",  "2×2 关键指标卡 · 大号数字"),
        ("09", "Timeline",       "水平时间轴",  "政策时间线 · 路线图 · 里程碑"),
        ("10", "Quote",          "引用证言",    "客户证言 · 引用块"),
        ("11", "Table",          "表格页",      "Header 绿色 · 斑马纹数据"),
        ("12", "Closing",        "结尾致谢",    "深色 · CTA · 联系方式"),
    ]

    col_n = 2
    row_n = (len(entries) + col_n - 1) // col_n
    ew = (CW - Mm(6)) / col_n
    ey = Mm(38)
    eh = (H - ey - Mm(6)) / row_n

    for i, (num, en_name, cn_name, desc) in enumerate(entries):
        col = i % col_n
        row = i // col_n
        ex = ML + col * (ew + Mm(6))
        sy = ey + row * eh

        # 序号 badge
        rect(slide, l=ex, t=sy + Mm(2), w=Mm(9), h=Mm(9), fill=P1)
        txb(slide, num, l=ex, t=sy + Mm(2),
            w=Mm(9), h=Mm(9), sz=9, bold=True, color=P5,
            align=PP_ALIGN.CENTER)
        # 名称
        txb(slide, f"{cn_name}  ·  {en_name}",
            l=ex + Mm(12), t=sy + Mm(2),
            w=ew - Mm(13), h=Mm(9), sz=12, bold=True, color=N9)
        # 描述
        txb(slide, desc,
            l=ex + Mm(12), t=sy + Mm(12),
            w=ew - Mm(13), h=Mm(8), sz=10, color=N4)

    logo_h(slide, reverse=False, r_mm=9, b_mm=3, h_mm=7)
    return slide


# ─── 主函数 ───────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("正在生成品牌 PPT 母版展示文件 ...")

    layout_00_index(prs)
    print("  ✓ 00 版式目录")

    layout_01_cover(prs)
    print("  ✓ 01 封面")

    layout_02_section(prs)
    print("  ✓ 02 章节分隔")

    layout_03_title_content(prs)
    print("  ✓ 03 标准内容")

    layout_04_two_col(prs)
    print("  ✓ 04 两列内容")

    layout_05_three_cards(prs)
    print("  ✓ 05 三列卡片")

    layout_06_text_image_right(prs)
    print("  ✓ 06 左文右图")

    layout_07_image_left_text(prs)
    print("  ✓ 07 左图右文")

    layout_08_big_stats(prs)
    print("  ✓ 08 大数据指标")

    layout_09_timeline(prs)
    print("  ✓ 09 水平时间轴")

    layout_10_quote(prs)
    print("  ✓ 10 引用证言")

    layout_11_table(prs)
    print("  ✓ 11 表格页")

    layout_12_closing(prs)
    print("  ✓ 12 结尾致谢")

    out = os.path.join(_DOCS, "brand_ppt_master.pptx")
    prs.save(out)
    print(f"\n✅ 已输出: {out}")
    print(f"   幻灯片数: {len(prs.slides)}  （1 目录 + 12 版式示范）")


if __name__ == "__main__":
    main()
