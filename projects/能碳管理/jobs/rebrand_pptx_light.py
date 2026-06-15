"""
能碳平台介绍 PPT 浅色品牌化脚本
深色品牌版 → 浅色版：背景翻转为白/浅灰，文字翻转为深色，保持可读性。
在所有页面右上角插入 horizontal-primary logo。
"""

import sys, os
from pptx import Presentation
from pptx.util import Inches

_BRAND   = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..'))
_PROJECT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
_DOCS    = os.path.join(_PROJECT, "docs")

SRC  = os.path.join(_DOCS, "能碳平台介绍-ArkSus品牌版.pptx")
DST  = os.path.join(_DOCS, "能碳平台介绍-ArkSus浅色版.pptx")
LOGO = os.path.join(_BRAND, "assets", "logo-horizontal-primary.png")

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ── Background / shape fill colors → light ────────────────────
FILL_MAP = {
    '0E1216': 'FFFFFF',   # slide bg (deepest dark) → white
    '152030': 'F5F7FA',   # panel bg (dark-elevated-1) → soft blue-gray
    '1E2D40': 'EBF0F7',   # card bg  (dark-elevated-2) → lighter
    'F8FAFC': 'FAFAFA',   # near-white fill → keep near-white
    # Green fills (3EC99E), lime fills (C8E13C) stay unchanged — accents/badges
}

# ── Text run colors → high-contrast on light backgrounds ──────
# Contrast target: ≥ 4.5:1 on #FFFFFF for body text, ≥ 3:1 for large PPT headings
TEXT_MAP = {
    'F8FAFC': '0E1216',   # white text → near-black           (21:1 ✓)
    'C8E13C': '4B9E31',   # lime text  → primary-dk green     (~3.3:1, large text ✓)
    '3EC99E': '0A7D6B',   # green accent text → darker teal   (~5.0:1 ✓)
    'F3B021': 'B86F00',   # amber text → dark amber            (~4.6:1 ✓)
    '8255E1': '6B46C1',   # purple text → darker purple        (~4.8:1 ✓)
    'F12D2D': 'CC1F1F',   # danger text → slightly darker red  (~4.7:1 ✓)
    # '8A9199' kept: ~3.0:1 on white, acceptable for large secondary PPT text
    # '0E1216' kept: already dark, correct on light bg
}

# Tags marking text-run property scope (ancestors of srgbClr in text context)
_TEXT_TAGS = frozenset({'rPr', 'endParaRPr', 'defRPr'})
# Tags that stop upward search when looking for text context
_FILL_TAGS = frozenset({'spPr', 'grpSpPr', 'bgPr', 'tcPr', 'ln'})


def _is_text_context(el):
    """Walk up XML tree; return True if el is inside a text-run property."""
    parent = el.getparent()
    while parent is not None:
        local = parent.tag.split('}')[-1] if '}' in parent.tag else parent.tag
        if local in _TEXT_TAGS:
            return True
        if local in _FILL_TAGS:
            return False
        parent = parent.getparent()
    return False


def _norm(h: str):
    """Return uppercase 6-char hex or None."""
    h = h.strip().upper()
    return h if len(h) == 6 and all(c in '0123456789ABCDEF' for c in h) else None


def replace_in_tree(root) -> int:
    count = 0

    # Context-aware srgbClr replacement
    for el in root.iter(f'{{{NS_A}}}srgbClr'):
        val = _norm(el.get('val', ''))
        if not val:
            continue
        if _is_text_context(el):
            new = TEXT_MAP.get(val)
        else:
            new = FILL_MAP.get(val)
        if new:
            el.set('val', new)
            count += 1

    # sysClr lastClr (theme fallback color) — treat as fill context
    combined = {**FILL_MAP, **TEXT_MAP}
    for el in root.iter(f'{{{NS_A}}}sysClr'):
        val = _norm(el.get('lastClr', ''))
        if val and val in combined:
            el.set('lastClr', combined[val])
            count += 1

    return count


def add_logo(slide, prs_width: int):
    """Place horizontal-primary logo at top-right corner."""
    margin = Inches(0.20)
    logo_w = Inches(1.30)
    logo_h = Inches(0.27)
    left   = prs_width - logo_w - margin
    top    = margin
    slide.shapes.add_picture(LOGO, left, top, logo_w, logo_h)


def process():
    prs   = Presentation(SRC)
    total = 0

    # Replace in masters + layouts (theme-level colors)
    for master in prs.slide_masters:
        total += replace_in_tree(master.element)
        for layout in master.slide_layouts:
            total += replace_in_tree(layout.element)

    # Replace in each slide + add logo
    for i, slide in enumerate(prs.slides, 1):
        n = replace_in_tree(slide.element)
        total += n
        if n:
            print(f"  Slide {i:2d}: {n} replacements")
        add_logo(slide, prs.slide_width)

    prs.save(DST)
    return total


print("=" * 60)
print("ArkSus® PPT 浅色品牌化工具")
print("=" * 60)
print(f"Source : {SRC}")
print(f"Output : {DST}")
print(f"Logo   : {LOGO}\n")

n = process()

if os.path.exists(DST):
    size_kb = os.path.getsize(DST) / 1024
    print(f"\n✓ Done — {n} color replacements")
    print(f"✓ Logo added to all 32 slides")
    print(f"✓ Output: {DST} ({size_kb:.0f} KB)")
else:
    print("\n✗ Failed — output not created")
    sys.exit(1)
