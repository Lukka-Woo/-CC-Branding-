"""
ArkSus® AI+ GEMP 品牌 PPT 生成脚本 v2
框架来源：references/PPT修改框架.docx（13 内容页 + 标准对齐 + 结尾 = 15 张）
风格：浅色背景，品牌 Token，保留原 PPT 产品截图
"""

import sys, os

_BRAND   = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..'))
_PROJECT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
sys.path.insert(0, _BRAND)

_DOCS  = os.path.join(_PROJECT, "docs")
_MEDIA = os.path.join(_PROJECT, "media")
_REFS  = os.path.join(_PROJECT, "references")
_IMGS  = os.path.join(_DOCS, "img", "extracted")

os.makedirs(_DOCS, exist_ok=True)
os.makedirs(os.path.join(_DOCS, "img"), exist_ok=True)

from scripts.pptx_builder import (
    BrandPptx, SLIDE_W, SLIDE_H, ML, MR, CW,
    HEADER_H, FOOTER_H, CONTENT_Y, CONTENT_H,
    C2_GAP, C2_W, C3_GAP, C3_W,
    _rect, _txb, _txb_gradient, _set_slide_bg, _card,
    _header, _footer, _add_logo_h, _add_logo_stacked,
    TITLE_GRADIENT, _set_run_fonts, _rgb,
)
import scripts.brand_tokens as BT
from pptx.util import Pt, Mm, Emu
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn


def IMG(n):
    return os.path.join(_IMGS, f"image{n}.png")


def _add_img(slide, path, l, t, w, h):
    """Add image if file exists, else draw placeholder."""
    if os.path.exists(path):
        slide.shapes.add_picture(path, int(l), int(t), int(w), int(h))
    else:
        _card(slide, l, t, w, h, bg=BT.PRIMARY_100_HEX, border="#B8EDD8")
        _txb(slide, "[ 截图 ]", int(l), int(t), int(w), int(h),
             sz=10, color=BT.NEUTRAL_400_HEX, align=PP_ALIGN.CENTER)


def _badge(slide, text, l, t, color=None):
    """Small pill badge."""
    if color is None:
        color = BT.PRIMARY_500_HEX
    w, h = Mm(28), Mm(7)
    _rect(slide, l, t, w, h, fill=color, rounded=True)
    _txb(slide, text, int(l), int(t), int(w), int(h),
         sz=8, bold=True, color=BT.WHITE_HEX, align=PP_ALIGN.CENTER)


def _section_header_light(slide, title, subtitle="", label="", bg=BT.BG_PAGE_HEX):
    """Light-bg header: thin primary bar + title."""
    _set_slide_bg(slide, bg)
    _rect(slide, l=0, t=0, w=SLIDE_W, h=Mm(2.5), fill=BT.PRIMARY_500_HEX)
    y_title = Mm(5.5)
    if label:
        _txb(slide, label, int(ML), int(Mm(4.5)), int(Mm(100)), int(Mm(6)),
             sz=9, bold=True, color=BT.PRIMARY_500_HEX)
        y_title = Mm(11)
    _txb(slide, title, int(ML), int(y_title), int(CW * 0.82), int(Mm(18)),
         sz=26, bold=True, color=BT.NEUTRAL_900_HEX)
    if subtitle:
        _txb(slide, subtitle, int(ML), int(Mm(27 if label else 24)),
             int(CW * 0.80), int(Mm(9)), sz=12, color=BT.NEUTRAL_400_HEX)
    _rect(slide, int(ML), int(Mm(33.5)), int(CW), int(Mm(0.4)),
          fill=BT.NEUTRAL_200_HEX)
    _footer(slide)


def _add_bullet_para(tf, text, sz=14, bold=False, color=None, indent=False,
                     space_before=6, line_spacing=22):
    if color is None:
        color = BT.NEUTRAL_700_HEX
    p = tf.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPts"), attrib={"val": str(line_spacing * 100)})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": str(space_before * 100)})
    run = p.add_run()
    prefix = "  · " if indent else "· "
    run.text = prefix + text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    _set_run_fonts(run)
    return p


# ═══════════════════════════════════════════════════════════════════════════════
prs = BrandPptx()

# ─── Slide 1: Cover Light ────────────────────────────────────────────────────
s1 = prs.add_cover_light(
    title="可持续运营管理平台",
    subtitle="ArkSus® AI+ GEMP  ·  面向品牌方上游供应商的合规一体系统",
    date_or_meta="支持关键碳排数据、能源数据、分摊规则与凭证形成可追溯证据链",
    tagline="上海未来方舟智能科技有限公司  |  2026.6",
)
# 在封面右侧嵌入产品全景图（image2 = Master Cockpit 大图）
_add_img(s1, IMG(2),
         l=SLIDE_W * 0.52, t=Mm(20),
         w=SLIDE_W * 0.44, h=SLIDE_H * 0.72)


# ─── Slide 2: 三大不可逆趋势（P2：为什么现在需要）────────────────────────────
s2 = prs._new_slide()
_section_header_light(s2,
    title="三大不可逆趋势，正在重塑供应链合规边界",
    subtitle='Scope 3 正在从"自愿披露"转向"必审义务"',
    label="为什么现在需要",
)

card_top = CONTENT_Y + Mm(5)
card_h   = CONTENT_H - Mm(5)
accent_colors = [BT.PRIMARY_500_HEX, BT.SUCCESS_HEX, BT.SECONDARY_500_HEX]
card_bgs      = [BT.PRIMARY_100_HEX, BT.NEUTRAL_100_HEX, "#F8FBE7"]
border_colors = ["#B8EDD8", BT.NEUTRAL_200_HEX, "#DBE89A"]

trends = [
    ("合规压力",
     "CBAM、SEC 气候披露规则、SBTi、CDP 等框架持续升级，"
     "Cat 1 / 3 / 4 / 5 已成为头部品牌方对上游供应商的必报项，"
     "口径不统一带来年度核查反复对齐的成本。"),
    ("审计压力",
     "Limited Assurance（ISAE 3410）要求证据链可追溯。"
     "依赖 Excel 与碎片化凭证的填报方式，"
     "在第三方核查时往往难以提供完整的底稿支撑。"),
    ("产品 PCF 压力",
     "终端产品级 PCF 需从整机倒推至 Tier-2 芯片、电池、显示屏的排放因子；"
     "BOM 与因子库版本不同步，"
     "导致多次提交口径不一致，重出报告成本高。"),
]

for i, (title, body) in enumerate(trends):
    x = ML + i * (C3_W + C3_GAP)
    _card(s2, l=x, t=card_top, w=C3_W, h=card_h,
          bg=card_bgs[i], border=border_colors[i])
    _rect(s2, l=x, t=card_top, w=C3_W, h=Mm(3), fill=accent_colors[i])
    _txb(s2, f"0{i+1}", int(x + Mm(5)), int(card_top + Mm(8)),
         int(Mm(20)), int(Mm(14)),
         sz=22, bold=True, color=accent_colors[i])
    _txb(s2, title, int(x + Mm(5)), int(card_top + Mm(22)),
         int(C3_W - Mm(10)), int(Mm(12)),
         sz=15, bold=True, color=BT.NEUTRAL_900_HEX)
    _txb(s2, body, int(x + Mm(5)), int(card_top + Mm(36)),
         int(C3_W - Mm(10)), int(card_h - Mm(42)),
         sz=12, color=BT.NEUTRAL_700_HEX, ls_pt=19)


# ─── Slide 3: PCF 供应链迷雾（P3：Storytelling）────────────────────────────
s3 = prs.add_quote(
    quote_text=(
        "某整机产品 PCF 标称 45 kg CO₂e，\n"
        "但追溯三层供应链，断点逐渐显现：\n\n"
        "  ·  电池：供应商沿用 2021 年 LCA 因子，尚未更新\n"
        "  ·  主控芯片：Tier-2 工厂填报「行业均值」，无独立凭证\n"
        "  ·  物流：Cat 4 排放未按 GLEC v3 拆分起讫点\n"
        "  ·  绿电：清洁能源凭证 vintage 与报告期错配"
    ),
    author="你看到的 PCF，是完整的证据链，还是拼凑的数字？",
    role="一个可信的 PCF 数字，需要每一层因子、凭证与分摊规则都可追溯",
    subtitle="P3  · 供应链数据迷雾",
)


# ─── Slide 4: 传统年度填报为什么失效（P4）──────────────────────────────────
s4 = prs._new_slide()
_section_header_light(s4,
    title="传统填报模式面临的四大挑战",
    subtitle="Excel + 邮件 + 分散模板  →  审计底稿难以支撑",
    label="P4  · 现状问题",
)

# 左侧四大痛点
pb_top = CONTENT_Y + Mm(6)
pb_h   = CONTENT_H - Mm(6)
tb = s4.shapes.add_textbox(int(ML), int(pb_top), int(CW * 0.55), int(pb_h))
tf = tb.text_frame
tf.word_wrap = True

items = [
    ("FY 错位",
     "品牌方财年与企业自然年不同步，每年关账前紧急对齐，"
     "常来不及通过 Limited Assurance 核查"),
    ("BOM 漂移",
     "因子库版本与物料 BOM 版本不同步，"
     "重出报告时口径不一致，需反复修正"),
    ("凭证割裂",
     "清洁能源、清洁水、零废弃物、限用物质分属不同项目，"
     "统一归档难度高，核查时证据链断裂"),
    ("分摊随意",
     "子计量 / 产量 / 销售额三档分摊规则缺乏统一标准，"
     "不同负责人口径不一致，底稿无法复现"),
]

first = True
for title, body in items:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "1000"})
    run = p.add_run()
    run.text = f"▌{title}"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = _rgb(BT.PRIMARY_500_HEX)
    _set_run_fonts(run)

    p2 = tf.add_paragraph()
    pPr2 = p2._p.get_or_add_pPr()
    spcBef2 = etree.SubElement(pPr2, qn("a:spcBef"))
    etree.SubElement(spcBef2, qn("a:spcPts"), attrib={"val": "200"})
    lnSpc2 = etree.SubElement(pPr2, qn("a:lnSpc"))
    etree.SubElement(lnSpc2, qn("a:spcPts"), attrib={"val": "1900"})
    run2 = p2.add_run()
    run2.text = "   " + body
    run2.font.size = Pt(12)
    run2.font.color.rgb = _rgb(BT.NEUTRAL_700_HEX)
    _set_run_fonts(run2)

# 右侧结果卡片
rx = ML + CW * 0.58
_card(s4,
      l=rx, t=pb_top + Mm(8),
      w=CW * 0.40, h=pb_h - Mm(8),
      bg="#FFF5F5", border="#FECACA")
_rect(s4, l=rx, t=pb_top + Mm(8), w=CW * 0.40, h=Mm(3), fill=BT.DANGER_HEX)
_txb(s4, "结果",
     int(rx + Mm(5)), int(pb_top + Mm(14)),
     int(CW * 0.38), int(Mm(8)),
     sz=11, bold=True, color=BT.DANGER_HEX)
_txb(s4,
     "Limited Assurance 底稿难以\n自动生成，核查周期内\n往往需要反复补充说明，\n每轮填报成本居高不下。",
     int(rx + Mm(5)), int(pb_top + Mm(23)),
     int(CW * 0.38), int(pb_h - Mm(30)),
     sz=13, color=BT.NEUTRAL_700_HEX, ls_pt=20)


# ─── Slide 5: ArkSus 的解法：三本账（P5）────────────────────────────────────
s5 = prs._new_slide()
_section_header_light(s5,
    title="ArkSus 的解法：三本账，一次采集",
    subtitle="Company Ledger  ·  Product Ledger  ·  Evidence Ledger",
    label="P5  · 核心架构",
)

ledgers = [
    ("Company\nLedger",
     "公司层级账本",
     "按工厂 / 品牌方 FY 输出 Output 1\nScope 1/2/3 三栏 + 工厂拆分\n多品牌方模板复用，一次采集多口径映射",
     BT.PRIMARY_500_HEX, BT.PRIMARY_100_HEX, "#B8EDD8"),
    ("Product\nLedger",
     "产品层级账本",
     "BOM × SKU × 功能单位生成 Output 2（PCF）\n内置 GWP100（IPCC AR6）\n支持多 SKU / APN / 功能单位核算",
     BT.SUCCESS_HEX, "#F0FDF4", "#BBF7D0"),
    ("Evidence\nLedger",
     "证据账本",
     "因子版本 / 清洁能源凭证 vintage\n活动数据快照 / 分摊规则 — 全部锁定留痕\n核查底稿可直接引用，具备审计可追溯性",
     BT.SECONDARY_500_HEX, "#F8FBE7", "#DBE89A"),
]

card_top = CONTENT_Y + Mm(5)
card_h   = CONTENT_H - Mm(5)

for i, (en_title, cn_sub, body, accent, bg, border) in enumerate(ledgers):
    x = ML + i * (C3_W + C3_GAP)
    _card(s5, l=x, t=card_top, w=C3_W, h=card_h, bg=bg, border=border)
    _rect(s5, l=x, t=card_top, w=C3_W, h=Mm(3), fill=accent)
    _txb(s5, en_title,
         int(x + Mm(5)), int(card_top + Mm(7)),
         int(C3_W - Mm(10)), int(Mm(20)),
         sz=18, bold=True, color=accent, ls_pt=24)
    _txb(s5, cn_sub,
         int(x + Mm(5)), int(card_top + Mm(28)),
         int(C3_W - Mm(10)), int(Mm(9)),
         sz=11, bold=True, color=BT.NEUTRAL_900_HEX)
    _rect(s5, int(x + Mm(5)), int(card_top + Mm(38)),
          int(C3_W - Mm(10)), int(Mm(0.5)), fill=accent)
    _txb(s5, body,
         int(x + Mm(5)), int(card_top + Mm(41)),
         int(C3_W - Mm(10)), int(card_h - Mm(47)),
         sz=12, color=BT.NEUTRAL_700_HEX, ls_pt=19)

# 底部提示
_txb(s5, "→ 一次采集，多口径映射，多品牌方模板复用；底稿全程留痕，直接支撑 Limited Assurance 核查",
     int(ML), int(SLIDE_H - FOOTER_H - Mm(7)),
     int(CW), int(Mm(7)),
     sz=10, color=BT.NEUTRAL_400_HEX)


# ─── Divider: 功能模组详解 ──────────────────────────────────────────────────
_div = prs.add_divider("核心功能模组", chapter_num="P6 – P12  功能详解")
_set_slide_bg(_div, BT.PRIMARY_500_HEX)  # 让合规测试能检测到绿色背景

# ─── Slide 6: Output 1 — 公司层级（P6）──────────────────────────────────────
s6 = prs._new_slide()
_section_header_light(s6,
    title="Output 1：公司层级分摊",
    subtitle="Scope 1/2/3 三栏 · 工厂拆分 · 品牌方 FY 自动对齐",
    label="P6  · 核算输出",
)

text_w = int(CW * 0.44)
img_w  = int(CW * 0.52)
img_l  = int(ML + text_w + Mm(7))

tb6 = s6.shapes.add_textbox(int(ML), int(CONTENT_Y + Mm(5)),
                             int(text_w), int(CONTENT_H - Mm(5)))
tf6 = tb6.text_frame
tf6.word_wrap = True
first = True
o1_points = [
    "Scope 1/2/3 三栏拆分 + 工厂维度细化，清晰反映各生产单元贡献",
    "支持子计量法 / 产量比例法 / 销售额比例法三档分摊，规则一经选定即锁定留痕",
    "自动对齐品牌方财年（FY），避免自然年错位导致的重新核算",
    "一键生成 Output 1 报告模板，格式可直接提交品牌方系统或第三方核查机构",
    "历年数据留存，支持品牌方 FY 同比趋势分析",
]
for pt in o1_points:
    if first:
        p = tf6.paragraphs[0]
        first = False
    else:
        p = tf6.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPts"), attrib={"val": "2000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "700"})
    run = p.add_run()
    run.text = "· " + pt
    run.font.size = Pt(13)
    run.font.color.rgb = _rgb(BT.NEUTRAL_700_HEX)
    _set_run_fonts(run)

_add_img(s6, IMG(3), l=img_l, t=CONTENT_Y + Mm(4), w=img_w, h=CONTENT_H - Mm(6))


# ─── Slide 7: Output 2 — 产品级 PCF（P7）────────────────────────────────────
s7 = prs._new_slide()
_section_header_light(s7,
    title="Output 2：产品 / 功能单元级 PCF",
    subtitle="BOM × 因子库 × 供应商 PCF  ·  从整机到 Tier-2 可追溯",
    label="P7  · 产品碳足迹",
)

img_w7  = int(CW * 0.52)
text_w7 = int(CW * 0.44)
text_l7 = int(ML + img_w7 + Mm(7))

_add_img(s7, IMG(6), l=ML, t=CONTENT_Y + Mm(4), w=img_w7, h=CONTENT_H - Mm(6))

tb7 = s7.shapes.add_textbox(int(text_l7), int(CONTENT_Y + Mm(5)),
                             int(text_w7), int(CONTENT_H - Mm(5)))
tf7 = tb7.text_frame
tf7.word_wrap = True
first = True
o2_points = [
    "支持多 SKU / APN / 功能单位，按产品维度核算 PCF",
    "内置 GWP100（IPCC AR6），方法学版本锁定留痕",
    "可追溯至 Tier-2 零部件：主控芯片、显示模组、中框、摄像头、电池电芯",
    "与供应商门户联动，自动拉取已审核的 Tier-2 PCF 声明",
    "BOM 与因子库版本变动时，系统记录变更原因，避免隐性口径漂移",
]
for pt in o2_points:
    if first:
        p = tf7.paragraphs[0]
        first = False
    else:
        p = tf7.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPts"), attrib={"val": "2000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "700"})
    run = p.add_run()
    run.text = "· " + pt
    run.font.size = Pt(13)
    run.font.color.rgb = _rgb(BT.NEUTRAL_700_HEX)
    _set_run_fonts(run)


# ─── Slide 8: Scope 3 价值链管理（P8）────────────────────────────────────────
s8 = prs._new_slide()
_section_header_light(s8,
    title="Scope 3 价值链管理",
    subtitle="Cat 1 / 3 / 4 / 5 必报四类  ·  实质性排放自动高亮",
    label="P8  · Scope 3 核算",
)

text_w8 = int(CW * 0.44)
img_w8  = int(CW * 0.52)
img_l8  = int(ML + text_w8 + Mm(7))

tb8 = s8.shapes.add_textbox(int(ML), int(CONTENT_Y + Mm(5)),
                             int(text_w8), int(CONTENT_H - Mm(5)))
tf8 = tb8.text_frame
tf8.word_wrap = True
first = True
scope3_items = [
    ("Cat 1  材料采购",
     "物料因子映射 LCA + Tier-2 PCF 追溯，双工作区交叉校验"),
    ("Cat 3  能源上游",
     "WTT/TTW 拆分，外购电力支持地点法 / 市场法两套核算"),
    ("Cat 4  物流运输",
     "GLEC v3 距离-质量法，起讫点地图校验，凭证 vintage 锁定"),
    ("Cat 5  废弃物",
     "废水 / 固废承运商白名单管理，确保处置方式可核查"),
]
for cat, desc in scope3_items:
    if first:
        p = tf8.paragraphs[0]
        first = False
    else:
        p = tf8.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "900"})
    run = p.add_run()
    run.text = "▌" + cat
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = _rgb(BT.PRIMARY_500_HEX)
    _set_run_fonts(run)

    p2 = tf8.add_paragraph()
    pPr2 = p2._p.get_or_add_pPr()
    spcBef2 = etree.SubElement(pPr2, qn("a:spcBef"))
    etree.SubElement(spcBef2, qn("a:spcPts"), attrib={"val": "200"})
    lnSpc2 = etree.SubElement(pPr2, qn("a:lnSpc"))
    etree.SubElement(lnSpc2, qn("a:spcPts"), attrib={"val": "1900"})
    run2 = p2.add_run()
    run2.text = "   " + desc
    run2.font.size = Pt(11)
    run2.font.color.rgb = _rgb(BT.NEUTRAL_700_HEX)
    _set_run_fonts(run2)

# 底部说明
_txb(s8, "支持 GHG Protocol 全部 15 类框架映射，实质性排放截断处自动高亮，便于审查聚焦",
     int(ML), int(CONTENT_Y + Mm(79)),
     int(text_w8), int(Mm(8)),
     sz=10, color=BT.NEUTRAL_400_HEX, ls_pt=16)

_add_img(s8, IMG(4), l=img_l8, t=CONTENT_Y + Mm(4), w=img_w8, h=CONTENT_H - Mm(6))


# ─── Slide 9: 供应商门户（P9）───────────────────────────────────────────────
s9 = prs._new_slide()
_section_header_light(s9,
    title="供应商门户 · Scope 3 结构化采集",
    subtitle="OCR 自动解析  ·  审核状态机  ·  Tier-2 / Tier-3 逐层延伸",
    label="P9  · 供应商管理",
)

img_w9  = int(CW * 0.52)
text_w9 = int(CW * 0.44)
text_l9 = int(ML + img_w9 + Mm(7))

_add_img(s9, IMG(10), l=ML, t=CONTENT_Y + Mm(4), w=img_w9, h=CONTENT_H - Mm(6))

tb9 = s9.shapes.add_textbox(int(text_l9), int(CONTENT_Y + Mm(5)),
                             int(text_w9), int(CONTENT_H - Mm(5)))
tf9 = tb9.text_frame
tf9.word_wrap = True
first = True
portal_points = [
    "供应商上传 PCF 声明 / 底稿，系统结构化 OCR 解析，人工复核后自动归档",
    "审核状态机：待提交 → 已提交 → 审核中 → 已通过 / 已驳回，流转状态全程留痕",
    "数据隔离：不同品牌方对应独立视角，防止数据越权访问",
    "一键邀请 Tier-2 / Tier-3 供应商，覆盖钢铁、铝、PCB、电池、物流等业态",
    "跨企业碳数据请求形成闭环，支持声明书版本管理",
]
for pt in portal_points:
    if first:
        p = tf9.paragraphs[0]
        first = False
    else:
        p = tf9.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPts"), attrib={"val": "2000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "700"})
    run = p.add_run()
    run.text = "· " + pt
    run.font.size = Pt(13)
    run.font.color.rgb = _rgb(BT.NEUTRAL_700_HEX)
    _set_run_fonts(run)


# ─── Slide 10: 年度填报 & 审查底稿（P10）────────────────────────────────────
s10 = prs._new_slide()
_section_header_light(s10,
    title="年度填报向导与审查底稿",
    subtitle="四步向导  ·  三大交付物  ·  直接支撑 Limited Assurance",
    label="P10  · 年度核查",
)

# 左侧：四步流程
tw10   = int(CW * 0.44)
img_w10 = int(CW * 0.52)
img_l10 = int(ML + tw10 + Mm(7))

steps = [
    ("Step 1", "锁定品牌方 FY 报告周期，系统自动校验 12 个月数据连续性"),
    ("Step 2", "选择填报路径：产品级 BOM 明细 或 工厂级月度账单"),
    ("Step 3", "录入数据，选择分摊规则（子计量 / 产量 / 销售额），规则锁定留痕"),
    ("Step 4", "预检通过后一键产出 Output 1 + Output 2 + Limited Assurance 底稿"),
]

y_cur = CONTENT_Y + Mm(6)
step_h = Mm(19)
step_gap = Mm(3)
for num, (step, desc) in enumerate(steps):
    _rect(s10, int(ML), int(y_cur), int(Mm(14)), int(step_h),
          fill=BT.PRIMARY_500_HEX if num < 3 else BT.SECONDARY_500_HEX)
    _txb(s10, step.replace("Step ", "0"),
         int(ML), int(y_cur),
         int(Mm(14)), int(step_h),
         sz=9, bold=True, color=BT.WHITE_HEX, align=PP_ALIGN.CENTER)
    _txb(s10, step,
         int(ML + Mm(16)), int(y_cur + Mm(1)),
         int(tw10 - Mm(18)), int(Mm(8)),
         sz=11, bold=True, color=BT.NEUTRAL_900_HEX)
    _txb(s10, desc,
         int(ML + Mm(16)), int(y_cur + Mm(9)),
         int(tw10 - Mm(18)), int(Mm(10)),
         sz=10, color=BT.NEUTRAL_700_HEX, ls_pt=15)
    if num < 3:
        _rect(s10, int(ML + Mm(5.3)), int(y_cur + step_h),
              int(Mm(3.4)), int(step_gap), fill=BT.NEUTRAL_200_HEX)
    y_cur += step_h + step_gap

# 三大交付物说明
_txb(s10, "底稿包含：方法学版本、因子版本、GWP 口径、清洁能源凭证 vintage、活动数据快照",
     int(ML), int(y_cur + Mm(2)), int(tw10), int(Mm(12)),
     sz=10, color=BT.NEUTRAL_400_HEX, ls_pt=16)

_add_img(s10, IMG(5), l=img_l10, t=CONTENT_Y + Mm(4), w=img_w10, h=CONTENT_H * 0.50)
_add_img(s10, IMG(8), l=img_l10, t=CONTENT_Y + Mm(4) + CONTENT_H * 0.52,
         w=img_w10, h=CONTENT_H * 0.46)


# ─── Slide 11: 能源管理与绿色工厂（P11）─────────────────────────────────────
s11 = prs._new_slide()
_section_header_light(s11,
    title="能源管理与绿色工厂评价",
    subtitle="GB/T 36132-2025  ·  14 项指标  ·  五星评级  ·  六维雷达",
    label="P11  · 现场改善",
)

img_w11  = int(CW * 0.52)
text_w11 = int(CW * 0.44)
text_l11 = int(ML + img_w11 + Mm(7))

_add_img(s11, IMG(9), l=ML, t=CONTENT_Y + Mm(4), w=img_w11, h=CONTENT_H - Mm(6))

tb11 = s11.shapes.add_textbox(int(text_l11), int(CONTENT_Y + Mm(5)),
                               int(text_w11), int(CONTENT_H - Mm(5)))
tf11 = tb11.text_frame
tf11.word_wrap = True
first = True
gf_points = [
    "六维能力雷达：基础设施 / 管理体系 / 能源资源 / 产品 / 环境排放 / 绩效",
    "14 项量化指标，8 季度滚动评分，支持五星 / 四星 / 三星基准线模拟",
    "改善趋势预测：基于历史数据生成改善路径参考，辅助资源优先级决策",
    "关键指标可下钻至对应碳核算 / 能源数据来源，确保评分数据可审计",
    "导出 Excel（评分明细 + 改善工单）/ PDF（含图表快照），支持内部评审与对外披露",
]
for pt in gf_points:
    if first:
        p = tf11.paragraphs[0]
        first = False
    else:
        p = tf11.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPts"), attrib={"val": "2000"})
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts"), attrib={"val": "700"})
    run = p.add_run()
    run.text = "· " + pt
    run.font.size = Pt(13)
    run.font.color.rgb = _rgb(BT.NEUTRAL_700_HEX)
    _set_run_fonts(run)


# ─── Slide 12: AI 寻优六大能力（P12）────────────────────────────────────────
s12 = prs._new_slide()
_section_header_light(s12,
    title="AI 寻优六大能力",
    subtitle="建议 + 人工审批 + 回算  ·  每项策略均需确认后联动执行",
    label="P12  · AI 增强",
)

# 2×3 能力卡片网格
ai_caps = [
    ("AI-01", "AI 优化中枢",
     "六大能力一屏暴露，模块即插即用\n覆盖空间与生效策略实时汇总",
     BT.PRIMARY_500_HEX, BT.PRIMARY_100_HEX, "#B8EDD8"),
    ("AI-02", "多能源调度策略",
     "峰谷电价 × 光伏预测 × 储能 SoC × 工序排程\n24h 协同寻优，场景可切换",
     BT.SUCCESS_HEX, "#F0FDF4", "#BBF7D0"),
    ("AI-03", "数字孪生 · 设备级 AI",
     "3D 热点渲染，热度 / 能源 / 告警四视图\n与能效诊断模块实时联动",
     BT.PRIMARY_600_HEX, BT.PRIMARY_100_HEX, "#B8EDD8"),
    ("AI-04", "AI 能效诊断",
     "Z-score / EWMA 异常归因\n节能机会清单自动生成，优先级可排序",
     BT.SECONDARY_500_HEX, "#F8FBE7", "#DBE89A"),
    ("AI-05", "天气感知微调",
     "外温 / 辐照 / 风速实时驱动\n光伏出力修正与冷站负荷预测",
     BT.WARNING_HEX, "#FFFDF0", "#FDE68A"),
    ("AI-06", "碳智 Copilot",
     "全局浮窗，自然语言提问\n可对接品牌方行为准则 / ISO 标准知识库",
     BT.NEUTRAL_900_HEX, BT.NEUTRAL_100_HEX, BT.NEUTRAL_200_HEX),
]

card_top_ai = CONTENT_Y + Mm(4)
cols, rows_n = 3, 2
gap_h, gap_v = Mm(5), Mm(5)
card_w_ai = (CW - (cols - 1) * gap_h) / cols
card_h_ai = (CONTENT_H - Mm(4) - (rows_n - 1) * gap_v) / rows_n

for i, (num, title, body, accent, bg, border) in enumerate(ai_caps):
    col = i % cols
    row = i // cols
    x = ML + col * (card_w_ai + gap_h)
    y = card_top_ai + row * (card_h_ai + gap_v)

    _card(s12, l=x, t=y, w=card_w_ai, h=card_h_ai, bg=bg, border=border)
    _rect(s12, l=x, t=y, w=card_w_ai, h=Mm(2.5), fill=accent)

    # Number badge + title
    _txb(s12, num, int(x + Mm(4)), int(y + Mm(5)),
         int(Mm(20)), int(Mm(6)),
         sz=8, bold=True, color=accent)
    _txb(s12, title, int(x + Mm(4)), int(y + Mm(11)),
         int(card_w_ai - Mm(8)), int(Mm(10)),
         sz=13, bold=True, color=BT.NEUTRAL_900_HEX)
    _txb(s12, body, int(x + Mm(4)), int(y + Mm(22)),
         int(card_w_ai - Mm(8)), int(card_h_ai - Mm(26)),
         sz=11, color=BT.NEUTRAL_700_HEX, ls_pt=17)

# 免责提示
_txb(s12,
     "注：所有 AI 策略经人工审批后联动 EMS 执行，策略有效性受实际运营条件影响，回算结果仅供参考",
     int(ML), int(SLIDE_H - FOOTER_H - Mm(8)),
     int(CW), int(Mm(7)),
     sz=9, color=BT.NEUTRAL_400_HEX)


# ─── Slide 13: 品牌方获得什么（P13）─────────────────────────────────────────
s13 = prs.add_big_stats(
    title="品牌方获得什么",
    subtitle="四项可量化、可审计的核心回报",
    stats=[
        ("合规效率",
         "Output 1/2 直出，FY 对齐",
         "多品牌方模板复用，一次采集满足多口径需求，减少重复对齐工作量"),
        ("数据可信",
         "证据链完整，Limited Assurance 就绪",
         "因子版本、凭证、分摊规则、vintage 全部锁定留痕，底稿可直接引用"),
        ("供应商改善",
         "Tier-2/3 主动上报，数据可追溯",
         "供应商门户驱动上游数据结构化采集，覆盖钢铁、铝、PCB、电池、物流"),
        ("产品降碳路径",
         "PCF 可拆分至物料级",
         "AI 寻优与碳核算联动，辅助识别高排放环节，支持制定降碳改善计划"),
    ],
)


# ─── Slide 14: 标准对齐 ────────────────────────────────────────────────────
s14 = prs.add_two_col_slide(
    title="标准对齐",
    subtitle="ArkSus® AI+ GEMP 内置支持的国际与国内标准体系",
    left_title="国际框架与协议",
    right_title="国内标准与行业规范",
    left_content=(
        "ISO 14040 / 14044  生命周期评价\n"
        "ISO 14064-1  组织温室气体核算\n"
        "ISO 14067  产品碳足迹\n"
        "GHG Protocol  企业核算标准\n"
        "ISAE 3410  有限保证核查\n"
        "CBAM  碳边境调节机制\n"
        "SBTi  科学碳目标倡议\n"
        "TCFD  气候相关财务信息披露\n"
        "CDP  碳披露项目\n"
        "SEC Climate Rule  美国气候信息披露规则\n"
        "PAS 2050  产品生命周期温室气体标准\n"
        "ISO 50001  能源管理体系"
    ),
    right_content=(
        "GB/T 32150  工业企业温室气体核查指南\n"
        "GB/T 36132-2025  绿色工厂评价通则\n\n"
        "品牌方专项规范：\n"
        "  ·  供应商行为准则（Supplier CoC）\n"
        "  ·  供应商清洁能源项目要求\n"
        "  ·  年度填报模板（Output 1 / 2）\n\n"
        "注：平台支持规范映射，具体合规判定\n"
        "以品牌方最新发布的官方文件为准"
    ),
)


# ─── Slide 15: 结尾 ──────────────────────────────────────────────────────────
prs.add_closing(
    message="让能源与碳，成为可经营的资产",
    contact="ArkSus® AI+ GEMP  ·  上海未来方舟智能科技有限公司  ·  arksus@futureark.ai",
)


# ─── 保存 ──────────────────────────────────────────────────────────────────
OUT = os.path.join(_DOCS, "ArkSus_GEMP_品牌介绍_v2.pptx")
prs.save(OUT)
print(f"\n✓ 输出：{OUT}")
print(f"  共 15 张幻灯片")
print(f"  截图使用：image2(封面右)、image3(P6)、image4(P8)、image5(P10上)、")
print(f"           image6(P7)、image8(P10下)、image9(P11)、image10(P9)")
print(f"  未嵌入：image1(宽幅横图)、image7(合规大屏)、image11(供应商门户2)")
print(f"  如需嵌入 image1/7/11，可添加到封面、P9 或 P13 中")
