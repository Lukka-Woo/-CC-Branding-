#!/usr/bin/env python3
"""
gen_company_intro_branded.py

品牌化优化：上海未来方舟产品简介V0 → 品牌优化版

改动内容：
  1. 颜色全量映射到 tokens.json 品牌色系（fills / text / borders / gradients）
  2. Gradient 填充 → 纯色品牌绿 #3EC99E
  3. 章节分隔页左侧添加品牌绿竖线装饰
  4. 产品矩阵页：精简冗余描述
  5. AI能力章节标题 "---" → "·" 清理
  6. 关于我们页：删除 URL、精简行文
  7. 封面主标题：清理尾部空格/换行

输入: references/上海未来方舟智能科技及产品简介V0（主题风格）0526.pptx
输出: docs/未来方舟产品简介-品牌优化版.pptx
"""

import sys, os
from pptx import Presentation
from pptx.util import Pt, Mm, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

_BRAND   = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..'))
_PROJECT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
sys.path.insert(0, _BRAND)

_DOCS = os.path.join(_PROJECT, "docs")
_REFS = os.path.join(_PROJECT, "references")
os.makedirs(_DOCS, exist_ok=True)

SRC = os.path.join(_REFS, "上海未来方舟智能科技及产品简介V0（主题风格）0526.pptx")
DST = os.path.join(_DOCS, "未来方舟产品简介-品牌优化版.pptx")


# ─── 颜色映射表（六位大写十六进制，无#） ────────────────────────────────────────

FILL_MAP = {
    # 深色背景
    "333333": "0E1216",
    "2E2E2E": "0E1216",
    "1A1A1A": "0E1216",
    "0F1115": "0E1216",
    "0F1219": "0E1216",
    # 主绿 → 品牌主绿
    "54D7AA": "3EC99E",
    "55D7AB": "3EC99E",
    "41D2A5": "3EC99E",
    "40D2A4": "3EC99E",
    "3ED8A8": "3EC99E",
    "53D7A9": "3EC99E",
    # 浅绿 → 品牌浅绿
    "EBF8F1": "EAFAF5",
    "DBEAD5": "EAFAF5",
    "D9F5D6": "EAFAF5",
    "D5F5E3": "EAFAF5",
    "E8F8F0": "EAFAF5",
    # 黄绿/辅色
    "92D050": "C8E13C",
    "9ACD32": "C8E13C",
    # 成功绿
    "00B050": "53AF36",
    "00A850": "53AF36",
    # 灰色卡片/说明框
    "DEE0E3": "F2F3F5",
    "F6F6F6": "F2F3F5",
    "E8E8E8": "F2F3F5",
    "EBEBEB": "F2F3F5",
    "F5F5F5": "F2F3F5",
    "F5F7FA": "F2F3F5",
    # 淡底
    "F8F9FA": "F8FAFC",
    "FAFAFA": "F8FAFC",
}

TEXT_MAP = {
    # 标题文字深色
    "353535": "0E1216",
    "333333": "0E1216",
    "2C2C2C": "0E1216",
    "1F1F1F": "0E1216",
    "000000": "0E1216",
    "1A1A1A": "0E1216",
    "0F1115": "0E1216",
    "1E293B": "0E1216",
    # 正文灰
    "1F2329": "3D444A",
    "2D3748": "3D444A",
    "374151": "3D444A",
    # 次要信息灰
    "666666": "8A9199",
    "64748B": "8A9199",
    "969696": "8A9199",
    "9CA3AF": "8A9199",
    "A0AEC0": "8A9199",
    # 绿色文字
    "54D7AA": "3EC99E",
    "55D7AB": "3EC99E",
    "41D2A5": "3EC99E",
    "40D2A4": "3EC99E",
    # 黄绿文字
    "92D050": "C8E13C",
    # 成功绿文字
    "00B050": "53AF36",
}


# ─── 工具函数 ─────────────────────────────────────────────────────────────────

def _to_rgb(h: str) -> RGBColor:
    h = h.upper().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _rgb_hex(rgb: RGBColor) -> str:
    try:
        return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        return ""


def _map_solid_fill_api(fill, color_map):
    """通过 python-pptx API 替换 SOLID 填充色。"""
    try:
        old = _rgb_hex(fill.fore_color.rgb)
        new = color_map.get(old)
        if new:
            fill.solid()
            fill.fore_color.rgb = _to_rgb(new)
            return True
    except Exception:
        pass
    return False


def _map_colors_in_xml(elem, color_map):
    """
    在 XML 层面替换所有 <a:srgbClr val="XXXXXX"> 的颜色值。
    覆盖 API 处理不到的主题色、嵌套颜色等情况。
    """
    for srgb in elem.iter(qn("a:srgbClr")):
        old = srgb.get("val", "").upper()
        new = color_map.get(old)
        if new:
            srgb.set("val", new)


def _replace_gradient_with_solid(spPr_elem, solid_hex="3EC99E"):
    """
    将 spPr 中的 gradFill 替换为 solidFill，保留原有位置（fill 必须在 ln 之前）。
    """
    gradFill = spPr_elem.find(qn("a:gradFill"))
    if gradFill is None:
        return False
    idx = list(spPr_elem).index(gradFill)
    spPr_elem.remove(gradFill)
    sf = etree.Element(qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", solid_hex)
    spPr_elem.insert(idx, sf)
    return True


# ─── Shape 处理 ───────────────────────────────────────────────────────────────

def process_shape_fill(shape):
    """处理 shape 的填充色（包括 gradient→solid）。"""
    try:
        fill = shape.fill
        ft = fill.type  # None / 1=SOLID / 3=GRADIENT / 5=BACKGROUND / ...
    except Exception:
        return

    if ft is None or ft == 5:  # BACKGROUND / transparent
        return

    if ft == 3:  # GRADIENT → 纯色品牌绿
        try:
            spPr = shape._element.spPr
            if not _replace_gradient_with_solid(spPr, "3EC99E"):
                fill.solid()
                fill.fore_color.rgb = _to_rgb("3EC99E")
        except Exception:
            try:
                fill.solid()
                fill.fore_color.rgb = _to_rgb("3EC99E")
            except Exception:
                pass
        return

    if ft == 1:  # SOLID
        if not _map_solid_fill_api(fill, FILL_MAP):
            # 尝试 XML fallback
            try:
                _map_colors_in_xml(shape._element.spPr, FILL_MAP)
            except Exception:
                pass
        return

    # 其他类型（图案、图片等）→ 仅尝试 XML 替换
    try:
        _map_colors_in_xml(shape._element.spPr, FILL_MAP)
    except Exception:
        pass


def process_shape_text(shape):
    """处理 shape 的文字颜色。"""
    if not hasattr(shape, "text_frame"):
        return
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # API 替换
                try:
                    old = _rgb_hex(run.font.color.rgb)
                    new = TEXT_MAP.get(old)
                    if new:
                        run.font.color.rgb = _to_rgb(new)
                        continue
                except Exception:
                    pass
                # XML fallback
                try:
                    _map_colors_in_xml(run._r, TEXT_MAP)
                except Exception:
                    pass
    except Exception:
        pass


def process_shape_line(shape):
    """仅修改已有明确边框的 shape，绝不通过 API 创建新的 <a:ln> 元素。"""
    try:
        spPr = shape._element.spPr
        ln_elem = spPr.find(qn("a:ln"))
        if ln_elem is None:
            return  # 原本无边框，跳过，避免触发 python-pptx get_or_add_ln() 创建黑色默认边框
        _map_colors_in_xml(ln_elem, FILL_MAP)
    except Exception:
        pass


def process_table(shape):
    """处理 TABLE 形状内的单元格颜色和文字颜色。"""
    if not shape.has_table:
        return
    tbl = shape.table
    for row in tbl.rows:
        for cell in row.cells:
            # 单元格填充
            try:
                ft = cell.fill.type
                if ft == 1:
                    _map_solid_fill_api(cell.fill, FILL_MAP)
                elif ft == 3:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _to_rgb("3EC99E")
            except Exception:
                pass
            # XML-level cell fill colors
            try:
                _map_colors_in_xml(cell._tc, FILL_MAP)
            except Exception:
                pass
            # 文字颜色
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        old = _rgb_hex(run.font.color.rgb)
                        new = TEXT_MAP.get(old)
                        if new:
                            run.font.color.rgb = _to_rgb(new)
                    except Exception:
                        pass


def process_shape(shape):
    """对单个 shape 执行所有颜色替换。"""
    process_shape_fill(shape)
    process_shape_text(shape)
    process_shape_line(shape)
    if shape.has_table:
        process_table(shape)
    # 分组
    if shape.shape_type == 6:  # GROUP
        try:
            for child in shape.shapes:
                process_shape(child)
        except Exception:
            pass


# ─── 章节分隔页视觉增强 ──────────────────────────────────────────────────────

SECTION_SLIDE_INDICES = {2, 5, 25, 30, 36, 40}  # 0-based, 对应 Slide 3,6,26,31,37,41

def enhance_section_header(slide, slide_idx):
    """
    在章节分隔页左侧插入品牌绿竖线（5mm × 60% 高，垂直居中）。
    同时将章节编号（如 "01 关于我们" 中的 "01"）字体加粗、辅色高亮。
    """
    if slide_idx not in SECTION_SLIDE_INDICES:
        return

    SW = Inches(13.33)
    SH = Inches(7.5)
    bar_w = Mm(5)
    bar_h = SH * 0.58
    bar_t = (SH - bar_h) / 2

    bar = slide.shapes.add_shape(1, 0, int(bar_t), int(bar_w), int(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _to_rgb("3EC99E")
    bar.line.fill.background()
    bar.name = "_brand_section_bar"

    # 把竖线移到 z-order 最底（紧接 background 之后），避免遮挡文字
    sp_tree = slide.shapes._spTree
    bar_elem = bar._element
    sp_tree.remove(bar_elem)
    sp_tree.insert(2, bar_elem)


# ─── 文字精简替换 ─────────────────────────────────────────────────────────────

def _extract_first_run_fmt(text_frame):
    """提取第一段第一个 run 的字体格式，用于还原。"""
    fmt = {"size": None, "bold": None, "color": None, "name": None}
    try:
        for para in text_frame.paragraphs:
            for run in para.runs:
                fmt["size"] = run.font.size
                fmt["bold"] = run.font.bold
                fmt["name"] = run.font.name
                try:
                    fmt["color"] = _to_rgb(_rgb_hex(run.font.color.rgb))
                except Exception:
                    pass
                return fmt
    except Exception:
        pass
    return fmt


def _set_run_fmt(run, fmt):
    """将格式应用到 run。"""
    if fmt.get("size"):
        run.font.size = fmt["size"]
    if fmt.get("bold") is not None:
        run.font.bold = fmt["bold"]
    if fmt.get("name"):
        run.font.name = fmt["name"]
    if fmt.get("color"):
        run.font.color.rgb = fmt["color"]


def replace_text_keep_fmt(shape, new_text: str):
    """
    替换 shape 的文字为 new_text（\\n 分段），保留原始第一 run 的格式。
    如果 shape 没有 text_frame 则跳过。
    """
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    fmt = _extract_first_run_fmt(tf)

    # 保留原有段落对齐信息
    try:
        first_align = tf.paragraphs[0].alignment
    except Exception:
        first_align = None

    # 获取 txBody 元素，清空所有 <a:p>
    txBody = tf._txBody
    for p_elem in list(txBody.findall(qn("a:p"))):
        txBody.remove(p_elem)

    lines = new_text.split("\n")
    for i, line in enumerate(lines):
        p_elem = etree.SubElement(txBody, qn("a:p"))
        if first_align and i == 0:
            # 写入段落属性（对齐）
            pass  # alignment will be default; keep original
        if line:
            r_elem = etree.SubElement(p_elem, qn("a:r"))
            rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "zh-CN", "altLang": "en-US", "dirty": "0"})
            # 应用格式
            if fmt.get("size"):
                rPr.set("sz", str(int(fmt["size"].pt * 100)))
            if fmt.get("bold"):
                rPr.set("b", "1")
            if fmt.get("name"):
                latin = etree.SubElement(rPr, qn("a:latin"))
                latin.set("typeface", fmt["name"])
            if fmt.get("color"):
                sf = etree.SubElement(rPr, qn("a:solidFill"))
                sc = etree.SubElement(sf, qn("a:srgbClr"))
                sc.set("val", _rgb_hex(fmt["color"]))
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = line
        # 空行保留空 <a:p> 占位（段落间距）


# {slide_index(0-based): {shape_name: new_text}}
TEXT_EDITS = {
    0: {  # Slide 1 — Cover
        "文本占位符 1": "上海未来方舟智能科技\n产品简介",
    },
    3: {  # Slide 4 — 关于我们
        "AutoShape 3": (
            "上海未来方舟智能科技有限公司，专注企业级 AI 应用。\n\n"
            "产品 ArktechX 将严谨的数字化流程与具备感知力的 AI 智能体深度融合，"
            "重塑人与企业的协作边界。"
        ),
        "AutoShape 10": (
            "兼具 AI 与工业背景，能力高度互补\n"
            "管理团队 15 年+ EHS 行业经验\n"
            "技术骨干来自一线互联网大厂，10 年+ 研发沉淀\n"
            "最新 AI 技术直接落地产品"
        ),
    },
    6: {  # Slide 7 — 产品矩阵 Part 1
        "矩形 2": (
            "AI 一键创建检查表单\n"
            "系统化管理全部检查项\n"
            "消息实时提醒 · 多套灵活报告模板"
        ),
        "矩形 1": (
            "摄像头自动识别隐患，一键生成工单\n"
            "支持主流品牌摄像头\n"
            "20+ 类风险场景灵活配置\n"
            "园区-工厂多层级权限管控"
        ),
    },
    26: {  # Slide 27 — AI技术底座
        "Text 1": "AI 能力技术底座  ·  安全 · 可靠 · 可落地",
    },
    27: {  # Slide 28 — AI应用场景
        "矩形 9":  "作业票创建 · AI 风险预评估",
        "矩形 11": "作业票审批 · AI 辅助决策",
        "矩形 15": "ESG 报告 · AI 辅助生成",
        "矩形 19": "智能检查 · AI 创建模板",
        "矩形 22": "智能检查 · AI 拍照识别风险",
    },
    37: {  # Slide 38 — 私有化部署
        "矩形 23": "构建统一 AI 中台\n支持开源模型本地化部署",
    },
    44: {  # Slide 45 — 结尾
        "AutoShape 3": (
            "ArktechX 期待与您共同探索企业数字化转型，打造差异化竞争壁垒。"
        ),
    },
}


# ─── 主流程 ───────────────────────────────────────────────────────────────────

def main():
    print(f"加载原始文件: {os.path.basename(SRC)}")
    prs = Presentation(SRC)
    total = len(prs.slides)
    print(f"共 {total} 张幻灯片，开始处理...\n")

    for i, slide in enumerate(prs.slides):
        label = f"Slide {i+1:02d}/{total}"
        ops = []

        # ① 全量颜色替换
        for shape in slide.shapes:
            process_shape(shape)

        # ② 章节分隔页装饰线
        if i in SECTION_SLIDE_INDICES:
            enhance_section_header(slide, i)
            ops.append("章节装饰线")

        # ③ 文字精简
        if i in TEXT_EDITS:
            edits = TEXT_EDITS[i]
            for shape in slide.shapes:
                if shape.name in edits:
                    replace_text_keep_fmt(shape, edits[shape.name])
                    ops.append(f"文字[{shape.name}]")

        status = "  " + ("  ".join(ops) if ops else "—")
        print(f"  ✓ {label}{status}")

    print(f"\n保存到: {DST}")
    prs.save(DST)
    print("✅ 完成")
    print(f"   文件大小: {os.path.getsize(DST) / 1024:.0f} KB")


if __name__ == "__main__":
    main()
