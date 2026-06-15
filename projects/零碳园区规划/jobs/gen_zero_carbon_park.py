import sys, os

_BRAND   = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..'))
_PROJECT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
sys.path.insert(0, _BRAND)

_DOCS  = os.path.join(_PROJECT, 'docs')
_MEDIA = os.path.join(_PROJECT, 'media')
os.makedirs(_DOCS, exist_ok=True)
os.makedirs(os.path.join(_DOCS, 'img'), exist_ok=True)

from scripts.pptx_builder import (
    BrandPptx,
    _header, _footer, _set_slide_bg, _txb, _rect,
    ML, CW, CONTENT_Y, CONTENT_H, SLIDE_W, SLIDE_H, FOOTER_H,
    C2_W, C2_GAP,
)
import scripts.brand_tokens as BT
from pptx.util import Mm
from pptx.enum.text import PP_ALIGN


# ── Custom slide: AI video detection grid ─────────────────────────────────────
def _add_ai_detection_grid(prs_obj,
                            title, subtitle,
                            row1_label, row1_items,
                            row2_label, row2_items):
    """
    row1_items / row2_items: list of (img_path, bold_label, detail_label)
    Lays out two rows of image tiles with labels.
    """
    slide = prs_obj._new_slide()
    _set_slide_bg(slide, BT.WHITE_HEX)
    _header(slide, title, subtitle=subtitle)
    _footer(slide)

    PAD   = Mm(1)
    start = CONTENT_Y + Mm(4)

    # ── Row 1 ─────────────────────────────────────────────────────────────────
    n1       = len(row1_items)
    gap1     = Mm(4)
    cell_w1  = int((CW - (n1 - 1) * gap1) // n1)
    img_w1   = cell_w1 - 2 * PAD
    img_h1   = int(img_w1 / 1.45)  # auto-scale from avg aspect ratio
    txt_h    = Mm(12)

    # section label + accent line
    _txb(slide, row1_label,
         l=ML, t=start, w=CW, h=Mm(10),
         sz=12, bold=True, color=BT.PRIMARY_500_HEX)
    _rect(slide, l=ML, t=start + Mm(10), w=CW, h=Mm(0.8),
          fill=BT.PRIMARY_500_HEX)

    img_top1 = start + Mm(12)
    for i, (img_path, lbl, detail) in enumerate(row1_items):
        cx = ML + i * (cell_w1 + gap1)
        # accent chip
        _rect(slide, l=cx, t=img_top1 - Mm(2.5), w=cell_w1, h=Mm(2.5),
              fill=BT.PRIMARY_500_HEX)
        # image or placeholder
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path,
                                     cx + PAD, img_top1,
                                     img_w1, img_h1)
        else:
            _rect(slide, cx + PAD, img_top1, img_w1, img_h1,
                  fill=BT.PRIMARY_100_HEX)
        ty = img_top1 + img_h1 + Mm(2)
        _txb(slide, lbl,
             l=cx, t=ty, w=cell_w1, h=Mm(7),
             sz=9, bold=True, color=BT.NEUTRAL_900_HEX,
             align=PP_ALIGN.CENTER)
        _txb(slide, detail,
             l=cx, t=ty + Mm(7), w=cell_w1, h=Mm(6),
             sz=8, color=BT.NEUTRAL_400_HEX,
             align=PP_ALIGN.CENTER)

    # ── Row 2 ─────────────────────────────────────────────────────────────────
    n2      = len(row2_items)
    gap2    = Mm(4)
    cell_w2 = int((CW - (n2 - 1) * gap2) // n2)
    img_w2  = cell_w2 - 2 * PAD
    img_h2  = int(img_w2 / 1.45)

    row2_start = img_top1 + img_h1 + txt_h + Mm(5)

    _txb(slide, row2_label,
         l=ML, t=row2_start, w=CW, h=Mm(10),
         sz=12, bold=True, color=BT.SECONDARY_500_HEX)
    _rect(slide, l=ML, t=row2_start + Mm(10), w=CW, h=Mm(0.8),
          fill=BT.SECONDARY_500_HEX)

    img_top2 = row2_start + Mm(12)
    for i, (img_path, lbl, detail) in enumerate(row2_items):
        cx = ML + i * (cell_w2 + gap2)
        _rect(slide, l=cx, t=img_top2 - Mm(2.5), w=cell_w2, h=Mm(2.5),
              fill=BT.SECONDARY_500_HEX)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path,
                                     cx + PAD, img_top2,
                                     img_w2, img_h2)
        else:
            _rect(slide, cx + PAD, img_top2, img_w2, img_h2,
                  fill=BT.NEUTRAL_100_HEX)
        ty = img_top2 + img_h2 + Mm(2)
        _txb(slide, lbl,
             l=cx, t=ty, w=cell_w2, h=Mm(7),
             sz=9, bold=True, color=BT.NEUTRAL_900_HEX,
             align=PP_ALIGN.CENTER)
        _txb(slide, detail,
             l=cx, t=ty + Mm(7), w=cell_w2, h=Mm(6),
             sz=8, color=BT.NEUTRAL_400_HEX,
             align=PP_ALIGN.CENTER)

    return slide


# ── Custom slide: scene with photo (left: pain points | right: photo + solution) ─
def _add_scene_with_photo(prs_obj, title, subtitle,
                           left_title, left_content,
                           right_title, right_content,
                           photo_path, aspect_ratio=1.5):
    slide = prs_obj._new_slide()
    _set_slide_bg(slide, BT.WHITE_HEX)
    _header(slide, title, subtitle=subtitle)
    _footer(slide)

    top   = CONTENT_Y + Mm(5)
    h_use = CONTENT_H - Mm(5)
    rx    = ML + C2_W + C2_GAP   # right column x-start

    # Left column: pain points (unchanged style)
    _txb(slide, left_title,
         l=ML, t=top, w=C2_W, h=Mm(12),
         sz=13, bold=True, color=BT.PRIMARY_500_HEX)
    _txb(slide, left_content,
         l=ML, t=top + Mm(14), w=C2_W, h=h_use - Mm(14),
         sz=14, color=BT.NEUTRAL_700_HEX, ls_pt=22)

    # Vertical divider
    _rect(slide, l=ML + C2_W + C2_GAP // 2, t=top + Mm(4),
          w=Mm(0.4), h=h_use - Mm(8), fill=BT.NEUTRAL_200_HEX)

    # Right column: brand accent + photo (capped at 70mm height)
    MAX_H = Mm(70)
    img_w_full = C2_W
    img_h_full = int(img_w_full / aspect_ratio)
    if img_h_full > MAX_H:
        img_h = MAX_H
        img_w = int(MAX_H * aspect_ratio)
    else:
        img_h = img_h_full
        img_w = img_w_full

    img_x = rx + (C2_W - img_w) // 2   # center if narrower than column

    _rect(slide, l=rx, t=top - Mm(2.5), w=C2_W, h=Mm(2.5),
          fill=BT.PRIMARY_500_HEX)
    if os.path.exists(photo_path):
        slide.shapes.add_picture(photo_path, img_x, top, img_w, img_h)
    else:
        _rect(slide, img_x, top, img_w, img_h, fill=BT.NEUTRAL_100_HEX)

    # Solution text below photo
    txt_y  = top + img_h + Mm(4)
    txt_h  = h_use - img_h - Mm(4) - Mm(12)
    _txb(slide, right_title,
         l=rx, t=txt_y, w=C2_W, h=Mm(11),
         sz=11, bold=True, color=BT.PRIMARY_500_HEX)
    _txb(slide, right_content,
         l=rx, t=txt_y + Mm(12), w=C2_W, h=txt_h,
         sz=9, color=BT.NEUTRAL_700_HEX, ls_pt=14)

    return slide


# ── Photo paths ──────────────────────────────────────────────────────────────
_P = os.path.join(_MEDIA, 'photos')

PHOTOS_ELEC = [
    (os.path.join(_P, '无尘室特气泄漏防护.png'),  '无尘室特气泄漏',    'Fab/洁净室专属'),
    (os.path.join(_P, '危化品泄漏:火险.png'),      '危化品泄漏/火险',   '化学品存储区'),
    (os.path.join(_P, '锂电池冒烟:火险.png'),      '锂电池冒烟/火险',   '储能柜早期预警'),
    (os.path.join(_P, '高压绝缘防护.png'),         '高压绝缘防护',      '配电间安全核查'),
    (os.path.join(_P, '高温设备超温:火险.png'),    '高温超温/火险',     '服务器·充电桩热点'),
    (os.path.join(_P, '暗室高功率:火险.png'),      '暗室高功率/火险',   '微电子洁净室'),
]

PHOTOS_ENG = [
    (os.path.join(_P, '动火.png'),       '动火作业',   '烟/火识别'),
    (os.path.join(_P, '高空.png'),       '高空作业',   'PPE/动作合规'),
    (os.path.join(_P, '吊装.png'),       '吊装作业',   '隔离区/动作'),
    (os.path.join(_P, '临电.png'),       '临时用电',   '标牌上锁核查'),
    (os.path.join(_P, '有限空间.png'),   '有限空间',   'PPE/异常倒地'),
]


# ── Build presentation ────────────────────────────────────────────────────────
prs = BrandPptx()

# ── Slide 1: Cover ─────────────────────────────────────────────────────────────
prs.add_cover(
    title='零碳园区智慧运营方案',
    subtitle='以 ArktechX · Arksus 平台为底座，构建铸造 · 水泥 · 数据中心 · 锂电池园区的全栈低碳管控体系',
    tagline='8 大业务域 · 34 核心模块 · Enterprise AI & EHS · 2026'
)

# ── Slide 2: 目录 ──────────────────────────────────────────────────────────────
prs.add_body_slide(
    title='目录',
    subtitle='本方案围绕安全、合规、可持续三大主线，展示 Arksus 平台在零碳园区中的落地路径',
    bullets=[
        '01  政策背景与园区挑战 — 为什么零碳园区是当务之急',
        '02  传统管理的核心痛点 — 三条断裂带',
        '03  平台全景 — 8 大业务域 · 34 核心模块地图',
        '04  安全场景（一）— 建设期承包商管控 · 日常巡检与整改',
        '05  安全场景（二）— AI 视频行业专属检测：数据中心/锂电池 & 铸造/水泥',
        '06  可持续场景 — 物联碳感知 · ESG 安全数据闭环',
        '07  合规 & 协同场景 — 人员资质 · 供应链 Scope 3 · 治理',
        '08  综合收益与三阶段实施路径',
    ]
)

# ── Divider: Chapter 1 ─────────────────────────────────────────────────────────
prs.add_divider('政策背景与园区挑战', chapter_num='CHAPTER 01')

# ── Slide 3: 政策压力 ──────────────────────────────────────────────────────────
prs.add_body_slide(
    title='高排放园区正面临三重前所未有的合规压力',
    subtitle='3060 双碳目标 · 强制 ESG 披露 · 绿色金融准入门槛同步收紧',
    bullets=[
        '【政策端】国家 3060 目标明确，工业园区为重点监管对象；主要高排放行业已纳入全国碳市场，碳配额年度清缴压力逐年加大',
        '【市场端】苹果 · 宝马等跨国企业要求供应商提交符合 GRI/TCFD 框架的碳排放数据（Scope 3 核查），无法提供者已开始影响合同续签',
        '【融资端】绿色信贷 · 绿色债券将园区碳强度作为关键授信指标；未达绿色认证标准的园区，融资成本显著高于同类',
        '【运营端】安全事故一票否决绿色认证进度——安全合规与低碳运营必须协同管控，不能分拆治理',
    ]
)

# ── Slide 4: 四类园区特征 ──────────────────────────────────────────────────────
prs.add_three_cards(
    title='四类高排放园区：挑战高度相似，数字化起点普遍偏低',
    subtitle='铸造厂 · 水泥厂 · 数据中心 · 锂电池工厂',
    cards=[
        {
            'tag': '重工业 · 高危作业',
            'title': '铸造 & 水泥',
            'body': (
                '• 高温 · 粉尘 · 有毒气体多重职业危害\n'
                '• 每日须管控：动火 · 高空 · 受限空间三类高危作业\n'
                '• AI 视频重点：明火识别 · 高空 PPE · 吊装隔离区 · 受限空间入场合规\n'
                '• 碳排放强度大，核算基础数据极度分散'
            ),
        },
        {
            'tag': '高能耗 · 高密度',
            'title': '数据中心',
            'body': (
                '• AI 视频重点：锂电池冒烟/火险 · 高压绝缘防护 · 高温设备超温\n'
                '• PUE 是核心绿色指标，监管与托管客户双向关注\n'
                '• 7x24h 运维巡检压力大，人员配置紧张\n'
                '• ESG 报告须向托管客户和监管机构双向披露'
            ),
        },
        {
            'tag': '新能源 · 强合规',
            'title': '锂电池工厂',
            'body': (
                '• AI 视频重点：热失控早期烟雾 · 高温超温 · 危化品泄漏/火险\n'
                '• 工程作业：动火 · 临时用电管控同步覆盖\n'
                '• 下游车企 Scope 3 碳足迹核查严格，影响合同续签\n'
                '• 工艺变更须严格管控，与安全 · 碳排放联动'
            ),
        },
    ]
)

# ── Divider: Chapter 2 ─────────────────────────────────────────────────────────
prs.add_divider('传统管理的核心痛点', chapter_num='CHAPTER 02')

# ── Slide 5: 痛点 ─────────────────────────────────────────────────────────────
prs.add_three_cards(
    title='现状诊断：三条断裂带，贯穿安全 · 碳管理 · 协同',
    subtitle='每一条痛点背后，都是真实的合规风险与运营成本',
    cards=[
        {
            'tag': '安全 & 合规',
            'title': '隐患靠纸质，整改无闭环',
            'body': (
                '巡检表单手写、隐患照片靠微信传、整改结果无系统追踪。\n'
                '高危作业票审批跑断腿，中途异常无法实时通知相关人员。\n\n'
                '承包商资质造假难察觉，证书过期无预警，\n'
                '事故发生后责任认定困难。'
            ),
        },
        {
            'tag': '碳管理 & ESG',
            'title': '碳账不透明，ESG 的「S」缺数据',
            'body': (
                '能耗数据分散在各子系统，碳排放核算依赖月末人工填报，\n'
                '数据误差大、不可追溯，无法通过第三方碳核查。\n\n'
                'ESG 报告「S（社会）」章节数据来源割裂：\n'
                '失工率 · PPE 合规率 · 整改关闭率分散在多系统，无法自动汇总。'
            ),
        },
        {
            'tag': '协同 & 效率',
            'title': '园区与租户割裂，治理靠微信群',
            'body': (
                '园区运营方与企业租户间无统一数字化连接，\n'
                '政策通知 · 整改督促 · ESG 数据收集靠微信群和邮件，\n'
                '版本混乱、无法留痕。\n\n'
                '绩效考核 · 流程审批分散在多个系统，\n'
                '无法支撑园区级 ESG 管治体系。'
            ),
        },
    ]
)

# ── Divider: Chapter 3 ─────────────────────────────────────────────────────────
prs.add_divider('平台模块全景', chapter_num='CHAPTER 03')

# ── Slide 6: 8 域架构 ─────────────────────────────────────────────────────────
prs.add_body_slide(
    title='8 大业务域 · 34 核心模块：从设备感知到 ESG 战略输出',
    subtitle='各域数据互通——安全事件 · 碳排放 · 供应链记录共享同一数据底座，「S」章节数据自动流入 ESG 报告',
    bullets=[
        '【工作台  3模块】  仪表盘与看板 · 个人任务中心 · AI 智能助理 — 个性化入口，AI 降低使用门槛',
        '【安全与风险  8模块】  风险管控 · 隐患排查 · 智能检查 · 数字化作业票 · 承包商管理 · AI 视频识别 · 物联智能 · 应急物资',
        '【行动与合规  2模块】  整改中心（跨模块整改闭环）· 合规与法规（法规知识库 + AI 匹配）',
        '【资产与知识  3模块】  人员资质与培训 · 设备合规台账 · 安全知识库',
        '【环境与可持续  2模块】  环境合规（实时监测 + 政府对接）· 可持续与 ESG（碳核算 + 报告自动化）',
        '【质量与产品  3模块】  质量异常追踪 · 工艺与产品变更管理 · 产品管理与供应链审查',
        '【系统设置  6模块】  组织权限 · 流程中心 · 自定义表单 · 数据字典 · 三方集成 · 服务管理',
        '【办公协同  7模块】  协同总览 · 待办与审批 · 公告通知 · 日程会议 · 文档协作 · U Connect · 绩效与激励',
    ]
)

# ── Slide 7: 模块卡片 A ────────────────────────────────────────────────────────
prs.add_three_cards(
    title='模块全景 · 安全基础三域',
    subtitle='34 个模块 · 8 大业务域 — 安全、合规、资产三个核心运营域',
    cards=[
        {
            'tag': '8 个模块',
            'title': '安全与风险',
            'body': (
                '风险管控 · 隐患排查 · 智能检查\n'
                '数字化作业票 · 承包商管理\n'
                'AI 视频识别 · 物联智能 · 应急物资\n\n'
                '所有发现的隐患统一流入整改中心，\n'
                '形成安全感知 → 告警 → 整改 → 关闭的完整闭环'
            ),
        },
        {
            'tag': '2 个模块',
            'title': '行动与合规中心',
            'body': (
                '整改中心 · 合规与法规\n\n'
                '整改中心承接全平台隐患，统一调度跟踪；\n'
                '合规与法规内置法规知识库，\n'
                'AI 检查时自动匹配法规条文，\n'
                '让每一条整改建议都有法规依据'
            ),
        },
        {
            'tag': '3 个模块',
            'title': '资产与知识库',
            'body': (
                '人员资质与培训 · 设备合规台账\n'
                '安全知识库\n\n'
                '人 + 设备 + 知识三维台账联动；\n'
                '资质绑定 · 复训提醒自动触发，\n'
                '设备证书到期预警，\n'
                '知识库支撑 AI 培训课件自动生成'
            ),
        },
    ]
)

# ── Slide 8: 模块卡片 B ────────────────────────────────────────────────────────
prs.add_three_cards(
    title='模块全景 · 战略输出三域',
    subtitle='34 个模块 · 8 大业务域 — 环境可持续、质量供应链、协同治理三个战略域',
    cards=[
        {
            'tag': '2 个模块 · 零碳核心',
            'title': '环境与可持续',
            'body': (
                '环境合规 · 可持续与 ESG\n\n'
                '安全数据（AI 视频 · 巡检 · 整改）\n'
                '自动汇入 ESG「S（社会）」章节指标，\n'
                '无需二次填报；\n'
                '碳排放核算 + GRI/TCFD 报告自动生成，\n'
                '编写时间减少 80%'
            ),
        },
        {
            'tag': '3 个模块 · Scope 3',
            'title': '质量与产品',
            'body': (
                '质量异常追踪 · 工艺与产品变更管理\n'
                '产品管理与供应链审查\n\n'
                '锂电池 / 制造业园区的供应链碳足迹核心；\n'
                '工艺变更与碳排放联动分析，\n'
                '满足下游车企 Scope 3 核查要求，\n'
                '支撑产品 PCF 报告编制'
            ),
        },
        {
            'tag': '13 个模块 · 数字底座',
            'title': '工作台 · 协同 · 系统',
            'body': (
                'AI 智能助理 · U Connect\n'
                '流程中心 · 待办与审批 · 绩效激励\n'
                '三方集成（ERP / BMS / 政府碳报平台）\n\n'
                'U Connect 连接园区与企业租户；\n'
                '三方集成打通现有 IT 环境；\n'
                'AI 智能助理贯穿所有域'
            ),
        },
    ]
)

# ── Divider: Chapter 4 ─────────────────────────────────────────────────────────
prs.add_divider('安全场景深度应用', chapter_num='CHAPTER 04')

# ── Slide 9: 场景1 承包商 ──────────────────────────────────────────────────────
_add_scene_with_photo(
    prs,
    title='场景一：绿色建设期承包商全链路管控',
    subtitle='铸造 · 水泥 · 锂电池工厂绿色改造期，外来承包商激增、高危作业票每日数十张',
    left_title='痛点与风险',
    left_content=(
        '绿色改造期间数十支承包队伍同时入场，\n'
        '涉及动火 · 高处 · 受限空间等高危作业：\n\n'
        '• 证书过期无人预警，资质造假难以发现\n'
        '• 作业票审批流程靠纸质来回跑\n'
        '• 入场安全培训流于形式，与作业许可脱节\n'
        '• 风险驾驶舱缺失，隐患分散不聚合\n\n'
        '安全事故一旦发生，绿色认证进度直接中断，\n'
        '影响园区融资节奏与绿色债券发行计划。'
    ),
    right_title='平台解题路径',
    right_content=(
        '▸ 承包商管理\n'
        '黑名单拦截违规人员；证书到期自动预警并锁定作业票申请权限；\n'
        '每次作业留痕，数据驱动安全评分，替代凭印象考核。\n\n'
        '▸ 数字化作业票（8 类预置）\n'
        '动火 / 高处 / 受限空间，手机端创建 + AI 风险预评估 + 电子签名；\n'
        '中途异常一键通知所有相关人员中止作业。\n\n'
        '▸ 人员资质与培训\n'
        '承包商入场前完成在线安全培训 + 考试，未通过者无法创建作业票，\n'
        '形成系统级硬性入场门槛。\n\n'
        '▸ 风险管控 + 隐患排查\n'
        '建设期风险驾驶舱实时汇总所有隐患，自动流转整改中心，全程留痕。'
    ),
    photo_path=os.path.join(_P, '场景一痛点:风险.png'),
    aspect_ratio=1.78,
)

# ── Slide 10: 场景2a 日常巡检 ────────────────────────────────────────────────
prs.add_two_col_slide(
    title='场景二（上）：日常安全巡检与隐患整改闭环',
    subtitle='数据中心 · 锂电池工厂——从「发现问题」到「关闭隐患」形成可追溯的完整记录',
    left_title='痛点与风险',
    left_content=(
        '传统人工巡检的三大短板：\n\n'
        '• 班次交接存在时间盲区，夜班覆盖率无法核验\n'
        '• 隐患照片靠微信传，整改结果无系统闭环，\n'
        '  下次巡检才知道同一问题是否真正关闭\n'
        '• 发现问题 → 整改落实之间缺乏统一调度中枢，\n'
        '  责任人不明、超期无人跟催\n\n'
        '安全事故往往发生在无人值守的瞬间，\n'
        '纸质记录无法还原事故前的隐患链条。'
    ),
    right_title='平台解题路径',
    right_content=(
        '▸ 智能检查\n'
        'AI 一键生成行业专属巡检模板（电气柜 · 消防 · 化学品储存等），\n'
        '手机端扫码执行，隐患照片直传，自动派发整改任务，\n'
        '整改完成后方可关闭，过程全程留痕可追溯。\n\n'
        '▸ 整改中心\n'
        '承接全平台所有隐患（人工巡检 + AI 视频 + 物联设备告警），\n'
        '统一分配责任人 · 截止时间 · 跟踪整改进度；\n'
        '逾期自动升级通知，合规审查时一键导出完整整改台账。\n\n'
        '▸ 物联智能（联动）\n'
        '现有烟感 · 温感 · 用电监测设备接入系统，\n'
        '异常信号直接触发告警并写入整改中心，形成设备感知 → 整改闭环。'
    )
)

# ── Slide 11: 场景2b AI 视频检测图片网格（新增）────────────────────────────────
_add_ai_detection_grid(
    prs,
    title='场景二（下）：AI 视频识别 · 行业专属检测场景配置',
    subtitle='非通用监控——依托现有摄像头基础设施，针对不同园区类型预置独立检测模型，开箱即用',
    row1_label='数据中心 · 锂电池工厂 · 微电子/Fab  |  电子 & 精密化工专属检测场景（6 类）',
    row1_items=PHOTOS_ELEC,
    row2_label='铸造厂 & 水泥厂  |  工程作业场景（动火 · 高空 · 吊装 · 临电 · 有限空间）',
    row2_items=PHOTOS_ENG,
)

# ── Divider: Chapter 5 ─────────────────────────────────────────────────────────
prs.add_divider('可持续场景：物联碳感知与 ESG 报告', chapter_num='CHAPTER 05')

# ── Slide 12: 场景3 物联 & 碳感知 ────────────────────────────────────────────
_add_scene_with_photo(
    prs,
    title='场景三：物联智能驱动的碳排放实时感知',
    subtitle='物理设备数据入口——没有物联接入，碳账就是月末人工填报',
    left_title='痛点与风险',
    left_content=(
        '零碳认证的核心是数据可信 · 可追溯：\n\n'
        '• 数据中心 PUE 须实时监控并向客户证明\n'
        '• 水泥 · 铸造吨产品碳排是碳市场配额核算基础\n'
        '• 锂电池工厂须向车企提供产品碳足迹（PCF）\n\n'
        '传统困境：\n'
        '• 能耗数据分散各子系统，无统一采集入口\n'
        '• 设备运行状态靠人工记录，误差大 · 滞后严重\n'
        '• 无法区分正常生产能耗与设备故障造成的浪费\n\n'
        '碳账数据一旦被质疑，绿色认证申请即告中止。'
    ),
    right_title='平台解题路径',
    right_content=(
        '▸ 物联智能\n'
        '兼容园区现有电表 · 热表 · 气表 · 传感器；\n'
        '设备数据实时采集，接入统一数据总线；\n'
        '异常能耗自动触发告警并关联设备合规台账。\n\n'
        '▸ 环境合规\n'
        '大气排放 · 废水 · 固废等环境指标实时监测，\n'
        '自动生成合规报告，直接对接政府环保监管平台；\n'
        '碳排放数据全程留痕，满足 ISO 14064 第三方核查。\n\n'
        '▸ 设备合规台账\n'
        '统一管理园区所有能耗设备：预防性维护减少非计划停机，\n'
        '设备证书到期自动提醒，\n'
        '设备更换 / 升级与碳减排效果联动分析。'
    ),
    photo_path=os.path.join(_P, '场景三：物联智能.png'),
    aspect_ratio=1.50,
)

# ── Slide 13: 场景4 ESG ────────────────────────────────────────────────────────
prs.add_two_col_slide(
    title='场景四：ESG 报告自动化 · 安全数据直接输入「S」章节',
    subtitle='ESG 报告「S（社会）」章节的核心数据来自安全模块——两套系统打通才能真正自动化',
    left_title='痛点与风险',
    left_content=(
        'ESG 披露已是刚性义务（上市公司强制、绿色信贷授信要求、大客户核查）。\n\n'
        '传统方式三大症结：\n'
        '• 数据收集靠人工汇总，3-4 个月耗费大量时间\n'
        '• GRI / TCFD / ISO 14064 框架切换困难，\n'
        '  一套数据无法复用为多份报告\n\n'
        '被忽视的「S（社会）」数据缺口：\n'
        '• 失工伤害率（LTIR）· 职业病发生率\n'
        '• PPE 合规率 · 整改关闭率\n'
        '• 安全培训完成率 · 承包商安全评分\n\n'
        '这些数据分散在巡检记录 · AI 视频台账 ·\n'
        '整改中心 · 培训系统里，无法自动汇总。'
    ),
    right_title='平台解题路径',
    right_content=(
        '▸ 安全数据 → ESG「S」章节（核心打通）\n'
        'AI 视频识别 PPE 合规率\n'
        '整改中心 → 整改关闭率 / 平均响应时间\n'
        '人员培训模块 → 安全培训完成率 / 考试通过率\n'
        '承包商管理 → 承包商安全评分分布\n'
        '以上数据在系统内自动汇总，直接写入 ESG「S」章节，\n'
        '无需二次填报，数据来源可追溯。\n\n'
        '▸ 可持续与 ESG\n'
        'AI 对话式收集能耗账单 / 检查记录 / 员工健康数据，\n'
        '一键生成 GRI / TCFD 框架完整报告；\n'
        '编写时间从 3-4 个月缩至数周；\n'
        '数据存储在园区私有环境，不上传第三方。\n\n'
        '▸ 三方集成\n'
        'ERP / BMS / 政府碳报平台数据自动同步，\n'
        '消除人工导出环节，保障碳排放数据一致性。'
    )
)

# ── Divider: Chapter 6 ─────────────────────────────────────────────────────────
prs.add_divider('合规与协同场景', chapter_num='CHAPTER 06')

# ── Slide 14: 场景5 人员资质 ──────────────────────────────────────────────────
_add_scene_with_photo(
    prs,
    title='场景五：人员资质与培训全生命周期合规',
    subtitle='从入职到离场，每一张证书 · 每一次培训都有数字档案',
    left_title='痛点与风险',
    left_content=(
        '高排放园区人员合规复杂度极高：\n\n'
        '• 铸造 / 水泥 / 锂电池工厂大量特种作业人员\n'
        '  （焊工 · 电工 · 起重工 · 危化品操作工）\n'
        '• 职业病危害类别多（粉尘 · 噪声 · 有毒有害），\n'
        '  健康监测周期长，易遗漏\n'
        '• 年度安全培训组织难 · 考试阅卷费时，\n'
        '  电子学习档案缺失，合规审查时手忙脚乱\n\n'
        '特种作业证书失效 · 职业病认定纠纷，\n'
        '企业将面临行政处罚与巨额赔偿。'
    ),
    right_title='平台解题路径',
    right_content=(
        '▸ 人员资质与培训\n'
        '特种作业证书统一管理，到期前自动提醒复训；\n'
        'AI 自动生成专项安全课件（焊接 / 化学品 / 电气等），\n'
        '游戏化学习提升完课率；电子档案完整留痕，\n'
        '一键生成合规培训报告应对审查。\n\n'
        '▸ 整改中心（人员行为联动）\n'
        'AI 视频识别发现行为违规（未戴安全帽 · 违章操作）\n'
        '直接触发定向培训推送，\n'
        '形成违规 → 整改 → 培训的闭环。\n\n'
        '▸ 绩效与激励\n'
        '安全行为积分与绩效挂钩，正向激励全员安全合规参与，\n'
        '同时将培训完成率自动汇入 ESG「S」章节。'
    ),
    photo_path=os.path.join(_P, '场景五：人员资质培训.png'),
    aspect_ratio=1.50,
)

# ── Slide 15: 场景6 供应链碳足迹 ─────────────────────────────────────────────
_add_scene_with_photo(
    prs,
    title='场景六：供应链碳足迹管控（Scope 3）',
    subtitle='锂电池 · 化工 · 精密制造：下游客户碳核查已成入场门槛',
    left_title='痛点与风险',
    left_content=(
        '供应链碳管控是新能源产业园区的必答题：\n\n'
        '• 车企要求电池供应商提交每批产品的\n'
        '  产品碳足迹（PCF）报告，无法提供则失去订单\n'
        '• 工艺变更（配方调整 · 设备更换）会直接影响\n'
        '  单位产品碳排放，变更管控缺失导致碳账失真\n'
        '• 原材料供应商碳强度差异大，\n'
        '  无系统化审查机制无法识别高碳风险供应商\n\n'
        '质量异常与碳排放数据割裂，\n'
        '无法还原产品全生命周期碳排放路径。'
    ),
    right_title='平台解题路径',
    right_content=(
        '▸ 工艺与产品变更管理\n'
        '所有工艺变更须经数字化审批流程，\n'
        '变更前自动评估对碳排放强度的影响，\n'
        '变更记录永久留存，支撑 PCF 报告可追溯性。\n\n'
        '▸ 产品管理与供应链审查\n'
        '供应商碳强度数据纳入准入审查；\n'
        '定期供应链碳排放评估报告自动生成，\n'
        '高碳风险供应商触发整改通知或黑名单机制。\n\n'
        '▸ 质量异常追踪\n'
        '质量缺陷与生产能耗异常联动分析，\n'
        '定位高碳排放的生产环节，\n'
        '为工艺优化和减排路径规划提供数据基础。\n\n'
        '▸ 可持续与 ESG\n'
        '供应链碳数据自动汇入 ESG 报告 Scope 3 章节，统一披露。'
    ),
    photo_path=os.path.join(_P, '场景六：碳排放.png'),
    aspect_ratio=1.78,
)

# ── Slide 16: 场景7 协同治理 ─────────────────────────────────────────────────
prs.add_two_col_slide(
    title='场景七：园区数字化协同治理',
    subtitle='园区运营方与企业租户之间，需要一套统一的数字化连接器',
    left_title='痛点与风险',
    left_content=(
        '零碳园区不只是技术问题，更是治理问题：\n\n'
        '• 园区运营方须统一向数十家租户企业\n'
        '  发布安全政策 · 碳排放要求 · 整改通知\n'
        '• ESG 数据收集靠逐家发邮件 · 填表格，\n'
        '  数据不一致 · 版本混乱\n'
        '• 跨部门审批（绿色认证 · 大型作业许可 ·\n'
        '  设备采购）流程冗长，纸质签字追不上\n'
        '• 没有统一绩效体系，安全与低碳工作\n'
        '  无法量化激励，难以形成文化'
    ),
    right_title='平台解题路径',
    right_content=(
        '▸ U Connect 协同\n'
        '园区运营方与所有企业租户建立统一数字连接；\n'
        '通知 · 整改要求 · ESG 数据请求一键推送，\n'
        '反馈记录全部留存，替代微信群和邮件链。\n\n'
        '▸ 待办与审批 + 流程中心\n'
        '绿色认证 · 设备更换 · 大型作业等审批流程\n'
        '通过流程中心自定义配置，手机端一键签批，\n'
        '全程可追踪 · 可审计。\n\n'
        '▸ 公告与通知 + 日程与会议\n'
        '碳配额清缴提醒 · 安全演练通知 · ESG 委员会会议\n'
        '统一在平台内管理，确保重要节点不遗漏。\n\n'
        '▸ 绩效与激励\n'
        '安全合规 · 减排贡献纳入量化绩效指标，\n'
        '正向激励零碳文化在全园区落地。'
    )
)

# ── Divider: Chapter 7 ─────────────────────────────────────────────────────────
prs.add_divider('综合价值收益与实施路径', chapter_num='CHAPTER 07')

# ── Slide 17: 价值收益 ────────────────────────────────────────────────────────
prs.add_big_stats(
    title='五维价值收益：安全 · 碳管理 · ESG · 供应链 · 品牌',
    subtitle='每一项收益均有具体模块能力支撑',
    stats=[
        ('80%+', 'ESG 报告编写效率提升', '可持续与ESG：3-4个月压缩至数周；安全数据自动汇入「S」章节，无需二次填报'),
        ('7x24h', 'AI 视频无盲区监控', '行业专属检测模型预置：数据中心锂电/高压/高温 + 工程作业动火/高空/有限空间'),
        ('Scope 3', '供应链碳足迹可追溯', '工艺变更管控 + 供应链审查 + 质量异常三模块联动，PCF 报告可出具'),
        ('34模块', '全链路数字化覆盖', '安全 · 碳 · 质量 · 协同数据统一底座，「S」章节数据自动流入 ESG 报告'),
    ]
)

# ── Slide 18: 实施路径 ────────────────────────────────────────────────────────
prs.add_timeline(
    title='三阶段实施路径：18 个月从安全基础到 ESG 战略全面落地',
    subtitle='分阶段推进，每阶段均有独立交付价值，私有化部署，数据不出园区',
    milestones=[
        (
            'Phase 1  0-6M',
            '安全基础层',
            '数字化作业票 · 承包商管理\n'
            '智能检查 · AI 视频识别（行业专属）\n'
            '整改中心 · 物联智能接入\n'
            '消除最高优先级安全风险'
        ),
        (
            'Phase 2  6-12M',
            '合规与碳数据层',
            '人员资质与培训 · 设备合规台账\n'
            '环境合规 · 质量与供应链模块\n'
            '三方集成打通 BMS / ERP\n'
            '绿色认证申请数据链路成形'
        ),
        (
            'Phase 3  12-18M',
            'ESG 战略层',
            '可持续与 ESG 全面启用\n'
            'U Connect 连接园区与租户\n'
            '绩效与激励落地\n'
            '首份 GRI / TCFD 报告对外发布'
        ),
    ]
)

# ── Slide 19: Quote ────────────────────────────────────────────────────────────
prs.add_quote(
    quote_text=(
        '零碳园区不是一个认证标牌，而是一套每天运转的数字化体系。\n'
        '安全不出事，碳账算得清，ESG 的「S」章节有真实数据支撑——\n'
        '这三件事同时做到，才是真正的零碳园区运营中枢。'
    ),
    author='ArktechX 零碳园区解决方案团队',
    role='8 大业务域 · 34 核心模块 · Enterprise AI & EHS · 2026'
)

# ── Slide 20: Closing ─────────────────────────────────────────────────────────
prs.add_closing(
    message='共建零碳园区新标准',
    contact='info@ark-techx.com  ·  17887924530  ·  www.ark-techx.com'
)

# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(_DOCS, '零碳园区智慧运营方案.pptx')
prs.save(out)
print('Generated:', out)
